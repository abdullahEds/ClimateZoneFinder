import streamlit as st
import base64
import pandas as pd
import plotly.express as px
import numpy as np
import io
import re
import pytz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import tempfile
import os


st.set_page_config(
    page_title="Climate Analytics Dashboard",
    # page_icon="🌍",
    layout="wide"
)

def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode()

logo_base64 = get_base64_image("images/EDSlogo.jpg")

# STEP 1: Create header using st.columns for proper layout
col_logo, col_title, col_home = st.columns([1, 4, 1])

with col_logo:
    st.markdown(
        f'<img src="data:image/png;base64,{logo_base64}" style="height: 80px; margin-top: 45px;">',
        unsafe_allow_html=True
    )

with col_title:
    st.markdown(
        '<h2 style="text-align: center; color: #a85c42; margin-top: 45px;">Climate Analytics Dashboard</h2>',
        unsafe_allow_html=True
    )

with col_home:
    pass

def get_base64_image(image_path):
    try:
        with open(image_path, "rb") as img_file:
            return base64.b64encode(img_file.read()).decode()
    except:
        try:
            with open(f"../{image_path}", "rb") as img_file:
                return base64.b64encode(img_file.read()).decode()
        except:
            return ""

st.markdown("""
    <style>
    
    /* Style the header columns container */
    div[data-testid="stHorizontalBlock"]:first-of-type {
        border-bottom: 1px solid #e6e6e6;
        padding-bottom: 20px;
        margin-bottom: 0px;
        background-color: white;
    }
    

    
    /* Logo hover effect */
    div[data-testid="stHorizontalBlock"]:first-of-type img:hover {
        transform: scale(1.05);
        opacity: 0.85;
        transition: all 0.3s ease;
    }
    

    </style>
""", unsafe_allow_html=True)

# === HEADER ===
# st.markdown("""
#     <style>
#     .header-container {
#         background: linear-gradient(135deg, #1a3a52 0%, #2c5aa0 100%);
#         padding: 0px;
#         border-radius: 0;
#         margin-top: 50px;
#         box-shadow: 0 4px 12px rgba(0,0,0,0.2);
#         border-bottom: 4px solid #ffffff;

#            /* 👈 pushes it below Streamlit toolbar */
#         left: 0;
#         right: 0;
#         z-index: 999;
#         width: 100%;
#         box-sizing: border-box;
#     }
#     style>
#     /* Hide top toolbar */
#     header[data-testid="stHeader"] {
#         display: none;
#     }

#     /* Hide hamburger menu */
#     #MainMenu {
#         visibility: hidden;
#     }

#     /* Hide footer */
#     footer {
#         visibility: hidden;
#     }

#     /* Remove top padding since header is gone */
#     .block-container {
#         padding-top: 1rem;
#     }

#     /* Optional: Remove deploy button */
#     div[data-testid="stToolbar"] {
#         display: none;
#     }
#     /* Adjust body spacing to avoid overlap */
#     .main > div {
#         padding-top: 180px;   /* Increase if needed */
#     }

#     .header-content {
#         display: flex;
#         align-items: center;
#         gap: 20px;
#     }

#     .header-icon {
#         font-size: 48px;
#         display: inline-block;
#     }

#     .header-title {
#         color: #ffffff;
#         font-size: 32px;
#         font-weight: 800;
#         margin: 0;
#         letter-spacing: 0.5px;
#     }
#     </style>

#     <div class="header-container">
#         <div class="header-content">
#             <div class="header-icon">🌍</div>
#             <div class="header-title">Climate Analytics Dashboard</div>
#         </div>
#     </div>
# """, unsafe_allow_html=True)

# === INITIALIZE SESSION STATE FOR TABS ===
if "active_tab" not in st.session_state:
    st.session_state.active_tab = "Annual Trend"

st.markdown("""
    <style>
    /* Main layout adjustments */
    .block-container {
        padding-top: 0rem !important;
    }
    
    section[data-testid="stSidebar"] {
        display: none !important;
    }
    
    button[kind="header"] {
        display: none !important;
    }
    
    .main .block-container {
        max-width: 100% !important;
        padding-left: 1.5rem !important;
        padding-right: 1.5rem !important;
    }
    
    
    
    .control-section-header {
        font-size: 15px;
        font-weight: 700;
        color: #2c3e50;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        margin-bottom: 12px;
        margin-top: 16px;
        display: flex;
        align-items: center;
        gap: 8px;
        width: 200px;
    }
    
    .control-section-header:first-child {
        margin-top: 0;
    }
    
    /* Upload zone styling */
    .upload-zone {
        border: 2px dashed #cbd5e0;
        border-radius: 6px;
        padding: 12px;
        text-align: center;
        background-color: #f7fafc;
        margin-top: 8px;
    }
    
    .upload-zone.success {
        border-color: #48bb78;
        background-color: #f0fff4;
    }
    
    /* File uploader styling */
    [data-testid="fileUploadDropzone"] {
        border-radius: 6px !important;
        border-color: #cbd5e0 !important;
    }
    
    /* Success message styling */
    .stAlert[data-baseweb="notification"] {
        background-color: #f0fff4 !important;
        border-left: 4px solid #48bb78 !important;
        border-radius: 4px !important;
    }
    
    /* Slider styling */
    .stSlider {
        margin-top: 12px;
    }
    
    /* Date input styling */
    .stDateInput {
        margin-top: 12px;
    }
    
    /* Selectbox styling */
    .stSelectbox {
        margin-top: 8px;
    }
    
    /* Section styling */
    .section-title {
        font-size: 18px;
        font-weight: 700;
        color: #2c3e50;
        margin-bottom: 16px;
        border-bottom: 2px solid #3498db;
        padding-bottom: 8px;
        display: inline-block;
    }
    
    /* KPI Cards */
    .kpi-card {
        background: white;
        padding: 16px;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
        text-align: center;
    }
    
    .kpi-label {
        font-size: 11px;
        font-weight: 700;
        color: #718096;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        margin-bottom: 8px;
        opacity: 0.9;
    }
    
    .kpi-value {
        font-size: 26px;
        font-weight: 700;
        color: #2c3e50;
        margin: 8px 0;
    }
    
    .kpi-meta {
        font-size: 11px;
        color: #718096;
        opacity: 0.85;
    }
    </style>
""", unsafe_allow_html=True)

def _ppt_remove_all_slides(prs):
    """Remove every slide from an open Presentation while preserving theme/layouts."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)


def generate_pptx_report(df: pd.DataFrame, start_date, end_date, start_hour: int, end_hour: int, selected_parameter: str, metadata: dict = None):
    """Generate a PowerPoint report using the Voha template with Dry Bulb, Relative Humidity and Sun Path sections."""

    # ── Resolve paths ──────────────────────────────────────────────────────────
    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    # ── Open template and strip existing slides ────────────────────────────────
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]

    # ── Colors matching the template ──────────────────────────────────────────
    TITLE_RED   = RGBColor(0xC0, 0x00, 0x00)   # #C00000 – template title colour
    DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER_CLR = RGBColor(0xC0, 0x00, 0x00)

    # Slide canvas size (inches)
    SW = prs.slide_width.inches   # ≈ 13.33
    SH = prs.slide_height.inches  # ≈ 7.50

    # Logo dimensions (maintain aspect ratio 550:308 ≈ 1.786)
    LOGO_H  = 0.40  # inches
    LOGO_W  = LOGO_H * (550 / 308)
    LOGO_L  = 0.18
    LOGO_T  = SH - LOGO_H - 0.12

    def _add_logo(slide):
        """Place EDSlogo.png at the bottom-left corner of the slide."""
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(LOGO_L), Inches(LOGO_T),
                width=Inches(LOGO_W), height=Inches(LOGO_H)
            )

    def _add_slide_title(slide, text, left=0.27, top=0.13, width=None, height=0.45):
        """Add section title in the template's dark-red style."""
        if width is None:
            width = SW - left - 0.3
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _add_divider(slide, top_inches):
        """Thin horizontal rule matching template style."""
        line = slide.shapes.add_shape(
            1, Inches(0.27), Inches(top_inches),
            Inches(SW - 0.54), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = DIVIDER_CLR
        line.line.fill.background()

    def _save_mpl_figure(fig) -> str:
        """Save a matplotlib figure to a temp PNG and return its path."""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches='tight', facecolor='white')
            return tmp.name

    def _err_box(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Visualization error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    # Filter data for period
    filtered_df = df[
        (df["datetime"].dt.date >= start_date) &
        (df["datetime"].dt.date <= end_date) &
        (df["hour"].between(start_hour, end_hour))
    ]

    # ═══════════════════════════════════════════════════════════════════════════
    # COVER SLIDE
    # ═══════════════════════════════════════════════════════════════════════════
    def _make_cover_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)

        # Title block background
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        # Main title
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(SW - 1.2), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Climate Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Subtitle - City/Location name
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.7))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        _location = metadata.get("location", "") if metadata else ""
        location_display = _city if _city else (_location if _location else "Location")
        run2.text = f"{location_display}"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        # Bottom info line
        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Dry Bulb Temperature  |  Relative Humidity  |  Sun Path"
        run3.font.size = Pt(11)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _make_cover_slide()

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 1 – DRY BULB TEMPERATURE
    # ═══════════════════════════════════════════════════════════════════════════

    # --- 1a: Annual Trend Chart ---
    def _make_dbt_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                temp_min=("dry_bulb_temperature", "min"),
                temp_max=("dry_bulb_temperature", "max"),
                temp_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
            comfort_line = daily_avg.rolling(window=7, center=True).mean()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=130)
            ax.fill_between(daily_stats["doy"], comfort_line - 3.5, comfort_line + 3.5,
                            alpha=0.18, color='gray', label='ASHRAE 80% Comfort')
            ax.fill_between(daily_stats["doy"], comfort_line - 2.5, comfort_line + 2.5,
                            alpha=0.28, color='gray', label='ASHRAE 90% Comfort')
            ax.fill_between(daily_stats["doy"], daily_stats["temp_min"], daily_stats["temp_max"],
                            alpha=0.30, color='#FFB3B3', label='Daily Temp Range')
            ax.plot(daily_stats["doy"], daily_stats["temp_avg"],
                    color='#C00000', linewidth=2.2, label='Daily Average', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Annual Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=5, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_trend_slide()

    # --- 1b: Monthly Averages Bar Chart ---
    def _make_dbt_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly = df.groupby("month").agg(
                t_min=("dry_bulb_temperature", "min"),
                t_max=("dry_bulb_temperature", "max"),
                t_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=130)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly["t_min"], bar_w, color='#90CAF9', label='Min Temp')
            ax.bar(x,          monthly["t_avg"], bar_w, color='#C00000', label='Avg Temp', alpha=0.85)
            ax.bar(x + bar_w,  monthly["t_max"], bar_w, color='#EF9A9A', label='Max Temp')

            # Comfort band zone
            ax.axhspan(20, 26, alpha=0.10, color='green', label='Comfort Band (20–26°C)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Monthly Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            # Annotate period months
            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            # Key stats text box (bottom right)
            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min: {filtered_df['dry_bulb_temperature'].min():.1f}°C  "
                    f"Avg: {filtered_df['dry_bulb_temperature'].mean():.1f}°C  "
                    f"Max: {filtered_df['dry_bulb_temperature'].max():.1f}°C  "
                    f" |  Ann. CDD24: {(df['dry_bulb_temperature'] - 24).clip(lower=0).sum():.0f}   "
                    f"HDD18: {(18 - df['dry_bulb_temperature']).clip(lower=0).sum():.0f}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_monthly_slide()

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 2 – RELATIVE HUMIDITY
    # ═══════════════════════════════════════════════════════════════════════════

    # --- 2a: Annual Trend Chart ---
    def _make_rh_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=130)
            ax.axhspan(75, 100, alpha=0.13, color='#FF6B6B', label='Condensation Risk (>75%)')
            ax.axhspan(60,  75, alpha=0.13, color='#FFA500', label='High RH (60–75%)')
            ax.axhspan(30,  60, alpha=0.13, color='#4ECDC4', label='Comfortable (30–60%)')
            ax.axhspan( 0,  30, alpha=0.13, color='#FFD93D', label='Low RH (<30%)')

            ax.fill_between(daily_stats["doy"], daily_stats["rh_min"], daily_stats["rh_max"],
                            alpha=0.28, color='#0099ff', label='Daily RH Range')
            ax.plot(daily_stats["doy"], daily_stats["rh_avg"],
                    color='#0066cc', linewidth=2.2, label='Daily Average RH', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 100)
            ax.set_title('Annual Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_trend_slide()

    # --- 2b: Monthly RH Bar Chart ---
    def _make_rh_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly_rh = df.groupby("month").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=130)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly_rh["rh_min"], bar_w, color='#AED6F1', label='Min RH')
            ax.bar(x,          monthly_rh["rh_avg"], bar_w, color='#0066cc', label='Avg RH', alpha=0.85)
            ax.bar(x + bar_w,  monthly_rh["rh_max"], bar_w, color='#5DADE2', label='Max RH')

            ax.axhspan(30, 60, alpha=0.10, color='green', label='Comfortable (30–60%)')
            ax.axhline(75, color='#E74C3C', linewidth=1.2, linestyle='--', label='Condensation Threshold (75%)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 110)
            ax.set_title('Monthly Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min RH: {filtered_df['relative_humidity'].min():.0f}%  "
                    f"Avg RH: {filtered_df['relative_humidity'].mean():.0f}%  "
                    f"Max RH: {filtered_df['relative_humidity'].max():.0f}%  "
                    f" |  High RH hrs (>60%): {len(filtered_df[filtered_df['relative_humidity'] > 60])}  "
                    f"Condensation risk hrs (>75%): {len(filtered_df[filtered_df['relative_humidity'] > 75])}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_monthly_slide()

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 3 – SUN PATH
    # ═══════════════════════════════════════════════════════════════════════════

    def _make_sun_path_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Sun Path Diagram")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        lat = _meta.get("latitude")
        lon = _meta.get("longitude")
        tz_str = _meta.get("timezone", "UTC")

        if lat is None or lon is None:
            _err_box(slide, "Latitude/Longitude not available from EPW metadata.")
            _add_logo(slide)
            return

        try:
            from pvlib import solarposition as _solpos_lib
            import pytz as _pytz

            try:
                _tz = _pytz.timezone(tz_str)
            except Exception:
                _tz = _pytz.UTC

            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left")
            sol = _solpos_lib.get_solarposition(times, lat, lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()
            sol["r"] = 90 - sol["apparent_elevation"]

            # Build a matplotlib polar sun path chart (clean, static, print-ready)
            fig = plt.figure(figsize=(8.5, 7.2), dpi=130, facecolor='white')
            ax = fig.add_subplot(111, projection='polar')
            ax.set_theta_zero_location('N')
            ax.set_theta_direction(-1)   # clockwise = compass convention
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(['90°\n(Zenith)', '75°', '60°', '45°', '30°', '15°', '0°\n(Horizon)'],
                               fontsize=7, color='#555')
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'], fontsize=10, fontweight='bold')
            ax.set_facecolor('#F0F4F8')
            ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)

            # Scatter all-year sun positions coloured by day-of-year
            sc = ax.scatter(
                np.radians(sol["azimuth"].values),
                sol["r"].values,
                c=sol.index.dayofyear,
                cmap='YlOrRd',
                s=1.0,
                alpha=0.55,
                vmin=1, vmax=365,
                linewidths=0,
                zorder=2,
            )
            cbar = fig.colorbar(sc, ax=ax, pad=0.10, fraction=0.035, shrink=0.75)
            cbar.set_label('Day of Year', fontsize=9)
            cbar.set_ticks([1, 91, 182, 273, 365])
            cbar.set_ticklabels(['1\n(Jan)', '91\n(Apr)', '182\n(Jul)', '273\n(Oct)', '365\n(Dec)'])

            # Key date arcs
            key_dates = [
                ("Mar 21 (Spring Equinox)", "2020-03-21", "#FF9500", 1.6),
                ("Jun 21 (Summer Solstice)", "2020-06-21", "#CC0000", 2.0),
                ("Dec 21 (Winter Solstice)", "2020-12-21", "#0066CC", 2.0),
            ]
            for lbl, dstr, col, lw in key_dates:
                dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
                ks = _solpos_lib.get_solarposition(dt, lat, lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.06), ncol=3,
                      frameon=True, fontsize=8, borderaxespad=0)
            ax.set_title(f'Sun Path  |  Lat: {lat:.2f}°  Lon: {lon:.2f}°',
                         fontsize=11, fontweight='bold', color='#333', pad=14)

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)

            # Centre the sun path diagram on the slide
            img_w = SW * 0.62
            img_h = SH * 0.83
            img_l = (SW - img_w) / 2
            img_t = 0.72
            slide.shapes.add_picture(tmp, Inches(img_l), Inches(img_t), width=Inches(img_w), height=Inches(img_h))
            os.unlink(tmp)

            # Annotation panel on the right
            ann_l = img_l + img_w + 0.2
            ann_w = SW - ann_l - 0.27

            def _ann(slide, text, top, size=9, bold=False, color=None):
                tb = slide.shapes.add_textbox(Inches(ann_l), Inches(top), Inches(ann_w), Inches(0.38))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color or DARK_GREY

            # _ann(slide, "Location", 0.75, size=8, bold=True, color=TITLE_RED)
            # _ann(slide, f"Latitude:   {lat:.3f}\u00b0", 1.05, size=8)
            # _ann(slide, f"Longitude: {lon:.3f}\u00b0", 1.36, size=8)

            # _ann(slide, "Key Dates", 1.80, size=8, bold=True, color=TITLE_RED)
            # _ann(slide, "\u25a0  Mar 21 – Spring Equinox", 2.10, size=8)
            # _ann(slide, "\u25a0  Jun 21 – Summer Solstice", 2.40, size=8)
            # _ann(slide, "\u25a0  Dec 21 – Winter Solstice", 2.70, size=8)

            # Summer / Winter altitude stats
            # summer_dt = pd.date_range("2020-06-21", periods=24, freq="h", tz=_tz)
            # ssum = _solpos_lib.get_solarposition(summer_dt, lat, lon)
            # noon_alt_sum = float(ssum.iloc[12]["apparent_elevation"])

            # winter_dt = pd.date_range("2020-12-21", periods=24, freq="h", tz=_tz)
            # swin = _solpos_lib.get_solarposition(winter_dt, lat, lon)
            # noon_alt_win = float(swin.iloc[12]["apparent_elevation"])

            # _ann(slide, "Noon Altitudes", 3.15, size=8, bold=True, color=TITLE_RED)
            # _ann(slide, f"Summer Solstice: {noon_alt_sum:.1f}\u00b0", 3.45, size=8)
            # _ann(slide, f"Winter Solstice: {noon_alt_win:.1f}\u00b0", 3.75, size=8)

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_sun_path_slide()

    # ═══════════════════════════════════════════════════════════════════════════
    # SHADING ANALYSIS SLIDE
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _plot_sun_path_shading(lat, lon, tz_str):
        """Generate a sun path diagram with horizontal overhang shading profile."""
        try:
            from pvlib import solarposition as _sp
            import pytz as _pz
            
            try:
                _tz = _pz.timezone(tz_str)
            except:
                _tz = _pz.UTC
            
            # Generate hourly data for full year
            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left")
            sol = _sp.get_solarposition(times, lat, lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()
            sol["r"] = 90 - sol["apparent_elevation"]
            
            # Create polar plot
            fig = plt.figure(figsize=(8.0, 7.0), dpi=130, facecolor='white')
            ax = fig.add_subplot(111, projection='polar')
            ax.set_theta_zero_location('N')
            ax.set_theta_direction(-1)
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(['90°', '75°', '60°', '45°', '30°', '15°', '0°'], fontsize=7, color='#555')
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'], fontsize=9, fontweight='bold')
            ax.set_facecolor('#F0F4F8')
            ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)
            
            # Plot sun positions colored by season
            sc = ax.scatter(
                np.radians(sol["azimuth"].values),
                sol["r"].values,
                c=sol.index.dayofyear,
                cmap='YlOrRd',
                s=0.8,
                alpha=0.6,
                vmin=1, vmax=365,
                linewidths=0,
                zorder=2,
            )
            cbar = fig.colorbar(sc, ax=ax, pad=0.10, fraction=0.035, shrink=0.7)
            cbar.set_label('Day of Year', fontsize=8)
            cbar.set_ticks([1, 91, 182, 273, 365])
            cbar.set_ticklabels(['Jan', 'Apr', 'Jul', 'Oct', 'Dec'], fontsize=7)
            
            # Plot key date arcs
            key_dates = [
                ("Equinox", "2020-03-21", "#FF9500", 1.5),
                ("Summer", "2020-06-21", "#CC0000", 1.8),
                ("Winter", "2020-12-21", "#0066CC", 1.8),
            ]
            for lbl, dstr, col, lw in key_dates:
                dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
                ks = _sp.get_solarposition(dt, lat, lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)
            
            # Add shading profile (horizontal overhang with D/H ratio ~0.8)
            # This represents an overhang that blocks summer sun but allows winter sun
            overhang_altitude = 35  # degrees - typical overhang cutoff
            theta_range = np.radians(np.linspace(45, 315, 100))  # South-facing (SE to SW)
            shading_altitude = np.ones_like(theta_range) * overhang_altitude
            
            ax.fill_between(theta_range, shading_altitude, 90, alpha=0.12, color='#8B4513', 
                            label='Shading Zone', zorder=1)
            ax.plot(theta_range, shading_altitude, color='#654321', linewidth=2.5, 
                   label='Overhang Profile', linestyle='--', zorder=3)
            
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.08), ncol=4,
                     frameon=True, fontsize=7, borderaxespad=0)
            ax.set_title(f'Sun Path with Shading Profile\nLat: {lat:.2f}°  Lon: {lon:.2f}°',
                        fontsize=10, fontweight='bold', color='#333', pad=12)
            
            plt.tight_layout()
            return fig
        except Exception as e:
            print(f"Shading diagram error: {e}")
            return None
    
    def _make_shading_summary_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Shading Strategy")
        _add_divider(slide, 0.62)

        # Try to add sun path shading diagram
        diagram_added = False
        _meta = metadata or {}
        _lat = _meta.get("latitude")
        _lon = _meta.get("longitude")
        _tz_str = _meta.get("timezone", "UTC")
        
        if _lat is not None and _lon is not None:
            try:
                shading_fig = _plot_sun_path_shading(_lat, _lon, _tz_str)
                if shading_fig is not None:
                    tmp_shading = _save_mpl_figure(shading_fig)
                    plt.close(shading_fig)
                    # Position diagram on left side of slide
                    slide.shapes.add_picture(tmp_shading, Inches(0.27), Inches(0.75), 
                                           width=Inches(3.5), height=Inches(5.0))
                    os.unlink(tmp_shading)
                    diagram_added = True
            except Exception as e:
                print(f"Shading diagram error: {e}")
        
        # Add text content (right side if diagram added, otherwise full width)
        if diagram_added:
            text_left = 3.9
            text_width = SW - text_left - 0.27
            tb = slide.shapes.add_textbox(Inches(text_left), Inches(0.80), 
                                          Inches(text_width), Inches(5.8))
        else:
            tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        
        tf = tb.text_frame
        tf.word_wrap = True

        # Solar Geometry & Shading Analysis
        p = tf.paragraphs[0]
        p.text = "Solar Geometry & Shading Analysis"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(4)

        try:
            # Extract metadata for location
            _meta = metadata or {}
            _lat = _meta.get("latitude")
            _lon = _meta.get("longitude")
            _tz_str = _meta.get("timezone", "UTC")
            
            if _lat is None or _lon is None:
                raise ValueError("Latitude/Longitude not available")
            
            from pvlib import solarposition as _solpos_lib_shade
            import pytz as _pytz_shade
            
            try:
                _tz_obj = _pytz_shade.timezone(_tz_str)
            except Exception:
                _tz_obj = _pytz_shade.UTC
            
            # Calculate solar angles for Key dates
            summer_dt = pd.date_range("2020-06-21", periods=24, freq="h", tz=_tz_obj)
            ssum = _solpos_lib_shade.get_solarposition(summer_dt, _lat, _lon)
            noon_alt_sum = float(ssum.iloc[12]["apparent_elevation"])
            
            winter_dt = pd.date_range("2020-12-21", periods=24, freq="h", tz=_tz_obj)
            swin = _solpos_lib_shade.get_solarposition(winter_dt, _lat, _lon)
            noon_alt_win = float(swin.iloc[12]["apparent_elevation"])
            
            p = tf.add_paragraph()
            p.text = f"• Summer Solstice (Jun 21): Solar altitude at noon = {noon_alt_sum:.1f}°"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(3)
            
            p = tf.add_paragraph()
            p.text = f"• Winter Solstice (Dec 21): Solar altitude at noon = {noon_alt_win:.1f}°"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(3)
            
            # Calculate shading metrics from filtered data
            ghi_col = filtered_df.get("global_horizontal_irradiance", pd.Series(0))
            if len(ghi_col) == 0:
                ghi_col = pd.Series(0, index=filtered_df.index)
            ghi_col = ghi_col.fillna(0)
            
            temp_col = filtered_df.get("dry_bulb_temperature", pd.Series(20))
            if len(temp_col) == 0:
                temp_col = pd.Series(20, index=filtered_df.index)
            temp_col = temp_col.fillna(20)
            
            # Shading thresholds: temp > 28°C AND GHI > 315 Wh/m²
            shading_needed = (temp_col > 28) & (ghi_col > 315)
            shading_hours = shading_needed.sum() / 2  # Each row is 30 minutes
            total_observation_hours = len(filtered_df) / 2
            
            if total_observation_hours > 0:
                shading_pct = (shading_hours / total_observation_hours) * 100
            else:
                shading_pct = 0
            
            p = tf.add_paragraph()
            p.text = f"• Shading required: {shading_hours:.0f} hours ({shading_pct:.1f}% of period)"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(6)
        except Exception as e:
            p = tf.add_paragraph()
            p.text = "• Solar altitude and shading data not available"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(6)

        # Shading Recommendations
        p = tf.add_paragraph()
        p.text = "Shading Recommendations"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        recommendations = [
            "• South-facing facades: Use horizontal overhangs (louvers) or shading devices to block summer sun while allowing winter sunlight penetration",
            "• East/West facades: Use vertical fins or combination of overhangs and fins to minimize morning/afternoon heat gain",
            "• North-facing facades: Minimal shading required; prioritize daylighting and views",
            "• Use high-performance glazing with low solar heat gain coefficient (SHGC) in high solar radiation areas",
            "• Consider automated shading systems for dynamic climate response throughout the year"
        ]

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.0
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        # Design Considerations
        p = tf.add_paragraph()
        p.text = "Design Considerations"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(3)

        considerations = [
            "• Optimal window-to-wall ratio: 30-40% for climate comfort; balance daylighting with thermal performance",
            "• Depth of shading device: D/H ratio (depth to height) between 0.5-1.0 for effective summer shading",
            "• Material selection: High-albedo surfaces reflect solar radiation; low-emissivity coatings minimize thermal transmission"
        ]

        for cons in considerations:
            p = tf.add_paragraph()
            p.text = cons
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.0
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_shading_summary_slide()

    # ═══════════════════════════════════════════════════════════════════════════
    # ANNEXURE SLIDE – About EDS, Disclaimer, Acknowledgement
    # ═══════════════════════════════════════════════════════════════════════════
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Annexure")
        _add_divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        # About EDS
        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "Environmental Design Solutions [EDS] is a sustainability advisory firm. Since 2002, EDS has worked on over 500 green building and energy efficiency projects worldwide. The team focuses on climate change mitigation, low-carbon design, building simulation, performance audits, and capacity building. EDS continues to contribute to the buildings community with useful tools through its IT services."
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)
        p.level = 0

        # Disclaimer
        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        disclaimer_items = [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "\u2022  EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "\u2022  EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "\u2022  Open-source resources used: Clima - Berkley, Streamlit, Python",
            "\u2022  EDS is not liable to inform Users about updates to the Application or underlying resources"
        ]

        for item in disclaimer_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            p.level = 0

        # Acknowledgement
        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        ack_items = [
            "\u2022  Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "\u2022  Streamlit, \u00a9 Streamlit Inc., licensed under Apache 2.0",
            "\u2022  Python \u00a9 Python Software Foundation, licensed under PSF License Version 2"
        ]

        for item in ack_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            p.level = 0

        _add_logo(slide)

    _make_annexure_slide()

    # ── Save ───────────────────────────────────────────────────────────────────
    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes

def generate_shading_pptx_report(
    df: pd.DataFrame,
    metadata: dict,
    temp_threshold: float = 28.0,
    rad_threshold: float = 315.0,
    lat: float = None,
    lon: float = None,
    tz_str: str = "UTC",
    design_cutoff_angle: float = 45.0,
):
    """Generate a Shading Analysis PowerPoint report using the Voha template."""

    # ── Resolve paths ─────────────────────────────────────────────
    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]
    TITLE_RED  = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY  = RGBColor(0x40, 0x40, 0x40)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches
    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(LOGO_L), Inches(LOGO_T),
                                     width=Inches(LOGO_W), height=Inches(LOGO_H))

    def _slide_title(slide, text, top=0.13):
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(top), Inches(SW - 0.54), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _divider(slide, top_inches):
        bar = slide.shapes.add_shape(1, Inches(0.27), Inches(top_inches),
                                     Inches(SW - 0.54), Inches(0.03))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TITLE_RED
        bar.line.fill.background()

    def _err(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = TITLE_RED

    def _save_fig(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches="tight", facecolor="white")
            return tmp.name

    _lat  = lat  if lat  is not None else (metadata.get("latitude")  or 0.0)
    _lon  = lon  if lon  is not None else (metadata.get("longitude") or 0.0)
    _tz   = tz_str or metadata.get("timezone", "UTC")

    # ═══════════════════════════════════════════════════════════
    # SLIDE 1 – COVER
    # ═══════════════════════════════════════════════════════════
    def _cover():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(SW), Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(SW - 1.2), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Shading Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.75), Inches(SW - 1.2), Inches(0.65))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        run2.text = (
            f"Location: {_city}   |  "
            # f"Temp threshold: {temp_threshold}\u00b0C   |  "
            # f"Radiation threshold: {rad_threshold} W/m\u00b2   |  "
            # f"Design cutoff angle: {design_cutoff_angle}\u00b0"
        )
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Thermal & Radiation Matrix  |  Sun Path (Shading Mode)  |  Orientation Analysis  |  Shading Masks"
        run3.font.size = Pt(10)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _cover()

    # ═══════════════════════════════════════════════════════════
    # SLIDE 2 – THERMAL & RADIATION MATRIX
    # ═══════════════════════════════════════════════════════════
    def _thermal_matrix_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Thermal & Radiation Matrix")
        _divider(slide, 0.62)

        try:
            temp_matrix, rad_matrix, overheat_mask = build_thermal_matrix(df, temp_threshold, rad_threshold)
            months_lbl = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
            hours_lbl  = [f"{h:02d}:00" for h in range(24)]

            fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.2), dpi=120)

            for ax, matrix, title, cmap, clabel in [
                (axes[0], temp_matrix,  f"Mean Dry-Bulb Temp (\u00b0C)  [threshold: {temp_threshold}\u00b0C]", "RdYlBu_r", "\u00b0C"),
                (axes[1], rad_matrix,   f"Mean GHI (W/m\u00b2)  [threshold: {rad_threshold} W/m\u00b2]",  "YlOrRd",    "W/m\u00b2"),
            ]:
                im = ax.imshow(matrix.values, aspect="auto", origin="upper", cmap=cmap)
                plt.colorbar(im, ax=ax, fraction=0.035, pad=0.03, label=clabel)
                ax.set_xticks(range(12))
                ax.set_xticklabels(months_lbl, fontsize=8)
                ax.set_yticks(range(24))
                ax.set_yticklabels(hours_lbl, fontsize=7)
                ax.set_title(title, fontsize=10, fontweight="bold", pad=8)
                ax.set_xlabel("Month", fontsize=9)
                ax.set_ylabel("Hour of Day", fontsize=9)

                # Outline overheating cells
                for h_i in range(24):
                    for m_i in range(12):
                        if overheat_mask.iloc[h_i, m_i]:
                            rect = plt.Rectangle(
                                (m_i - 0.5, h_i - 0.5), 1, 1,
                                fill=False, edgecolor="black", linewidth=1.6
                            )
                            ax.add_patch(rect)

            fig.suptitle(
                "Overheating Hours  (black border = both thresholds exceeded)",
                fontsize=11, fontweight="bold", y=1.01, color="#333"
            )
            fig.patch.set_facecolor("white")
            plt.tight_layout()

            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _thermal_matrix_slide()

    # ═══════════════════════════════════════════════════════════
    # SLIDE 3 – SUN PATH (SHADING MODE)
    # ═══════════════════════════════════════════════════════════
    def _sun_path_shading_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Sun Path – Shading Analysis")
        _divider(slide, 0.62)

        try:
            import pytz as _pytz
            from pvlib import solarposition as _sp

            try:
                tz = _pytz.timezone(_tz)
            except Exception:
                tz = _pytz.UTC

            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=tz, inclusive="left")
            sol = _sp.get_solarposition(times, _lat, _lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()

            # Merge EPW data
            epw = df.set_index("datetime").copy()
            if epw.index.tz is None:
                epw.index = epw.index.tz_localize(tz)
            else:
                epw.index = epw.index.tz_convert(tz)
            epw.index = epw.index.map(lambda x: x.replace(year=2020))

            sol = sol.join(epw[["dry_bulb_temperature","global_horizontal_irradiance"]], how="left")
            sol["global_horizontal_irradiance"] = sol["global_horizontal_irradiance"].fillna(0)
            sol["dry_bulb_temperature"] = sol["dry_bulb_temperature"].fillna(
                sol["dry_bulb_temperature"].median()
            )
            shading_needed = (
                (sol["dry_bulb_temperature"] > temp_threshold) &
                (sol["global_horizontal_irradiance"] > rad_threshold)
            )

            fig = plt.figure(figsize=(9, 7.2), dpi=130, facecolor="white")
            ax = fig.add_subplot(111, projection="polar")
            ax.set_theta_zero_location("N")
            ax.set_theta_direction(-1)
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(["90\u00b0","75\u00b0","60\u00b0","45\u00b0","30\u00b0","15\u00b0","0\u00b0"],
                               fontsize=7, color="#555")
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(["N","NE","E","SE","S","SW","W","NW"], fontsize=10, fontweight="bold")
            ax.set_facecolor("#F0F4F8")
            ax.grid(True, alpha=0.35, linestyle="--", linewidth=0.6)

            r = 90 - sol["apparent_elevation"].values
            theta = np.radians(sol["azimuth"].values)

            # Background – no-shading hours
            mask_ok = ~shading_needed.values
            ax.scatter(theta[mask_ok], r[mask_ok], c="#FFF9C4", s=1.2,
                       alpha=0.45, linewidths=0, label="No shading needed", zorder=2)

            # Foreground – shading required hours
            mask_sh = shading_needed.values
            ax.scatter(theta[mask_sh], r[mask_sh], c="#E65100", s=2.5,
                       alpha=0.75, linewidths=0, label="Shading required", zorder=3)

            # Key date arcs
            for lbl, dstr, col, lw in [
                ("Mar 21", "2020-03-21", "#FF9500", 1.4),
                ("Jun 21", "2020-06-21", "#CC0000", 1.8),
                ("Dec 21", "2020-12-21", "#0066CC", 1.8),
            ]:
                dt = pd.date_range(dstr, periods=288, freq="5min", tz=tz)
                ks = _sp.get_solarposition(dt, _lat, _lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            shading_pct = mask_sh.sum() / len(mask_sh) * 100 if len(mask_sh) else 0
            ax.set_title(
                f"Sun Path – Shading Mode   ({shading_pct:.1f}% of daytime hours require shading)",
                fontsize=10, fontweight="bold", color="#333", pad=14
            )
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.06), ncol=3,
                      frameon=True, fontsize=8)

            plt.tight_layout()
            tmp = _save_fig(fig)
            plt.close(fig)

            iw = SW * 0.58
            ih = SH * 0.84
            il = (SW - iw) / 2
            slide.shapes.add_picture(tmp, Inches(il), Inches(0.72),
                                     width=Inches(iw), height=Inches(ih))
            os.unlink(tmp)

            # Stats panel – right side
            pl = il + iw + 0.18
            pw = SW - pl - 0.27

            def _ann(text, top, size=9, bold=False, color=None):
                tb = slide.shapes.add_textbox(Inches(pl), Inches(top), Inches(pw), Inches(0.38))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color or DARK_GREY

            # _ann("Location", 0.75, bold=True, color=TITLE_RED)
            # _ann(f"Lat: {_lat:.3f}\u00b0", 1.08)
            # _ann(f"Lon: {_lon:.3f}\u00b0", 1.38)
            # _ann("Thresholds", 1.78, bold=True, color=TITLE_RED)
            # _ann(f"Temp > {temp_threshold}\u00b0C", 2.10)
            # _ann(f"GHI > {rad_threshold} W/m\u00b2", 2.40)
            # n_sh = int(mask_sh.sum())
            # n_total = len(mask_sh)
            # _ann("Shading Stats", 2.80, bold=True, color=TITLE_RED)
            # _ann(f"Total daytime pts: {n_total}", 3.12)
            # _ann(f"Shading required: {n_sh}", 3.42)
            # _ann(f"({shading_pct:.1f}% of daytime)", 3.72)

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _sun_path_shading_slide()

    # ═══════════════════════════════════════════════════════════
    # SLIDE 4 – ORIENTATION SHADING ANALYSIS TABLE
    # ═══════════════════════════════════════════════════════════
    def _orientation_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, f"Orientation Shading Analysis  (Design cutoff: {design_cutoff_angle}\u00b0)")
        _divider(slide, 0.62)

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz)
            if solar_pos.empty:
                _err(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_df = build_orientation_table(solar_pos, design_cutoff_angle)

            # ── Build matplotlib table ──────────────────────────────
            fig, ax = plt.subplots(figsize=(13, 5.8), dpi=120)
            ax.axis("off")

            col_labels = ["Orientation","Rays\nHitting","Min VSA","Max |HSA|",
                          "D/H\nOverhang","D/W\nFin","Protection %"]
            table_data = []
            row_colors = []
            for _, row in orient_df.iterrows():
                pct = row["Protection (%)"]
                if pct is None:
                    c = "#f5f5f5"
                elif pct >= 95:
                    c = "#e8f5e9"
                elif pct >= 80:
                    c = "#fff3e0"
                else:
                    c = "#ffebee"
                row_colors.append([c] * 7)

                dh = f"{row['D/H (Overhang)']:.3f}" if row["D/H (Overhang)"] is not None else "—"
                dw = f"{row['D/W (Fin)']:.3f}"       if row["D/W (Fin)"] is not None else "—"
                vsa = f"{row['Min VSA (°)']:.1f}\u00b0"  if row["Min VSA (°)"] is not None else "—"
                hsa = f"{row['Max |HSA| (°)']:.1f}\u00b0" if row["Max |HSA| (°)"] is not None else "—"
                pct_s = f"{pct:.1f}%" if pct is not None else "—"
                table_data.append([
                    row["Orientation"],
                    str(row["Rays Hitting"]),
                    vsa, hsa, dh, dw, pct_s
                ])

            tbl = ax.table(
                cellText=table_data,
                colLabels=col_labels,
                cellColours=row_colors,
                loc="center",
                cellLoc="center",
            )
            tbl.auto_set_font_size(False)
            tbl.set_fontsize(11)
            tbl.scale(1, 2.1)

            for j in range(len(col_labels)):
                cell = tbl[0, j]
                cell.set_facecolor("#1a3a52")
                cell.set_text_props(color="white", fontweight="bold")

            # Left-align orientation column
            for i in range(1, len(table_data) + 1):
                tbl[i, 0].get_text().set_ha("left")

            fig.patch.set_facecolor("white")
            ax.set_title(
                f"{len(solar_pos)} overheating daytime sun positions  |  "
                f"Temp > {temp_threshold}\u00b0C  &  GHI > {rad_threshold} W/m\u00b2",
                fontsize=10, color="#555", pad=10
            )

            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _orientation_slide()

    # ═══════════════════════════════════════════════════════════
    # SLIDE 5 – SHADING MASK DIAGRAMS (8 orientations, 2×4 grid)
    # ═══════════════════════════════════════════════════════════
    def _shading_masks_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Shading Mask Diagrams")
        _divider(slide, 0.62)

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz)
            if solar_pos.empty:
                _err(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_items = list(_ORIENTATIONS.items())  # 8 items

            # 2 rows × 4 cols grid
            n_cols = 4
            cell_w = (SW - 0.54) / n_cols
            cell_h = (SH - 1.05) / 2  # 2 rows, leave room for title + logo
            top_start = 0.75

            for idx, (oname, faz) in enumerate(orient_items):
                col_i = idx % n_cols
                row_i = idx // n_cols
                cell_l = 0.27 + col_i * cell_w
                cell_t = top_start + row_i * cell_h

                geom = compute_shading_geometry(solar_pos, faz)
                facing = geom[geom["hits_facade"]]
                other  = geom[~geom["hits_facade"]]

                fig = plt.figure(figsize=(3.0, 2.8), dpi=110, facecolor="white")
                ax = fig.add_subplot(111, projection="polar")
                ax.set_theta_zero_location("N")
                ax.set_theta_direction(-1)
                ax.set_ylim(0, 90)
                ax.set_yticks([])
                ax.set_xticks(np.radians([0, 90, 180, 270]))
                ax.set_xticklabels(["N","E","S","W"], fontsize=8, fontweight="bold")
                ax.set_facecolor("#f0f8ff")
                ax.grid(True, alpha=0.30, linewidth=0.5)

                if not other.empty:
                    ax.scatter(
                        np.radians(other["solar_azimuth"]), 90 - other["solar_altitude"],
                        s=2, c="lightgrey", alpha=0.5, linewidths=0, zorder=2
                    )
                if not facing.empty:
                    ax.scatter(
                        np.radians(facing["solar_azimuth"]), 90 - facing["solar_altitude"],
                        s=4, c="#E65100", alpha=0.75, linewidths=0, zorder=3
                    )

                # Shading cutoff arc
                rel_az_r = np.linspace(-89, 89, 179)
                tan_co = np.tan(np.radians(design_cutoff_angle))
                co_alt = np.degrees(np.arctan(tan_co * np.cos(np.radians(rel_az_r))))
                co_az  = faz + rel_az_r
                valid  = co_alt > 0
                if valid.any():
                    ax.plot(np.radians(co_az[valid]), 90 - co_alt[valid],
                            color="#1565C0", linewidth=1.4, linestyle="--", zorder=4)

                # Facade direction
                ax.plot(np.radians([faz, faz]), [0, 85],
                        color="#388E3C", linewidth=1.4, zorder=5)

                ax.set_title(oname, fontsize=7, fontweight="bold", pad=4, color="#222")
                plt.tight_layout(pad=0.3)

                tmp = _save_fig(fig)
                plt.close(fig)
                slide.shapes.add_picture(
                    tmp, Inches(cell_l), Inches(cell_t),
                    width=Inches(cell_w - 0.05), height=Inches(cell_h - 0.05)
                )
                os.unlink(tmp)

            # Legend strip at bottom
            leg_tb = slide.shapes.add_textbox(
                Inches(0.27), Inches(SH - LOGO_H - 0.55), Inches(SW - 0.54), Inches(0.35)
            )
            leg_tf = leg_tb.text_frame
            leg_p = leg_tf.paragraphs[0]
            leg_run = leg_p.add_run()
            leg_run.text = (
                "\u25cf Overheating rays (hits facade)  "
                "\u25cf Overheating (other side)  "
                "- - Cutoff arc (VSA cut-off)  "
                "\u2014 Facade direction"
            )
            leg_run.font.size = Pt(8)
            leg_run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _shading_masks_slide()

    # ═══════════════════════════════════════════════════════════
    # ANNEXURE SLIDE – About EDS, Disclaimer, Acknowledgement
    # ═══════════════════════════════════════════════════════════
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Annexure")
        _divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        # About EDS
        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "Environmental Design Solutions [EDS] is a sustainability advisory firm. Since 2002, EDS has worked on over 500 green building and energy efficiency projects worldwide. The team focuses on climate change mitigation, low-carbon design, building simulation, performance audits, and capacity building. EDS continues to contribute to the buildings community with useful tools through its IT services."
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)
        p.level = 0

        # Disclaimer
        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        disclaimer_items = [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "\u2022  EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "\u2022  EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "\u2022  Open-source resources used: Clima - Berkley, Streamlit, Python",
            "\u2022  EDS is not liable to inform Users about updates to the Application or underlying resources"
        ]

        for item in disclaimer_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            p.level = 0

        # Acknowledgement
        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        ack_items = [
            "\u2022  Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "\u2022  Streamlit, \u00a9 Streamlit Inc., licensed under Apache 2.0",
            "\u2022  Python \u00a9 Python Software Foundation, licensed under PSF License Version 2"
        ]

        for item in ack_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)
            p.level = 0

        _add_logo(slide)

    _make_annexure_slide()

    # ── Save ──────────────────────────────────────────────────
    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes


def convert_epw_timezone(tz_offset):
    """Convert EPW numeric timezone to valid pytz timezone string."""
    # Common mappings for EPW numeric timezone offsets
    tz_map = {
        5.5: "Asia/Kolkata",
        0: "UTC",
        -5: "Etc/GMT+5",
        -6: "Etc/GMT+6",
        -7: "Etc/GMT+7",
        -8: "Etc/GMT+8",
        1: "Europe/London",
        2: "Europe/Paris",
    }
    try:
        tz_float = float(tz_offset)
        if tz_float in tz_map:
            return tz_map[tz_float]
    except (ValueError, TypeError):
        pass
    return "UTC"


def parse_epw(epw_text: str) -> tuple:
    """Parse EPW formatted text and return a tuple of (DataFrame, metadata).
    
    Returns:
        tuple: (df, metadata) where df has datetime, dry_bulb_temperature, relative_humidity, hour
               and metadata contains latitude, longitude, timezone
    """
    # split into lines and find first data row (starts with a year integer)
    lines = [ln.strip() for ln in epw_text.splitlines() if ln.strip() != ""]
    # Extract metadata from header (first line)
    metadata = {"latitude": None, "longitude": None, "timezone": "UTC", "city": None, "location": None}
    if len(lines) > 0:
        header = lines[0].split(",")
        try:
            # EPW header format: LOCATION,CITY,STATE,COUNTRY,DATA SOURCE,WMO STATION #,LATITUDE,LONGITUDE,TIMEZONE,ELEVATION
            # Column indices:      0        1     2     3       4           5              6         7          8        9
            if len(header) >= 2:
                metadata["location"] = header[0].strip()
                metadata["city"] = header[1].strip()
            if len(header) >= 8:
                metadata["latitude"] = float(header[6].strip())
                metadata["longitude"] = float(header[7].strip())
            if len(header) >= 9:
                metadata["timezone"] = convert_epw_timezone(header[8].strip())
        except (ValueError, IndexError, TypeError):
            pass
    
    data_start = None
    for i, ln in enumerate(lines):
        # consider a line a data line if first token is a 4-digit year
        toks = ln.split(",")
        if len(toks) > 1 and re.fullmatch(r"\d{4}", toks[0].strip()):
            data_start = i
            break
    
    if data_start is None:
        raise ValueError("Could not locate EPW data rows")
    
    data_str = "\n".join(lines[data_start:])
    df_raw = pd.read_csv(io.StringIO(data_str), header=None)
    
    # EPW standard columns (0-based index):
    # 0=year,1=month,2=day,3=hour,4=minute,5=data source,6=dry bulb (C),7=dew point,8=relative humidity (%),
    # 14=direct_normal_irradiance (Wh/m²), 15=diffuse_horizontal_irradiance (Wh/m²)
    col_map = {
        "year": 0,
        "month": 1,
        "day": 2,
        "hour": 3,
        "minute": 4,
        "dry_bulb_temperature": 6,
        "relative_humidity": 8,
        "direct_normal_irradiance": 14,
        "diffuse_horizontal_irradiance": 15,
    }
    
    # safety: ensure DataFrame has enough columns
    max_needed = max(col_map.values())
    if df_raw.shape[1] <= max_needed:
        raise ValueError("EPW data appears to have insufficient columns")
    
    df = pd.DataFrame()
    df["year"] = pd.to_numeric(df_raw.iloc[:, col_map["year"]], errors="coerce").astype("Int64")
    df["month"] = pd.to_numeric(df_raw.iloc[:, col_map["month"]], errors="coerce").astype("Int64")
    df["day"] = pd.to_numeric(df_raw.iloc[:, col_map["day"]], errors="coerce").astype("Int64")
    df["hour_raw"] = pd.to_numeric(df_raw.iloc[:, col_map["hour"]], errors="coerce").astype("Int64")
    df["minute"] = pd.to_numeric(df_raw.iloc[:, col_map["minute"]], errors="coerce").astype("Int64")
    
    # EPW hours are 1-24 representing the hour ending; convert to 0-23 by subtracting 1
    df["hour"] = (df["hour_raw"].fillna(1).astype(int) - 1) % 24
    
    df["dry_bulb_temperature"] = pd.to_numeric(df_raw.iloc[:, col_map["dry_bulb_temperature"]], errors="coerce")
    df["relative_humidity"] = pd.to_numeric(df_raw.iloc[:, col_map["relative_humidity"]], errors="coerce")
    df["direct_normal_irradiance"] = pd.to_numeric(df_raw.iloc[:, col_map["direct_normal_irradiance"]], errors="coerce")
    df["diffuse_horizontal_irradiance"] = pd.to_numeric(df_raw.iloc[:, col_map["diffuse_horizontal_irradiance"]], errors="coerce").fillna(0)
    # Global Horizontal Irradiance from EPW column 13 (Wh/m²)
    df["global_horizontal_irradiance"] = pd.to_numeric(df_raw.iloc[:, 13], errors="coerce").fillna(0)
    
    # build datetime (note: this may produce NaT for invalid rows)
    df["datetime"] = pd.to_datetime(
        dict(year=df["year"], month=df["month"], day=df["day"], hour=df["hour"], minute=df["minute"]),
        errors="coerce",
    )
    
    df = df.dropna(subset=["datetime"]).reset_index(drop=True)
    return df[["datetime", "dry_bulb_temperature", "relative_humidity", "direct_normal_irradiance", "diffuse_horizontal_irradiance", "global_horizontal_irradiance", "hour"]], metadata


def plot_sun_path(data: pd.DataFrame, metadata: dict, chart_type: str = "Sun Path") -> dict:
    """Generate and display an interactive Sun Path Diagram using Plotly.
    
    Args:
        data: DataFrame with datetime and weather data
        metadata: Dictionary containing latitude, longitude, timezone
        chart_type: Type of chart to display - "Direct Normal Radiation", "Global Horizontal Radiation", "Dry Bulb Temperature", or "Sun Path"
    
    Returns:
        Dictionary containing metrics for Shading diagram, or empty dict for other chart types
    """
    try:
        from pvlib import solarposition
        import plotly.graph_objects as go
    except ImportError:
        st.error("Required packages not found. Please ensure pvlib and plotly are installed.")
        return
    
    # Extract location information from metadata
    lat = metadata.get("latitude")
    lon = metadata.get("longitude")
    tz_str = metadata.get("timezone", "UTC")
    
    if lat is None or lon is None:
        st.error("Location information (latitude/longitude) not found in EPW file.")
        return
    
    try:
        # Get timezone object using pytz
        try:
            tz = pytz.timezone(tz_str)
        except:
            # Fallback: try to parse numeric timezone
            try:
                tz_offset = float(tz_str)
                hours = int(tz_offset)
                minutes = int((tz_offset - hours) * 60)
                sign = "+" if tz_offset >= 0 else "-"
                tz_for_localize = f"UTC{sign}{abs(hours):02d}:{abs(minutes):02d}"
                tz = pytz.timezone(tz_for_localize)
            except:
                tz = pytz.UTC
        
        # Create full year of hourly times
        times = pd.date_range(
            "2020-01-01 00:00:00",
            "2021-01-01 00:00:00",
            freq="h",
            tz=tz,
            inclusive="left"
        )
        
        # Calculate solar positions for entire year
        solpos = solarposition.get_solarposition(times, lat, lon)
        
        # Filter to daytime only (elevation > 0)
        solpos = solpos[solpos["apparent_elevation"] > 0]
        
        if solpos.empty:
            st.error("No daytime solar positions found. Check timezone or location.")
            return
        
        # Convert to polar coordinates (zenith distance = 90 - elevation)
        solpos["r"] = 90 - solpos["apparent_elevation"]
        solpos["theta"] = solpos["azimuth"]
        
        # Merge with EPW data to get temperature and radiation
        # First, make EPW data timezone-aware to match solpos
        epw_data = data.set_index("datetime")
        if epw_data.index.tz is None:
            epw_data.index = epw_data.index.tz_localize(tz)
        else:
            epw_data.index = epw_data.index.tz_convert(tz)
        
        # Normalize EPW data to use 2020 as the year (to match solpos calculations)
        # This allows proper merging by month/day/hour
        epw_data.index = epw_data.index.map(lambda x: x.replace(year=2020))
        
        solpos_merged = solpos.copy()
        solpos_merged = solpos_merged.join(epw_data[["dry_bulb_temperature", "direct_normal_irradiance", "diffuse_horizontal_irradiance"]], how="left")
        
        # Calculate Global Horizontal Irradiance (GHI) from DNI and DHI using solar altitude angle
        # GHI = DNI * sin(altitude) + DHI
        altitude_rad = np.radians(solpos_merged["apparent_elevation"])
        dni = solpos_merged["direct_normal_irradiance"].fillna(0)
        dhi = solpos_merged["diffuse_horizontal_irradiance"].fillna(0)
        solpos_merged["global_horizontal_irradiance"] = (
            np.maximum(0, np.sin(altitude_rad)) * dni + dhi
        )
        
        # Create figure
        fig = go.Figure()
        
        # Determine color scale and colorbar title based on chart type
        # All charts use YlOrRd colorscale for consistency with DNR chart
        if chart_type == "Direct Normal Radiation":
            color_data = solpos_merged.get("direct_normal_irradiance", pd.Series(0, index=solpos_merged.index))
            color_data = color_data.fillna(0)
            colorscale = "YlOrRd"  # Yellow-Orange-Red
            colorbar_title = "DNR (Wh/m²)"
            colorbar_min = 0
            colorbar_max = 1000
            # Ticks: 0, 200, 400, 600, 800, 1000
            tickvals = [0, 200, 400, 600, 800, 1000]
            ticktext = ["0", "200", "400", "600", "800", "1000"]
            hover_template = (
                "<b>%{customdata[0]} %{customdata[1]}</b><br>"
                "Time: %{customdata[2]}<br>"
                "Altitude: %{customdata[3]:.1f}°<br>"
                "Azimuth: %{customdata[4]:.1f}°<br>"
                "DNR: %{customdata[5]:.0f} Wh/m²"
                "<extra></extra>"
            )
        elif chart_type == "Global Horizontal Radiation":
            color_data = solpos_merged.get("global_horizontal_irradiance", pd.Series(0, index=solpos_merged.index))
            color_data = color_data.fillna(0)
            colorscale = "YlOrRd"  # Yellow-Orange-Red
            colorbar_title = "GHR (Wh/m²)"
            colorbar_min = 0
            colorbar_max = 1200
            # Ticks: 0, 200, 400, 600, 800, 1000, 1200
            tickvals = [0, 200, 400, 600, 800, 1000, 1200]
            ticktext = ["0", "200", "400", "600", "800", "1000", "1200"]
            hover_template = (
                "<b>%{customdata[0]} %{customdata[1]}</b><br>"
                "Time: %{customdata[2]}<br>"
                "Altitude: %{customdata[3]:.1f}°<br>"
                "Azimuth: %{customdata[4]:.1f}°<br>"
                "GHR: %{customdata[5]:.0f} Wh/m²"
                "<extra></extra>"
            )
        elif chart_type == "Dry Bulb Temperature":
            color_data = solpos_merged.get("dry_bulb_temperature", pd.Series(20, index=solpos_merged.index))
            color_data = color_data.fillna(20)
            colorscale = "YlOrRd"  # Use same colorscale as DNR
            colorbar_title = "Temperature (°C)"
            colorbar_min = 5
            colorbar_max = 40
            # Ticks: 5, 10, 15, 20, 25, 30, 35, 40
            tickvals = [5, 10, 15, 20, 25, 30, 35, 40]
            ticktext = ["5", "10", "15", "20", "25", "30", "35", "40"]
            hover_template = (
                "<b>%{customdata[0]} %{customdata[1]}</b><br>"
                "Time: %{customdata[2]}<br>"
                "Altitude: %{customdata[3]:.1f}°<br>"
                "Azimuth: %{customdata[4]:.1f}°<br>"
                "Temperature: %{customdata[5]:.1f}°C"
                "<extra></extra>"
            )
        elif chart_type == "Shading":
            ghi = solpos_merged.get("global_horizontal_irradiance", pd.Series(0, index=solpos_merged.index)).fillna(0)
            temp = solpos_merged.get("dry_bulb_temperature", pd.Series(20, index=solpos_merged.index)).fillna(20)
            # Boolean: 1 = shading needed, 0 = no shading needed
            shading_needed = ((temp > 28) & (ghi > 315)).astype(int)
            color_data = shading_needed
            colorscale = [[0, "#FFF9C4"], [1, "#E65100"]]  # pale yellow -> deep orange
            colorbar_title = "Shading Need"
            colorbar_min = 0
            colorbar_max = 1
            tickvals = [0, 1]
            ticktext = ["No Shading", "Shading Required"]
            hover_template = (
                "<b>%{customdata[0]} %{customdata[1]}</b><br>"
                "Time: %{customdata[2]}<br>"
                "Altitude: %{customdata[3]:.1f}°<br>"
                "Azimuth: %{customdata[4]:.1f}°<br>"
                "Temp: %{customdata[5]:.1f}°C<br>"
                "GHI: %{customdata[6]:.0f} Wh/m²<br>"
                "Shading: %{customdata[7]}"
                "<extra></extra>"
            )
        else:  # Sun Path - color by day of year
            color_data = solpos_merged.index.dayofyear
            colorscale = "YlOrRd"  # Use same colorscale as DNR
            colorbar_title = "Day of Year"
            colorbar_min = 1
            colorbar_max = 365
            # Ticks: 1, 91, 182, 273, 365 (approx quarterly)
            tickvals = [1, 91, 182, 273, 365]
            ticktext = ["1", "91", "182", "273", "365"]
            hover_template = (
                "<b>%{customdata[0]} %{customdata[1]}</b><br>"
                "Time: %{customdata[2]}<br>"
                "Altitude: %{customdata[3]:.1f}°<br>"
                "Azimuth: %{customdata[4]:.1f}°"
                "<extra></extra>"
            )
        
        # ========================
        # Add analemma curves (one per hour of day)
        # ========================
        first_trace = True  # Track first trace to show colorbar
        for hour in range(24):
            subset_idx = solpos_merged.index.hour == hour
            subset = solpos_merged[subset_idx]
            subset_colors = color_data[subset_idx]
            
            if len(subset) == 0:
                continue
            
            # Format hover text components
            month_names = subset.index.strftime("%B")
            hour_formatted = subset.index.strftime("%I:00 %p")
            temp_values = subset.get("dry_bulb_temperature", pd.Series(np.nan, index=subset.index))
            dnr_values = subset.get("direct_normal_irradiance", pd.Series(np.nan, index=subset.index))
            ghr_values = subset.get("global_horizontal_irradiance", pd.Series(np.nan, index=subset.index))
            
            if chart_type == "Direct Normal Radiation":
                customdata = np.stack(
                    (
                        month_names,
                        subset.index.day,
                        hour_formatted,
                        subset["apparent_elevation"],
                        subset["azimuth"],
                        dnr_values.fillna(0),
                    ),
                    axis=-1,
                )
            elif chart_type == "Global Horizontal Radiation":
                customdata = np.stack(
                    (
                        month_names,
                        subset.index.day,
                        hour_formatted,
                        subset["apparent_elevation"],
                        subset["azimuth"],
                        ghr_values.fillna(0),
                    ),
                    axis=-1,
                )
            elif chart_type == "Dry Bulb Temperature":
                customdata = np.stack(
                    (
                        month_names,
                        subset.index.day,
                        hour_formatted,
                        subset["apparent_elevation"],
                        subset["azimuth"],
                        temp_values.fillna(20),
                    ),
                    axis=-1,
                )
            elif chart_type == "Shading":
                shading_temp = subset.get("dry_bulb_temperature", pd.Series(np.nan, index=subset.index)).fillna(20)
                shading_ghi = subset.get("global_horizontal_irradiance", pd.Series(np.nan, index=subset.index)).fillna(0)
                shading_labels = np.where(
                    (shading_temp > 28) & (shading_ghi > 315), "Required", "Not Required"
                )
                customdata = np.stack(
                    (
                        month_names,
                        subset.index.day,
                        hour_formatted,
                        subset["apparent_elevation"],
                        subset["azimuth"],
                        shading_temp,
                        shading_ghi,
                        shading_labels,
                    ),
                    axis=-1,
                )
            else:  # Sun Path - only position data
                customdata = np.stack(
                    (
                        month_names,
                        subset.index.day,
                        hour_formatted,
                        subset["apparent_elevation"],
                        subset["azimuth"],
                    ),
                    axis=-1,
                )
            
            fig.add_trace(
                go.Scatterpolar(
                    r=subset["r"],
                    theta=subset["theta"],
                    mode="lines+markers",
                    marker=dict(
                        size=4,
                        color=subset_colors.values,
                        colorscale=colorscale,
                        cmin=colorbar_min,
                        cmax=colorbar_max,
                        showscale=(first_trace and chart_type != "Shading"),  # Hide colorbar for Shading (binary)
                        colorbar=dict(
                            title=dict(
                                text=colorbar_title,
                                side="right",
                                font=dict(size=12, color="#333"),
                            ),
                            thickness=20,
                            len=0.7,
                            x=1.1,
                            tickvals=tickvals,
                            ticktext=ticktext,
                            tickmode="array"
                        ) if chart_type != "Shading" else None,
                        opacity=0.7,
                        line=dict(width=0.5)
                    ),
                    line=dict(width=1, color="rgba(100, 100, 100, 0.3)"),
                    showlegend=False,
                    customdata=customdata,
                    hovertemplate=hover_template,
                )
            )
            first_trace = False  # Only show colorbar on first trace
        
        # ========================
        # Shading metrics calculation (for display below chart)
        # ========================
        shading_metrics = {}
        if chart_type == "Shading":
            # Calculate metrics
            ghi_col = solpos_merged.get("global_horizontal_irradiance", pd.Series(0, index=solpos_merged.index)).fillna(0)
            temp_col = solpos_merged.get("dry_bulb_temperature", pd.Series(20, index=solpos_merged.index)).fillna(20)
            
            # Total sunshine hours (GHI > 300 Wh/m²)
            sunshine_mask = ghi_col > 300
            sunshine_hours = sunshine_mask.sum() / 2  # Each data point is 30 minutes (0.5 hours)
            
            # Required shading hours (temp > 28°C AND GHI > 315 Wh/m²)
            shading_mask = (temp_col > 28) & (ghi_col > 315)
            shading_hours = shading_mask.sum() / 2  # Each data point is 30 minutes (0.5 hours)
            
            shading_metrics = {
                "total_sunshine_hours": sunshine_hours,
                "required_shading_hours": shading_hours
            }
            
            # Add legend items for binary shading colors
            fig.add_trace(
                go.Scatterpolar(
                    r=[None],
                    theta=[None],
                    mode="markers",
                    marker=dict(size=12, color="#FFF9C4", symbol="square"),
                    name="No Shading",
                    showlegend=True,
                    hoverinfo="skip"
                )
            )
            fig.add_trace(
                go.Scatterpolar(
                    r=[None],
                    theta=[None],
                    mode="markers",
                    marker=dict(size=12, color="#E65100", symbol="square"),
                    name="Shading Required",
                    showlegend=True,
                    hoverinfo="skip"
                )
            )
        
        # ========================
        # Add special solar dates (equinoxes and solstices)
        # ========================
        key_dates = {
            "Mar 21 (Spring)": ("2020-03-21", "#FF9500"),
            "Jun 21 (Summer)": ("2020-06-21", "#FF0000"),
            "Dec 21 (Winter)": ("2020-12-21", "#0066CC"),
        }
        
        for label, (date_str, color) in key_dates.items():
            date = pd.Timestamp(date_str)
            day_times = pd.date_range(
                date,
                date + pd.Timedelta("1D"),
                freq="5min",
                tz=tz,
                inclusive="left"
            )
            
            sol = solarposition.get_solarposition(day_times, lat, lon)
            sol = sol[sol["apparent_elevation"] > 0]
            
            if sol.empty:
                continue
            
            # Add EPW data to special date curves
            sol_merged = sol.copy()
            sol_merged = sol_merged.join(epw_data[["dry_bulb_temperature", "direct_normal_irradiance", "diffuse_horizontal_irradiance"]], how="left")
            
            # Calculate Global Horizontal Irradiance for special dates
            altitude_rad = np.radians(sol_merged["apparent_elevation"])
            dni = sol_merged["direct_normal_irradiance"].fillna(0)
            dhi = sol_merged["diffuse_horizontal_irradiance"].fillna(0)
            sol_merged["global_horizontal_irradiance"] = (
                np.maximum(0, np.sin(altitude_rad)) * dni + dhi
            )
            
            r = 90 - sol["apparent_elevation"]
            theta = sol["azimuth"]
            # draw the key date curve
            fig.add_trace(
                go.Scatterpolar(
                    r=r,
                    theta=theta,
                    mode="lines",
                    line=dict(width=0.8, color=color),
                    name=label,
                    showlegend=True
                )
            )

            # =========================
            # Add hour labels on Summer curve
            # =========================
            if label == "Jun 21 (Summer)":

                hour_labels = [6, 9, 12, 15, 18]

                label_times = [
                    pd.Timestamp(date_str).tz_localize(tz) + pd.Timedelta(hours=h)
                    for h in hour_labels
                ]

                label_sol = solarposition.get_solarposition(label_times, lat, lon)
                label_sol = label_sol[label_sol["apparent_elevation"] > 0]

                if not label_sol.empty:

                    r_labels = 90 - label_sol["apparent_elevation"]
                    theta_labels = label_sol["azimuth"]

                    fig.add_trace(
                        go.Scatterpolar(
                            r=r_labels,
                            theta=theta_labels,
                            mode="text",
                            text=[str(h) for h in hour_labels],
                            textfont=dict(size=11, color="black"),
                            showlegend=False,
                            hoverinfo="skip"
                        )
                    )
            temp = sol_merged.get("dry_bulb_temperature", pd.Series(np.nan, index=sol_merged.index))
            dnr = sol_merged.get("direct_normal_irradiance", pd.Series(np.nan, index=sol_merged.index))
            
            # Build customdata and hovertemplate based on chart type
            ghr = sol_merged.get("global_horizontal_irradiance", pd.Series(np.nan, index=sol_merged.index))
            if chart_type == "Direct Normal Radiation":
                special_customdata = np.stack((sol["apparent_elevation"], dnr.fillna(0)), axis=-1)
                special_hovertemplate = (
                    f"<b>{label}</b><br>"
                    "Altitude: %{customdata[0]:.1f}°<br>"
                    "Azimuth: %{theta:.1f}°<br>"
                    "DNR: %{customdata[1]:.0f} Wh/m²"
                    "<extra></extra>"
                )
            elif chart_type == "Global Horizontal Radiation":
                special_customdata = np.stack((sol["apparent_elevation"], ghr.fillna(0)), axis=-1)
                special_hovertemplate = (
                    f"<b>{label}</b><br>"
                    "Altitude: %{customdata[0]:.1f}°<br>"
                    "Azimuth: %{theta:.1f}°<br>"
                    "GHR: %{customdata[1]:.0f} Wh/m²"
                    "<extra></extra>"
                )
            elif chart_type == "Dry Bulb Temperature":
                special_customdata = np.stack((sol["apparent_elevation"], temp.fillna(20)), axis=-1)
                special_hovertemplate = (
                    f"<b>{label}</b><br>"
                    "Altitude: %{customdata[0]:.1f}°<br>"
                    "Azimuth: %{theta:.1f}°<br>"
                    "Temperature: %{customdata[1]:.1f}°C"
                    "<extra></extra>"
                )
            elif chart_type == "Shading":
                s_ghi = sol_merged.get("global_horizontal_irradiance", pd.Series(np.nan, index=sol_merged.index)).fillna(0)
                special_customdata = np.stack((sol["apparent_elevation"], temp.fillna(20), s_ghi), axis=-1)
                special_hovertemplate = (
                    f"<b>{label}</b><br>"
                    "Altitude: %{customdata[0]:.1f}°<br>"
                    "Azimuth: %{theta:.1f}°<br>"
                    "Temp: %{customdata[1]:.1f}°C<br>"
                    "GHI: %{customdata[2]:.0f} Wh/m²"
                    "<extra></extra>"
                )
            else:  # Sun Path
                special_customdata = np.stack((sol["apparent_elevation"],), axis=-1)
                special_hovertemplate = (
                    f"<b>{label}</b><br>"
                    "Altitude: %{customdata[0]:.1f}°<br>"
                    "Azimuth: %{theta:.1f}°"
                    "<extra></extra>"
                )
            
            # fig.add_trace(
            #     go.Scatterpolar(
            #         r=r,
            #         theta=theta,
            #         mode="lines",
            #         line=dict(width=3, color=color),
            #         name=label,
            #         showlegend=True,
            #         hovertemplate=special_hovertemplate,
            #         customdata=special_customdata,
            #     )
            # )
        
        # ========================
        # Configure polar layout
        # ========================
        fig.update_layout(
            # title={
            #     'text': f'☀️ Sun Path - {chart_type}',
            #     'font': {'size': 24, 'color': '#8B4513'},
            #     'x': 0.5,
            #     'xanchor': 'center'
            # },
            polar=dict(
                bgcolor="rgba(240, 240, 240, 0.3)",
                radialaxis=dict(
                    visible=True,
                    range=[0, 90],
                    showticklabels=False,
                    showline=False,
                    ticks="", 
                    # ticks="outside",
                    # tickfont=dict(size=10),
                    gridcolor="rgba(128, 128, 128, 0.2)",
                    # tickvals=[0, 15, 30, 45, 60, 75, 90],
                    # ticktext=["90° (Zenith)", "75°", "60°", "45°", "30°", "15°", "0° (Horizon)"],
                ),
                angularaxis=dict(
                    tickfont=dict(size=11),
                    rotation=90,
                    direction="clockwise",
                    gridcolor="rgba(128, 128, 128, 0.3)",
                    tickvals=[0, 45, 90, 135, 180, 225, 270, 315],
                    ticktext=["N", "NE", "E", "SE", "S", "SW", "W", "NW"],
                ),
            ),
            
            showlegend=True,
            legend=dict(
                x=0.02,
                y=0.98,
                bgcolor="rgba(255, 255, 255, 0.9)",
                bordercolor="black",
                borderwidth=1,
                font=dict(size=10)
            ),
            hovermode="closest",
            height=700,
            margin=dict(l=80, r=140, t=100, b=80),
            plot_bgcolor="white",
            paper_bgcolor="white",
            font=dict(family="Arial, sans-serif", size=12, color="black")
        )
        fig.add_trace(
        go.Scatterpolar(
            r=[0, 15, 30, 45, 60, 75, 90],
            theta=[45]*7,   # NE direction
            mode="lines+text",
            line=dict(color="black", width=1),
            text=[
                "90° (Zenith)",
                "75°",
                "60°",
                "45°",
                "30°",
                "15°",
                "0° (Horizon)"
            ],
            textposition="middle right",
            textfont=dict(size=10),
            showlegend=False,
            hoverinfo="skip"
        )
    )
        
        # Display with Streamlit
        st.plotly_chart(fig, use_container_width=True)
        
        # Return metrics if available
        return shading_metrics if chart_type == "Shading" else {}
        
    except Exception as e:
        st.error(f"Error generating Sun Path Diagram: {str(e)}")
        return {}

# =====================================================================================
# SHADING ANALYSIS HELPER FUNCTIONS
# =====================================================================================

_MONTH_SHORT = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]

_ORIENTATIONS = {
    "North (0°)": 0,
    "North-East (45°)": 45,
    "East (90°)": 90,
    "South-East (135°)": 135,
    "South (180°)": 180,
    "South-West (225°)": 225,
    "West (270°)": 270,
    "North-West (315°)": 315,
}


def build_thermal_matrix(epw_df: pd.DataFrame, temp_threshold: float, rad_threshold: float):
    """Build a 24×12 matrix of mean temperature and GHI per (hour, month).

    Returns:
        temp_matrix:   24×12 DataFrame of mean dry bulb temperature
        rad_matrix:    24×12 DataFrame of mean GHI
        overheat_mask: bool 24×12 DataFrame where both thresholds are exceeded
    """
    d = epw_df.copy()
    d["_month"] = d["datetime"].dt.month
    d["_hour"] = d["hour"].astype(int)

    temp_pivot = (
        d.groupby(["_hour", "_month"])["dry_bulb_temperature"]
        .mean()
        .unstack("_month")
    )
    rad_pivot = (
        d.groupby(["_hour", "_month"])["global_horizontal_irradiance"]
        .mean()
        .unstack("_month")
    )

    temp_pivot.columns = _MONTH_SHORT
    rad_pivot.columns = _MONTH_SHORT
    temp_pivot.index.name = "Hour"
    rad_pivot.index.name = "Hour"

    overheat_mask = (temp_pivot > temp_threshold) & (rad_pivot > rad_threshold)
    return temp_pivot, rad_pivot, overheat_mask


def get_overheating_hours(epw_df: pd.DataFrame, temp_threshold: float, rad_threshold: float) -> pd.DataFrame:
    """Return EPW rows where temperature > temp_threshold AND GHI > rad_threshold."""
    mask = (
        (epw_df["dry_bulb_temperature"] > temp_threshold) &
        (epw_df["global_horizontal_irradiance"] > rad_threshold)
    )
    return epw_df[mask].copy().reset_index(drop=True)


def compute_solar_angles(overheat_df: pd.DataFrame, lat: float, lon: float, tz_str: str) -> pd.DataFrame:
    """Compute solar altitude and azimuth for each overheating timestamp using pvlib.

    Returns overheat_df rows above the horizon with solar_altitude and solar_azimuth columns.
    """
    from pvlib import solarposition

    try:
        tz = pytz.timezone(tz_str)
    except Exception:
        tz = pytz.UTC

    times = pd.DatetimeIndex(overheat_df["datetime"].values)
    if times.tz is None:
        times = times.tz_localize(tz)

    # Normalise to year 2020 for consistency with plot_sun_path
    times_2020 = times.map(lambda t: t.replace(year=2020))
    solpos = solarposition.get_solarposition(times_2020, lat, lon)

    result = overheat_df.copy()
    result["solar_altitude"] = solpos["apparent_elevation"].values
    result["solar_azimuth"] = solpos["azimuth"].values

    # Keep only sun-above-horizon positions
    result = result[result["solar_altitude"] > 0].copy().reset_index(drop=True)
    return result


def compute_shading_geometry(solar_positions: pd.DataFrame, facade_azimuth: float) -> pd.DataFrame:
    """Compute VSA, HSA and facade-hit flag for a given facade azimuth.

    Formulas:
        relative_azimuth = solar_azimuth − facade_azimuth
        VSA = arctan( tan(altitude) / cos(relative_azimuth) )
        HSA = arctan( sin(relative_azimuth) / tan(altitude) )
    """
    alt_rad = np.radians(solar_positions["solar_altitude"].values)
    rel_az = solar_positions["solar_azimuth"].values - facade_azimuth
    rel_az = ((rel_az + 180) % 360) - 180          # normalise to −180…+180
    rel_az_rad = np.radians(rel_az)

    hits_facade = np.abs(rel_az) < 90

    tan_alt = np.tan(alt_rad)
    cos_rel = np.cos(rel_az_rad)
    sin_rel = np.sin(rel_az_rad)

    cos_rel_safe = np.where(np.abs(cos_rel) < 1e-9, np.sign(cos_rel + 1e-18) * 1e-9, cos_rel)
    tan_alt_safe = np.where(np.abs(tan_alt) < 1e-9, 1e-9, tan_alt)

    vsa_deg = np.degrees(np.arctan(tan_alt / cos_rel_safe))
    hsa_deg = np.degrees(np.arctan(sin_rel / tan_alt_safe))

    result = solar_positions.copy()
    result["relative_azimuth"] = rel_az
    result["VSA"] = vsa_deg
    result["HSA"] = hsa_deg
    result["hits_facade"] = hits_facade
    return result


def build_orientation_table(solar_positions: pd.DataFrame, design_cutoff_angle: float) -> pd.DataFrame:
    """Build the 8-orientation shading analysis table."""
    rows = []
    for orient_name, facade_az in _ORIENTATIONS.items():
        geom = compute_shading_geometry(solar_positions, facade_az)
        facing = geom[geom["hits_facade"]]

        if facing.empty:
            rows.append({
                "Orientation": orient_name,
                "Rays Hitting": 0,
                "Min VSA (°)": None,
                "Max |HSA| (°)": None,
                "D/H (Overhang)": None,
                "D/W (Fin)": None,
                "Protection (%)": 100.0,
            })
            continue

        min_vsa = float(facing["VSA"].min())
        max_abs_hsa = float(facing["HSA"].abs().max())

        # D/H = 1/tan(Min VSA)  — min VSA is the hardest ray to block (shallowest sun)
        # Capped at 1.0 (overhang deeper than window height is impractical to display)
        tan_min_vsa = np.tan(np.radians(min_vsa))
        dh = round(min(1.0, float(1.0 / tan_min_vsa)) if abs(tan_min_vsa) > 1e-6 else 1.0, 3)

        # D/W = tan(Max |HSA|)  — max HSA is the hardest lateral ray to block
        # Capped at 1.0 (fin deeper than window width is impractical to display)
        dw = round(min(1.0, float(np.tan(np.radians(max_abs_hsa)))), 3)

        # Protection: rays with VSA >= design_cutoff_angle are blocked by the overhang
        # (steeper sun = shallower shadow angle = overhang blocks it)
        blocked = int((facing["VSA"] >= design_cutoff_angle).sum())
        protection_pct = (blocked / len(facing)) * 100

        rows.append({
            "Orientation": orient_name,
            "Rays Hitting": int(len(facing)),
            "Min VSA (°)": round(min_vsa, 1),
            "Max |HSA| (°)": round(max_abs_hsa, 1),
            "D/H (Overhang)": dh,
            "D/W (Fin)": dw,
            "Protection (%)": round(protection_pct, 1),
        })

    return pd.DataFrame(rows)


def make_shading_mask_chart(solar_positions: pd.DataFrame, facade_azimuth: float,
                             design_cutoff_angle: float):
    """Create a small polar chart showing overheating sun positions and shading cutoff arc."""
    import plotly.graph_objects as go

    geom = compute_shading_geometry(solar_positions, facade_azimuth)
    facing = geom[geom["hits_facade"]]
    other = geom[~geom["hits_facade"]]

    fig = go.Figure()

    # Non-facing overheating positions (grey background context)
    if not other.empty:
        fig.add_trace(go.Scatterpolar(
            r=(90 - other["solar_altitude"]).tolist(),
            theta=other["solar_azimuth"].tolist(),
            mode="markers",
            marker=dict(size=3, color="lightgrey", opacity=0.5),
            showlegend=False,
            hoverinfo="skip",
        ))

    # Facing overheating positions (deep orange / red)
    if not facing.empty:
        fig.add_trace(go.Scatterpolar(
            r=(90 - facing["solar_altitude"]).tolist(),
            theta=facing["solar_azimuth"].tolist(),
            mode="markers",
            marker=dict(size=4, color="#E65100", opacity=0.75),
            showlegend=False,
            hovertemplate=(
                "Alt: %{text}<br>Az: %{theta:.0f}°<extra></extra>"
            ),
            text=[f"{a:.1f}°" for a in facing["solar_altitude"]],
        ))

    # Shading cutoff arc (VSA = design_cutoff_angle contour)
    rel_az_range = np.linspace(-89, 89, 179)
    tan_cutoff = np.tan(np.radians(design_cutoff_angle))
    cutoff_alt = np.degrees(np.arctan(tan_cutoff * np.cos(np.radians(rel_az_range))))
    cutoff_az_abs = facade_azimuth + rel_az_range
    valid = cutoff_alt > 0
    if valid.any():
        fig.add_trace(go.Scatterpolar(
            r=(90 - cutoff_alt[valid]).tolist(),
            theta=cutoff_az_abs[valid].tolist(),
            mode="lines",
            line=dict(color="#1565C0", width=2, dash="dash"),
            showlegend=False,
            hoverinfo="skip",
        ))

    # Facade direction indicator (green radial line)
    fig.add_trace(go.Scatterpolar(
        r=[0, 85],
        theta=[facade_azimuth, facade_azimuth],
        mode="lines",
        line=dict(color="#388E3C", width=2),
        showlegend=False,
        hoverinfo="skip",
    ))

    fig.update_layout(
        polar=dict(
            bgcolor="rgba(240, 248, 255, 0.6)",
            radialaxis=dict(visible=False, range=[0, 90]),
            angularaxis=dict(
                tickfont=dict(size=7),
                rotation=90,
                direction="clockwise",
                tickvals=[0, 90, 180, 270],
                ticktext=["N", "E", "S", "W"],
            ),
        ),
        showlegend=False,
        height=180,
        margin=dict(l=10, r=10, t=10, b=10),
        paper_bgcolor="rgba(0,0,0,0)",
    )
    return fig


# === MAIN LAYOUT ===
col_left, col_right = st.columns([0.85, 2.15], gap="small")

with col_left:
    # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
    
    # Upload EPW File Section
    # st.markdown('<div class="control-section-header">📤 Upload EPW File</div>', unsafe_allow_html=True)
    # st.markdown('<div class="upload-zone">Limit 200MB per file · EPW</div>', unsafe_allow_html=True)
    st.write("##### 📤 Upload EPW File")
    uploaded = st.file_uploader("", type=["epw"], label_visibility="collapsed", width=300)
    
    # Parameter Selection
    # st.markdown('<div class="control-section-header">⚙️ Parameter</div>', unsafe_allow_html=True)
    st.write("##### ⚙️ Parameter")
    selected_parameter = st.selectbox(
        "Select parameter",
        ["Temperature", "Humidity", "Sun Path"],
        label_visibility="collapsed",
        key="parameter_selector",
        width=300
    )

if uploaded is None:
    with col_left:
        st.info("Please upload an .epw file to analyze.", width=300)
    st.stop()

try:
    raw = uploaded.getvalue().decode("utf-8", errors="replace")
    df, metadata = parse_epw(raw)
    
    # compute derived date fields used by the charts
    df["doy"] = df["datetime"].dt.dayofyear
    df["day"] = df["datetime"].dt.day
    df["month"] = df["datetime"].dt.month
    df["month_name"] = df["datetime"].dt.strftime("%b")
    
    with col_left:
        # st.markdown("""
        #     <div style="
        #         background-color: #f0fff4;
        #         border-left: 4px solid #48bb78;
        #         padding: 12px;
        #         border-radius: 4px;
        #         margin: 8px 0;
        #         width: 300px;
        #     ">
        #         <div style="color: #22543d; font-weight: 600; font-size: 12px;">✅ EPW parsed successfully</div>
        #     </div>
        # """, unsafe_allow_html=True)
        
        # Time range (hour of use) slider for diurnal/peak analysis
        if selected_parameter == "Sun Path":
            # For Sun Path, always use full 0-23 range
            hour_range = (0, 23)

            # ── Shading Analysis Parameters ─────────────────────────────────────
            st.markdown('<div class="control-section-header">🌡️ Overheating Thresholds</div>', unsafe_allow_html=True)
            temp_threshold = st.number_input(
                "Temperature Threshold (°C)",
                value=28.0,
                step=0.5,
                key="temp_threshold",
                help="Hours with dry-bulb temperature above this value are flagged",
                width=300,
            )
            rad_threshold = st.number_input(
                "Radiation Threshold (W/m²)",
                value=315.0,
                step=10.0,
                key="rad_threshold",
                help="Hours with Global Horizontal Irradiance above this value are flagged",
                width=300,
            )

            st.markdown('<div class="control-section-header">📍 Location</div>', unsafe_allow_html=True)
            _lat_default = float(metadata.get("latitude") or 0.0)
            _lon_default = float(metadata.get("longitude") or 0.0)
            shading_lat = st.number_input(
                "Latitude", value=_lat_default, step=0.1, key="shading_lat",
                help="Auto-read from EPW metadata",
                width=300
            )
            shading_lon = st.number_input(
                "Longitude", value=_lon_default, step=0.1, key="shading_lon",
                help="Auto-read from EPW metadata",
                width=300
            )

            

            st.markdown('<div class="control-section-header">📐 Design Parameters</div>', unsafe_allow_html=True)
            design_cutoff_angle = st.number_input(
                "Design Cutoff Angle (°)",
                value=45.0,
                step=1.0,
                min_value=5.0,
                max_value=89.0,
                key="design_cutoff_angle",
                help="Solar altitude cutoff used to assess shading device protection",
                width=300
            )
        else:
            st.markdown('<div class="control-section-header">⏰ Time Range (Hours)</div>', unsafe_allow_html=True)
            hour_range = st.slider(
                "Select hours (start - end)",
                min_value=0,
                max_value=23,
                value=(8, 18),
                step=1,
                key="hour_range",
                label_visibility="collapsed",
                width=300
            )
        
        # Date range selection
        st.markdown('<div class="control-section-header">📅 Date Range</div>', unsafe_allow_html=True)
        
        months_list = ["January", "February", "March", "April", "May", "June", 
                       "July", "August", "September", "October", "November", "December"]
        
        # Initialize session state for month selection
        if "start_month_idx" not in st.session_state:
            st.session_state.start_month_idx = 0
        if "end_month_idx" not in st.session_state:
            st.session_state.end_month_idx = 11
        
        # Create two columns for start and end month dropdowns
        month_col1, month_col2, col3 = st.columns([1,1,0.5], gap="small")
        
        with month_col1:
            start_month = st.selectbox(
                "From",
                options=range(len(months_list)),
                format_func=lambda x: months_list[x],
                key="start_month_select",
                label_visibility="collapsed",
                width=150
            )
            st.session_state.start_month_idx = start_month
        
        with month_col2:
            # End month should only show months from start_month onwards
            end_month_options = list(range(start_month, len(months_list)))
            # Ensure selected end_month is valid
            if st.session_state.end_month_idx < start_month:
                st.session_state.end_month_idx = start_month
            
            end_month = st.selectbox(
                "To",
                options=end_month_options,
                format_func=lambda x: months_list[x],
                key="end_month_select",
                index=min(st.session_state.end_month_idx - start_month, len(end_month_options) - 1),
                label_visibility="collapsed",
                width=150,
            )
            st.session_state.end_month_idx = end_month
        
        with col3:
            pass
        
        # ── Report (PowerPoint) ──────────────────────────────────
        st.markdown('<div class="control-section-header">📊 Report (PowerPoint)</div>', unsafe_allow_html=True)

        _active_chart = st.session_state.get("sun_chart_type", "Sun Path")
        _is_shading_mode = (selected_parameter == "Sun Path" and _active_chart == "Shading")

        try:
            start_month_num = st.session_state.start_month_idx + 1
            end_month_num   = st.session_state.end_month_idx + 1
            year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024

            start_date = pd.to_datetime(f"{year}-{start_month_num}-01").date()
            if end_month_num == 12:
                end_date = pd.to_datetime(f"{year}-12-31").date()
            else:
                end_date = (pd.to_datetime(f"{year}-{end_month_num+1}-01") - pd.Timedelta(days=1)).date()

            start_hour, end_hour = st.session_state.get("hour_range", (8, 18))

            if _is_shading_mode:
                # ── Shading Analysis Report ──────────────────────────
                _temp_thr = float(st.session_state.get("temp_threshold", 28.0))
                _rad_thr  = float(st.session_state.get("rad_threshold", 315.0))
                _sh_lat   = float(st.session_state.get("shading_lat",  metadata.get("latitude")  or 0.0))
                _sh_lon   = float(st.session_state.get("shading_lon",  metadata.get("longitude") or 0.0))
                _cutoff   = float(st.session_state.get("design_cutoff_angle", 45.0))
                _tz_str   = metadata.get("timezone", "UTC")

                shading_report = generate_shading_pptx_report(
                    df,
                    metadata,
                    temp_threshold=_temp_thr,
                    rad_threshold=_rad_thr,
                    lat=_sh_lat,
                    lon=_sh_lon,
                    tz_str=_tz_str,
                    design_cutoff_angle=_cutoff,
                )
                st.download_button(
                    label="⬇️ Download Shading Report",
                    data=shading_report,
                    file_name="Shading_Analysis_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_shading_report",
                    width=300,
                )
            else:
                # ── Climate Analysis Report ──────────────────────────
                report_bytes = generate_pptx_report(
                    df,
                    start_date,
                    end_date,
                    start_hour,
                    end_hour,
                    selected_parameter,
                    metadata=metadata,
                )
                st.download_button(
                    label="⬇️ Download Climate Report",
                    data=report_bytes,
                    file_name=f"Climate_Analysis_Report_{start_date.strftime('%Y%m%d')}_to_{end_date.strftime('%Y%m%d')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_report",
                    width=300,
                )

        except Exception as e:
            st.error(f"\u274c Failed to generate report: {str(e)}")
        
except Exception as e:
    with col_left:
        st.error(f"❌ Failed to parse EPW: {e}")
    st.stop()

def calculate_ashrae_comfort(df: pd.DataFrame) -> tuple:
    """
    Calculate ASHRAE adaptive comfort bands.
    Returns (comfort_80_lower, comfort_80_upper, comfort_90_lower, comfort_90_upper) as daily rolling averages.
    """
    # Simple ASHRAE adaptive comfort model (simplified)
    # For illustration: ranges relative to outdoor mean monthly temperature
    daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
    
    # 80% acceptability: ±3.5°C from comfort line
    # 90% acceptability: ±2.5°C from comfort line
    comfort_line = daily_avg.rolling(window=7, center=True).mean()
    
    comfort_80_lower = comfort_line - 3.5
    comfort_80_upper = comfort_line + 3.5
    comfort_90_lower = comfort_line - 2.5
    comfort_90_upper = comfort_line + 2.5
    
    return comfort_80_lower, comfort_80_upper, comfort_90_lower, comfort_90_upper

# Compute ASHRAE comfort bands
ashrae_80_lower, ashrae_80_upper, ashrae_90_lower, ashrae_90_upper = calculate_ashrae_comfort(df)

with col_right:
    # === INTERACTIVE TABS AT TOP ===
    st.markdown("""
        <style>
        .tab-container {
            display: flex;
            gap: 0;
            background-color: #f8f9fa;
            padding: 0;
            margin: -1rem -1rem 1.5rem -1rem;
            border-bottom: 2px solid #e9ecef;
        }
        .tab-button {
            padding: 12px 24px;
            cursor: pointer;
            font-size: 14px;
            font-weight: 600;
            color: #495057;
            background-color: #f8f9fa;
            border: none;
            border-bottom: 3px solid transparent;
            transition: all 0.3s ease;
        }
        .tab-button:hover {
            background-color: #e9ecef;
        }
        .tab-button.active {
            color: #2c3e50;
            border-bottom-color: #3498db;
            background-color: white;
        }
        </style>
    """, unsafe_allow_html=True)
    
    # Check if Sun Path is selected - show it directly without tabs
    if selected_parameter == "Sun Path":
        st.markdown("<h3 style='text-align: left; margin-top: 20px;'>Sun Path Analysis</h3>", unsafe_allow_html=True)
        
        # Add dropdown to select chart type
        col_selector = st.columns([1, 4])
        with col_selector[0]:
            chart_type = st.selectbox(
                "Chart Type:",
                ["Sun Path","Dry Bulb Temperature", "Direct Normal Radiation", "Global Horizontal Radiation","Shading"],
                key="sun_chart_type",
            )
        
        metrics = plot_sun_path(df, metadata, chart_type)
        
        # Display metrics for Shading chart
        if chart_type == "Shading" and metrics:
            st.markdown("---")
            st.markdown("<h4 style='text-align: left;'>📊 Shading Metrics</h4>", unsafe_allow_html=True)
            metric_col1, metric_col2 = st.columns(2)
            
            with metric_col1:
                st.metric(
                    label="Total Sunshine Hours",
                    value=f"{metrics.get('total_sunshine_hours', 0):.1f}",
                    help="Hours with Global Horizontal Irradiance > 300 Wh/m²"
                )
            
            with metric_col2:
                st.metric(
                    label="Required Shading Hours",
                    value=f"{metrics.get('required_shading_hours', 0):.1f}",
                    help="Hours where Temperature > 28°C AND GHI > 315 Wh/m²"
                )

        # =====================================================================
        # EXTENDED SHADING ANALYSIS – shown when "Shading" chart type is active
        # =====================================================================
        if chart_type == "Shading":
            import plotly.graph_objects as go

            # Read threshold/location values set by the sidebar inputs
            _temp_thr = float(st.session_state.get("temp_threshold", 28.0))
            _rad_thr = float(st.session_state.get("rad_threshold", 315.0))
            _lat = float(st.session_state.get("shading_lat", metadata.get("latitude") or 0.0))
            _lon = float(st.session_state.get("shading_lon", metadata.get("longitude") or 0.0))
            _cutoff = float(st.session_state.get("design_cutoff_angle", 45.0))
            _tz_str = metadata.get("timezone", "UTC")

            # ── 1. THERMAL & RADIATION MATRIX ─────────────────────────────────
            st.markdown("---")
            st.markdown(
                "<h4 style='margin-bottom:4px;'>Thermal &amp; Radiation Matrix"
                f"<span style='font-size:13px; font-weight:400; color:#666; margin-left:12px;'>"
                f"Black border = Temp &gt; {_temp_thr}°C AND GHI &gt; {_rad_thr} W/m²</span></h4>",
                unsafe_allow_html=True,
            )

            try:
                temp_matrix, rad_matrix, overheat_mask = build_thermal_matrix(df, _temp_thr, _rad_thr)
                hours_labels = [f"{h:02d}:00" for h in range(24)]

                # Build hover text
                hover_text = []
                for h_idx in range(24):
                    row_txt = []
                    for m_idx in range(12):
                        t_v = temp_matrix.iloc[h_idx, m_idx]
                        r_v = rad_matrix.iloc[h_idx, m_idx]
                        warn = " ⚠️ OVERHEATING" if overheat_mask.iloc[h_idx, m_idx] else ""
                        row_txt.append(
                            f"<b>{hours_labels[h_idx]}, {_MONTH_SHORT[m_idx]}</b><br>"
                            f"Avg Temp: {t_v:.1f}°C<br>"
                            f"Avg GHI: {r_v:.0f} W/m²{warn}"
                        )
                    hover_text.append(row_txt)

                fig_matrix = go.Figure(go.Heatmap(
                    z=temp_matrix.values,
                    x=_MONTH_SHORT,
                    y=hours_labels,
                    colorscale="RdYlBu_r",
                    colorbar=dict(
                        title=dict(text="°C", side="right"),
                        thickness=14,
                        len=0.8,
                    ),
                    text=hover_text,
                    hovertemplate="%{text}<extra></extra>",
                    showscale=True,
                ))

                # Overlay orange outlines on overheating cells
                shapes = []
                for h_idx in range(24):
                    for m_idx in range(12):
                        if overheat_mask.iloc[h_idx, m_idx]:
                            shapes.append(dict(
                                type="rect",
                                xref="x", yref="y",
                                x0=m_idx - 0.5, x1=m_idx + 0.5,
                                y0=h_idx - 0.5, y1=h_idx + 0.5,
                                line=dict(color="#080808", width=2.5),
                                fillcolor="rgba(0,0,0,0)",
                            ))

                n_overheat = int(overheat_mask.values.sum())
                fig_matrix.update_layout(
                    xaxis_title="Month",
                    yaxis_title="Hour",
                    height=540,
                    shapes=shapes,
                    template="plotly_white",
                    margin=dict(l=70, r=90, t=20, b=50),
                    annotations=[dict(
                        x=1.13, y=0.5, xref="paper", yref="paper",
                        text=f"<b>{n_overheat}</b><br>cells<br>overheat",
                        showarrow=False,
                        font=dict(size=11, color="#E65100"),
                        align="center",
                    )],
                )

                st.plotly_chart(fig_matrix, use_container_width=True)

            except Exception as _e:
                st.error(f"Could not build Thermal Matrix: {_e}")

            # ── 2. ORIENTATION SHADING ANALYSIS ───────────────────────────────
            st.markdown("---")
            st.markdown(
                "<h4 style='margin-bottom:4px;'>🏗️ Orientation Shading Analysis"
                f"<span style='font-size:13px; font-weight:400; color:#666; margin-left:12px;'>"
                f"Design cutoff angle: {_cutoff}°</span></h4>",
                unsafe_allow_html=True,
            )

            try:
                overheat_df_sa = get_overheating_hours(df, _temp_thr, _rad_thr)
                if overheat_df_sa.empty:
                    st.info("No overheating hours found with current thresholds.")
                else:
                    with st.spinner("Computing solar positions for overheating hours…"):
                        solar_pos = compute_solar_angles(overheat_df_sa, _lat, _lon, _tz_str)

                    if solar_pos.empty:
                        st.info("No daytime overheating sun positions found.")
                    else:
                        n_total = len(solar_pos)
                        st.caption(
                            f"**{n_total}** overheating daytime sun positions "
                            f"({len(overheat_df_sa)} EPW rows → {n_total} above horizon)"
                        )

                        orient_df = build_orientation_table(solar_pos, _cutoff)

                        # Protection percentage coloured badge
                        def _protection_badge(pct):
                            if pct is None:
                                return "—"
                            color = "#4caf50" if pct >= 95 else ("#ff9800" if pct >= 80 else "#f44336")
                            return (
                                f'<span style="background:{color}; color:white; padding:2px 8px; '
                                f'border-radius:10px; font-weight:600; font-size:12px;">'
                                f'{pct:.1f}%</span>'
                            )

                        html_rows = ""
                        for _, _row in orient_df.iterrows():
                            badge = _protection_badge(_row.get("Protection (%)"))
                            _dh = f"{_row['D/H (Overhang)']:.3f}" if _row["D/H (Overhang)"] is not None else "—"
                            _dw = f"{_row['D/W (Fin)']:.3f}" if _row["D/W (Fin)"] is not None else "—"
                            _min_vsa = f"{_row['Min VSA (°)']:.1f}°" if _row["Min VSA (°)"] is not None else "—"
                            _max_hsa = f"{_row['Max |HSA| (°)']:.1f}°" if _row["Max |HSA| (°)"] is not None else "—"
                            html_rows += (
                                "<tr>"
                                f"<td style='padding:6px 12px; font-weight:600;'>{_row['Orientation']}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{_row['Rays Hitting']}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{_min_vsa}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{_max_hsa}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{_dh}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{_dw}</td>"
                                f"<td style='padding:6px 12px; text-align:center;'>{badge}</td>"
                                "</tr>"
                            )

                        st.markdown(f"""
<style>
.shading-table {{ border-collapse: collapse; width: 100%; font-size: 13px; }}
.shading-table th {{
    background: #1a3a52; color: white; padding: 8px 12px;
    text-align: center; font-weight: 600; letter-spacing: 0.3px;
}}
.shading-table th:first-child {{ text-align: left; }}
.shading-table tr:nth-child(even) {{ background: #f5f7fa; }}
.shading-table tr:hover {{ background: #e8f4f8; }}
.shading-table td {{ border-bottom: 1px solid #e0e0e0; }}
</style>
<table class="shading-table">
<thead>
  <tr>
    <th>Orientation</th>
    <th>Rays Hitting</th>
    <th>Min VSA</th>
    <th>Max |HSA|</th>
    <th>D/H (Overhang)</th>
    <th>D/W (Fin)</th>
    <th>Protection %</th>
  </tr>
</thead>
<tbody>
  {html_rows}
</tbody>
</table>
<p style='font-size:11px; color:#888; margin-top:6px;'>
  <b>D/H</b> = horizontal overhang depth-to-height ratio &nbsp;|&nbsp;
  <b>D/W</b> = vertical fin depth-to-width ratio &nbsp;|&nbsp;
  Protection % = overheating rays blocked at design cutoff angle &nbsp;
  <span style='color:#4caf50; font-weight:600;'>■</span> ≥95% &nbsp;
  <span style='color:#ff9800; font-weight:600;'>■</span> 80–95% &nbsp;
  <span style='color:#f44336; font-weight:600;'>■</span> &lt;80%
</p>
""", unsafe_allow_html=True)

                        # ── 3. SHADING MASK MINI DIAGRAMS ─────────────────────
                        st.markdown(
                            "<h5 style='margin-top:22px; margin-bottom:4px;'>"
                            "🔵 Shading Mask Diagrams</h5>"
                            "<p style='font-size:12px; color:#666; margin-bottom:10px;'>"
                            "<span style='color:#E65100; font-weight:600;'>●</span> Overheating (hits façade) &nbsp;|&nbsp; "
                            "<span style='color:lightgrey; font-weight:600;'>●</span> Overheating (other side) &nbsp;|&nbsp; "
                            "<span style='color:#1565C0; font-weight:600;'>- -</span> VSA cutoff arc &nbsp;|&nbsp; "
                            "<span style='color:#388E3C; font-weight:600;'>—</span> Façade direction</p>",
                            unsafe_allow_html=True,
                        )

                        orient_pairs = list(_ORIENTATIONS.items())
                        for _row_start in range(0, 8, 4):
                            dial_cols = st.columns(4)
                            for _col_idx, (_oname, _faz) in enumerate(
                                orient_pairs[_row_start: _row_start + 4]
                            ):
                                with dial_cols[_col_idx]:
                                    st.markdown(
                                        f"<p style='font-size:11px; font-weight:600;"
                                        f" text-align:center; margin-bottom:2px;'>{_oname}</p>",
                                        unsafe_allow_html=True,
                                    )
                                    _mini_fig = make_shading_mask_chart(
                                        solar_pos, _faz, _cutoff
                                    )
                                    _safe_key = (
                                        _oname.replace(" ", "_")
                                              .replace("(", "")
                                              .replace(")", "")
                                              .replace("°", "deg")
                                    )
                                    st.plotly_chart(
                                        _mini_fig,
                                        use_container_width=True,
                                        key=f"mini_{_safe_key}",
                                    )

            except Exception as _e:
                st.error(f"Could not build Orientation Shading Analysis: {_e}")

    else:
                # Create tab buttons
        tabs_col1, tabs_col2, tabs_col3, tabs_col4, tabs_col5 = st.columns(5, gap="small")
        
        with tabs_col1:
            if st.button("Annual Trend", key="tab_annual", use_container_width=True):
                st.session_state.active_tab = "Annual Trend"
        
        with tabs_col2:
            if st.button("Monthly Trend", key="tab_monthly", use_container_width=True):
                st.session_state.active_tab = "Monthly Trend"
        
        with tabs_col3:
            if st.button("Diurnal Profile", key="tab_diurnal", use_container_width=True):
                st.session_state.active_tab = "Diurnal Profile"
        
        with tabs_col4:
            if st.button("Comfort Analysis", key="tab_comfort", use_container_width=True):
                st.session_state.active_tab = "Comfort Analysis"
        
        with tabs_col5:
            if st.button("Energy Metrics", key="tab_energy", use_container_width=True):
                st.session_state.active_tab = "Energy Metrics"
        
        # Add visual styling to show active tab
        st.markdown(f"""
            <script>
            var active_tab = '{st.session_state.active_tab}';
            </script>
        """, unsafe_allow_html=True)
        # Get the actual year from the data
        year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
        
        # Create start and end dates based on selected months
        start_month_num = st.session_state.start_month_idx + 1
        end_month_num = st.session_state.end_month_idx + 1
        
        start_date = pd.to_datetime(f"{year}-{start_month_num}-01").date()
        
        # For end_date, get the last day of the end month
        if end_month_num == 12:
            end_date = pd.to_datetime(f"{year}-12-31").date()
        else:
            end_date = (pd.to_datetime(f"{year}-{end_month_num+1}-01") - pd.Timedelta(days=1)).date()
        
        start_hour, end_hour = st.session_state.get("hour_range", (8, 18))
        
        # Compute daily min/max and average
        daily_stats = df.groupby("doy").agg({
            "dry_bulb_temperature": ["min", "max", "mean"],
            "relative_humidity": ["min", "max", "mean"],
        }).reset_index()
        
        daily_stats.columns = ["doy", "temp_min", "temp_max", "temp_avg", "rh_min", "rh_max", "rh_avg"]
        daily_stats["datetime_start"] = pd.to_datetime(daily_stats["doy"], format="%j", errors="coerce")
        
        # Map doy back to actual year-dates for chart display
        year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
        daily_stats["datetime"] = pd.to_datetime(
            daily_stats["doy"].astype(str) + f"-{year}", format="%j-%Y", errors="coerce"
        )
        
        # Add a day-month only column for display (without year)
        daily_stats["datetime_display"] = daily_stats["datetime"].dt.strftime("%b %d")
        
        # Map comfort bands to doy
        comfort_df = pd.DataFrame({
            "doy": ashrae_80_lower.index,
            "comfort_80_lower": ashrae_80_lower.values,
            "comfort_80_upper": ashrae_80_upper.values,
            "comfort_90_lower": ashrae_90_lower.values,
            "comfort_90_upper": ashrae_90_upper.values,
        })
        
        # Merge comfort data with daily stats
        daily_stats = daily_stats.merge(comfort_df, on="doy", how="left")
    
        # === TAB 1: ANNUAL TREND ===
        if st.session_state.active_tab == "Annual Trend":
            if selected_parameter == "Temperature":
                import plotly.graph_objects as go
                
                # Calculate day of year for selected month range for greying
                start_month_num = st.session_state.start_month_idx + 1
                end_month_num = st.session_state.end_month_idx + 1
                start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
                if end_month_num == 12:
                    end_doy = 366  # Use 366 to catch Dec 31
                else:
                    end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
                
                fig_yearly = go.Figure()
                
                # GREYED OUT: Before selected range
                if start_doy > 1:
                    before_data = daily_stats[daily_stats["doy"] < start_doy]
                    fig_yearly.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_yearly.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Period",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    fig_yearly.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected range
                active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
                # Add ASHRAE 80% band for ACTIVE range
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_80_upper"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_80_lower"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="ASHRAE adaptive comfort (80%)",
                    fillcolor="rgba(128, 128, 128, 0.2)",
                    hoverinfo="skip",
                ))
                
                # Add ASHRAE 90% band
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_90_upper"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_90_lower"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="ASHRAE adaptive comfort (90%)",
                    fillcolor="rgba(128, 128, 128, 0.4)",
                    hoverinfo="skip",
                ))
                
                # Add temperature range (min/max) for ACTIVE range
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    name="Dry bulb temperature Range",
                    fillcolor="rgba(255, 173, 173, 0.4)",
                    customdata=active_range["temp_max"],
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
                ))
                
                # Add average line for ACTIVE range
                fig_yearly.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_avg"],
                    mode="lines",
                    name="Average Dry bulb temperature",
                    line=dict(color="#d32f2f", width=2),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
                ))
                
                # GREYED OUT: After selected range
                if end_doy < 365:
                    after_data = daily_stats[daily_stats["doy"] > end_doy]
                    if not after_data.empty:
                        fig_yearly.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_max"],
                            fill=None,
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_yearly.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_min"],
                            fill="tonexty",
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            fillcolor="rgba(180, 180, 180, 0.15)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_yearly.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_avg"],
                            mode="lines",
                            line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                
                
                fig_yearly.update_layout(
                    title="Annual Dry Bulb Temperature Trend",
                    xaxis_title=None,
                    yaxis_title="Temperature (°C)",
                    hovermode="x unified",
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1
                    ),
                    xaxis_rangeslider_visible=False,
                    height=450,
                    template="plotly_white",
                    margin=dict(b=80),
                )
    
                st.plotly_chart(fig_yearly, use_container_width=True)
                
                # Get filtered data for period
                filtered_df = df[
                    (df["datetime"].dt.date >= start_date) &
                    (df["datetime"].dt.date <= end_date) &
                    (df["hour"].between(start_hour, end_hour))
                ]
                
                if not filtered_df.empty:
                    # Get min/max values and their corresponding rows with datetime/hour info
                    min_idx = filtered_df["dry_bulb_temperature"].idxmin()
                    max_idx = filtered_df["dry_bulb_temperature"].idxmax()
                    
                    min_row = filtered_df.loc[min_idx]
                    max_row = filtered_df.loc[max_idx]
                    
                    temp_min = min_row["dry_bulb_temperature"]
                    temp_max = max_row["dry_bulb_temperature"]
                    temp_avg = filtered_df["dry_bulb_temperature"].mean()
                    diurnal_range = temp_max - temp_min
                    
                    # Extract date and hour information for min/max
                    min_date_str = min_row["datetime"].strftime("%b %d")
                    min_hour = int(min_row["hour"])
                    
                    max_date_str = max_row["datetime"].strftime("%b %d")
                    max_hour = int(max_row["hour"])
                    
                    # Calculate additional metrics using full year data
                    # HDD18 (Heating Degree Days at 18°C base)
                    hdd18 = (18 - df["dry_bulb_temperature"]).clip(lower=0).sum()
                    
                    # CDD24 (Cooling Degree Days at 24°C base)
                    cdd24 = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                    
                    # Comfort metrics (using full year data for comfort bands)
                    def get_comfort_band_range(temps):
                        """Get comfort band as ±3.5°C from mean"""
                        mean_temp = temps.mean()
                        return mean_temp - 3.5, mean_temp + 3.5
                    
                    comfort_lower, comfort_upper = get_comfort_band_range(df["dry_bulb_temperature"])
                    comfort_hours = len(df[(df["dry_bulb_temperature"] >= comfort_lower) & (df["dry_bulb_temperature"] <= comfort_upper)])
                    comfort_80_percent = (comfort_hours / len(df)) * 100
                    
                    # 1% Cooling (99th percentile)
                    cooling_1pct = df["dry_bulb_temperature"].quantile(0.99)
                    
                    # Overheat hours (hours above 28°C threshold)
                    overheat_hrs = len(df[df["dry_bulb_temperature"] > 28])
                    
                    # Cold hours (hours below 12°C threshold)
                    cold_hrs = len(df[df["dry_bulb_temperature"] < 12])
                    
                    # Temperature metrics cards - Row 1
                    kpi_col1, kpi_col2, kpi_col3, kpi_col4, kpi_col5 = st.columns(5)
                    
                    with kpi_col1:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #f59e0b;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Min Temp</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_min:.2f} °C</div>
                                <div style="font-size: 11px; color: #718096;">{min_date_str} · {min_hour:02d}:00</div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col2:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #ef4444;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Max Temp</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_max:.2f} °C</div>
                                <div style="font-size: 11px; color: #718096;">{max_date_str} · {max_hour:02d}:00</div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col3:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #8b5cf6;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">Avg Temp</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_avg:.2f} °C</div>
                                <div style="font-size: 11px; color: #718096;">All year average</div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col4:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #3b82f6;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Diurnal Range</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{diurnal_range:.2f} °C</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col5:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #06b6d4;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">1% Cooling</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cooling_1pct:.2f} °C</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    # Row 2 - Additional metrics
                    kpi_col6, kpi_col7, kpi_col8, kpi_col9, kpi_col10 = st.columns(5)
                    
                    with kpi_col6:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #dc2626;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #dc2626; text-transform: uppercase; letter-spacing: 0.5px;">HDD18</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{hdd18:.0f}</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col7:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #0891b2;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #0891b2; text-transform: uppercase; letter-spacing: 0.5px;">CDD24</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cdd24:.0f}</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col8:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #06b6d4;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">Comfort 80%</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{comfort_80_percent:.0f} %</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col9:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #8b5cf6;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">Overheat Hrs</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{overheat_hrs}</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi_col10:
                        st.markdown(f"""
                            <div style="
                                background: white;
                                padding: 16px;
                                border-radius: 8px;
                                border-left: 4px solid #3b82f6;
                                box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                                text-align: center;
                            ">
                                <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Cold Hrs</div>
                                <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cold_hrs}</div>
                                <div style="font-size: 11px; color: #718096;"></div>
                            </div>
                        """, unsafe_allow_html=True)
                else:
                    st.info(f"No data available between {start_hour:02d}:00 and {end_hour:02d}:00 in the selected date range.")
            
            elif selected_parameter == "Humidity":
                import plotly.graph_objects as go
                
                # Define humidity comfort bands (typical comfort range 30-65%)
                humidity_comfort_lower = 30
                humidity_comfort_upper = 65
                
                fig_yearly = go.Figure()
                
                # Add humidity comfort band
                fig_yearly.add_trace(go.Scatter(
                    x=daily_stats["datetime_display"],
                    y=[humidity_comfort_upper] * len(daily_stats),
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_yearly.add_trace(go.Scatter(
                    x=daily_stats["datetime_display"],
                    y=[humidity_comfort_lower] * len(daily_stats),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="Humidity comfort band",
                    fillcolor="rgba(128, 128, 128, 0.2)",
                    hoverinfo="skip",
                ))
                
                # Add relative humidity range (min/max)
                fig_yearly.add_trace(go.Scatter(
                    x=daily_stats["datetime_display"],
                    y=daily_stats["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_yearly.add_trace(go.Scatter(
                    x=daily_stats["datetime_display"],
                    y=daily_stats["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    name="Relative humidity Range",
                    fillcolor="rgba(0, 150, 255, 0.3)",
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.1f}%<extra></extra>",
                ))
                
                # Add average line
                fig_yearly.add_trace(go.Scatter(
                    x=daily_stats["datetime_display"],
                    y=daily_stats["rh_avg"],
                    mode="lines",
                    name="Average Relative humidity",
                    line=dict(color="#00a8ff", width=2),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.1f}%<extra></extra>",
                ))
                
                fig_yearly.update_layout(
                    title="Annual Profile – Relative Humidity",
                    xaxis_title="Day",
                    yaxis_title="Relative Humidity (%)",
                    hovermode="x unified",
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1
                    ),
                    xaxis_rangeslider_visible=False,
                    height=450,
                    template="plotly_white",
                    margin=dict(b=80),
                )
                
                st.plotly_chart(fig_yearly, use_container_width=True)
                
                # Calculate humidity metrics
                comfort_rh_lower = 40  # Comfort zone lower bound
                comfort_rh_upper = 60  # Comfort zone upper bound
                
                # Get min/max RH values
                rh_min = df["relative_humidity"].min()
                rh_max = df["relative_humidity"].max()
                rh_avg = df["relative_humidity"].mean()
                
                # Find date/time of peak RH
                max_rh_idx = df["relative_humidity"].idxmax()
                max_rh_row = df.loc[max_rh_idx]
                max_rh_date_str = max_rh_row["datetime"].strftime("%b %d")
                max_rh_hour = int(max_rh_row["hour"])
                
                # Comfort hours (40-60% RH)
                comfort_rh_hours = len(df[(df["relative_humidity"] >= comfort_rh_lower) & (df["relative_humidity"] <= comfort_rh_upper)])
                comfort_rh_percent = (comfort_rh_hours / len(df)) * 100
                
                # High humidity hours (> 60% RH)
                high_humidity_hrs = len(df[df["relative_humidity"] > 60])
                
                # Condensation risk hours (> 75% RH sustained)
                condensation_risk_hrs = len(df[df["relative_humidity"] > 75])
                
                # Low humidity hours (< 30% RH)
                low_humidity_hrs = len(df[df["relative_humidity"] < 30])
                
                # Mold risk hours (> 60% RH sustained for extended periods)
                mold_risk_hrs = len(df[df["relative_humidity"] > 60])
                
                # HVAC RH control (percentage within comfort band)
                hvac_rh_control = comfort_rh_percent
                
                # Overhumidification hours (> 70% RH)
                overhumidification_hrs = len(df[df["relative_humidity"] > 70])
                
                # Humidity metrics cards - Row 1
                rh_col1, rh_col2, rh_col3, rh_col4, rh_col5 = st.columns(5)
                
                with rh_col1:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #f59e0b;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Comfort 40-60%</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{comfort_rh_percent:.0f} %</div>
                            <div style="font-size: 11px; color: #718096;">Occupied RH Hrs</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col2:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #ef4444;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Peak RH (Occupied)</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_max:.1f} %</div>
                            <div style="font-size: 11px; color: #718096;">All year</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col3:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #8b5cf6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">High Humidity Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{high_humidity_hrs}</div>
                            <div style="font-size: 11px; color: #718096;">&gt; 60% RH</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col4:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #06b6d4;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">Condensation Risk Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{condensation_risk_hrs}</div>
                            <div style="font-size: 11px; color: #718096;">Surface Temp &lt; Dew Point</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col5:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #3b82f6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Avg RH</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_avg:.1f} %</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                # Row 2 - Additional metrics
                rh_col6, rh_col7, rh_col8, rh_col9, rh_col10 = st.columns(5)
                
                with rh_col6:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #f59e0b;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Low Humidity Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{low_humidity_hrs}</div>
                            <div style="font-size: 11px; color: #718096;">&lt; 30% RH</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col7:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #ef4444;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Mold Risk Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{mold_risk_hrs}</div>
                            <div style="font-size: 11px; color: #718096;">&gt; 60% RH Sustained</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col8:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #06b6d4;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">HVAC RH Control</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{hvac_rh_control:.0f} %</div>
                            <div style="font-size: 11px; color: #718096;">Outside RH vs Inside RH</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col9:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #3b82f6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Overhumidification Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{overhumidification_hrs}</div>
                            <div style="font-size: 11px; color: #718096;">System Failure Indicator</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with rh_col10:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #0891b2;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #0891b2; text-transform: uppercase; letter-spacing: 0.5px;">Min RH</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_min:.1f} %</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
            
            else:
                st.info("Sun Path analysis is not yet implemented.")
    
        # === TAB 2: MONTHLY TREND ===
        elif st.session_state.active_tab == "Monthly Trend":
            if selected_parameter == "Temperature":
                import plotly.graph_objects as go
                
                # Calculate monthly statistics
                monthly_stats = df.groupby("month").agg({
                    "dry_bulb_temperature": ["min", "max", "mean"],
                    "relative_humidity": ["min", "max", "mean"],
                }).reset_index()
                
                monthly_stats.columns = ["month", "temp_min", "temp_max", "temp_avg", "rh_min", "rh_max", "rh_avg"]
                month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                monthly_stats["month_name"] = monthly_stats["month"].apply(lambda x: month_names[x-1])
                
                fig_monthly = go.Figure()
                
                # Get selected month range
                start_month = st.session_state.start_month_idx + 1
                end_month = st.session_state.end_month_idx + 1
                
                # GREYED OUT: Before selected range
                if start_month > 1:
                    before_months = monthly_stats[monthly_stats["month"] < start_month]
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Period",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["temp_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected range
                active_months = monthly_stats[(monthly_stats["month"] >= start_month) & (monthly_stats["month"] <= end_month)]
                
                # Add temperature range (min/max) for active range
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    name="Monthly Temperature Range",
                    fillcolor="rgba(255, 173, 173, 0.4)",
                    customdata=active_months["temp_max"],
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
                ))
                
                # Add average line for active range
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["temp_avg"],
                    mode="lines+markers",
                    name="Monthly Average Temperature",
                    line=dict(color="#d32f2f", width=2),
                    marker=dict(size=8),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
                ))
                
                # GREYED OUT: After selected range
                if end_month < 12:
                    after_months = monthly_stats[monthly_stats["month"] > end_month]
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["temp_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                
                fig_monthly.update_layout(
                    title="Monthly Temperature Trend",
                    xaxis_title="Month",
                    yaxis_title="Temperature (°C)",
                    hovermode="x unified",
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1
                    ),
                    height=450,
                    template="plotly_white",
                    margin=dict(b=80),
                )
                
                st.plotly_chart(fig_monthly, use_container_width=True)
                
                # Display monthly KPI metrics
                st.markdown("#### Monthly Temperature Summary")
                
                # Create a dataframe for monthly metrics
                kpi_data = monthly_stats[["month_name", "temp_min", "temp_max", "temp_avg"]].copy()
                kpi_data.columns = ["Month", "Min (°C)", "Max (°C)", "Avg (°C)"]
                
                # Display as a nice table
                st.dataframe(
                    kpi_data,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Min (°C)": st.column_config.NumberColumn(format="%.2f"),
                        "Max (°C)": st.column_config.NumberColumn(format="%.2f"),
                        "Avg (°C)": st.column_config.NumberColumn(format="%.2f"),
                    }
                )
                
            elif selected_parameter == "Humidity":
                import plotly.graph_objects as go
                
                # Get selected month range for greying
                start_month = st.session_state.start_month_idx + 1
                end_month = st.session_state.end_month_idx + 1
                
                # Calculate monthly humidity statistics
                monthly_stats = df.groupby("month").agg({
                    "relative_humidity": ["min", "max", "mean"],
                }).reset_index()
                
                monthly_stats.columns = ["month", "rh_min", "rh_max", "rh_avg"]
                month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                monthly_stats["month_name"] = monthly_stats["month"].apply(lambda x: month_names[x-1])
                
                # Split data into before/active/after ranges
                before_months = monthly_stats[monthly_stats["month"] < start_month]
                active_months = monthly_stats[(monthly_stats["month"] >= start_month) & (monthly_stats["month"] <= end_month)]
                after_months = monthly_stats[monthly_stats["month"] > end_month]
                
                fig_monthly = go.Figure()
                
                # GREYED OUT: Before selected range
                if not before_months.empty:
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=[65] * len(before_months),
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=[30] * len(before_months),
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Period",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=before_months["month_name"],
                        y=before_months["rh_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected month range
                # Add humidity comfort band for active range
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=[65] * len(active_months),
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=[30] * len(active_months),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="Humidity comfort band (30-65%)",
                    fillcolor="rgba(128, 128, 128, 0.2)",
                    hoverinfo="skip",
                ))
                
                # Add humidity range (min/max) for active range
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    name="Monthly Humidity Range",
                    fillcolor="rgba(0, 150, 255, 0.3)",
                    customdata=active_months["rh_max"],
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.1f}%<br>Max: %{customdata:.1f}%<extra></extra>",
                ))
                
                # Add average line for active range
                fig_monthly.add_trace(go.Scatter(
                    x=active_months["month_name"],
                    y=active_months["rh_avg"],
                    mode="lines+markers",
                    name="Monthly Average Humidity",
                    line=dict(color="#00a8ff", width=2),
                    marker=dict(size=8),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.1f}%<extra></extra>",
                ))
                
                # GREYED OUT: After selected range
                if not after_months.empty:
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=[65] * len(after_months),
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=[30] * len(after_months),
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_monthly.add_trace(go.Scatter(
                        x=after_months["month_name"],
                        y=after_months["rh_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                fig_monthly.update_layout(
                    title="Monthly Relative Humidity Trend",
                    xaxis_title="Month",
                    yaxis_title="Relative Humidity (%)",
                    hovermode="x unified",
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1
                    ),
                    height=450,
                    template="plotly_white",
                    margin=dict(b=80),
                )
                
                st.plotly_chart(fig_monthly, use_container_width=True)
                
                # Display monthly KPI metrics
                st.markdown("#### Monthly Humidity Summary")
                
                # Create a dataframe for monthly metrics
                kpi_data = monthly_stats[["month_name", "rh_min", "rh_max", "rh_avg"]].copy()
                kpi_data.columns = ["Month", "Min (%)", "Max (%)", "Avg (%)"]
                
                # Display as a nice table
                st.dataframe(
                    kpi_data,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Min (%)": st.column_config.NumberColumn(format="%.1f"),
                        "Max (%)": st.column_config.NumberColumn(format="%.1f"),
                        "Avg (%)": st.column_config.NumberColumn(format="%.1f"),
                    }
                )
            
            else:
                st.info("Monthly trend analysis is not yet implemented for Sun Path.")
    
        # === TAB 3: DIURNAL PROFILE ===
        elif st.session_state.active_tab == "Diurnal Profile":
            import plotly.graph_objects as go
            
            if selected_parameter == "Temperature":
                # Create hourly averages for each month
                hourly_stats = df.groupby(["month", "hour"]).agg({
                    "dry_bulb_temperature": ["min", "max", "mean"],
                }).reset_index()
                
                hourly_stats.columns = ["month", "hour", "temp_min", "temp_max", "temp_avg"]
                hourly_stats["month_name"] = hourly_stats["month"].apply(lambda x: ["Jan", "Feb", "Mar", "Apr", "May", "Jun", 
                                                                                      "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"][x-1])
                
                fig_diurnal = go.Figure()
                
                # Get hourly averages across ALL months for diurnal profile
                # Diurnal Profile uses TIME filter only, not month filter
                avg_hourly = hourly_stats.groupby("hour").agg({
                    "temp_min": "min",
                    "temp_max": "max",
                    "temp_avg": "mean"
                }).reset_index()
                
                # Get selected hour range
                start_hour_sel, end_hour_sel = st.session_state.get("hour_range", (8, 18))
                
                # GREYED OUT: Before selected hours
                if start_hour_sel > 0:
                    before_hours = avg_hourly[avg_hourly["hour"] < start_hour_sel]
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours["hour"],
                        y=before_hours["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours["hour"],
                        y=before_hours["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Hours",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours["hour"],
                        y=before_hours["temp_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected hours
                active_hours = avg_hourly[(avg_hourly["hour"] >= start_hour_sel) & (avg_hourly["hour"] <= end_hour_sel)]
                
                # Add temperature range for active hours
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours["hour"],
                    y=active_hours["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours["hour"],
                    y=active_hours["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    name="Daily Range",
                    fillcolor="rgba(255, 173, 173, 0.3)",
                    customdata=active_hours["temp_max"],
                    hovertemplate="<b>Hour %{x}:00</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
                ))
                
                # Add average for active hours
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours["hour"],
                    y=active_hours["temp_avg"],
                    mode="lines+markers",
                    name="Average Temperature",
                    line=dict(color="#d32f2f", width=2),
                    marker=dict(size=6),
                    hovertemplate="<b>Hour %{x}:00</b><br>Avg: %{y:.2f}°C<extra></extra>",
                ))
                
                # GREYED OUT: After selected hours
                if end_hour_sel < 23:
                    after_hours = avg_hourly[avg_hourly["hour"] > end_hour_sel]
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours["hour"],
                        y=after_hours["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours["hour"],
                        y=after_hours["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours["hour"],
                        y=after_hours["temp_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                
                fig_diurnal.update_layout(
                    title="Diurnal Temperature Profile",
                    xaxis_title="Hour of Day",
                    yaxis_title="Temperature (°C)",
                    hovermode="x unified",
                    showlegend=True,
                    template="plotly_white",
                    height=450,
                )
                
                st.plotly_chart(fig_diurnal, use_container_width=True)
            
            elif selected_parameter == "Humidity":
                # Create hourly humidity averages for each month
                hourly_humidity = df.groupby(["month", "hour"]).agg({
                    "relative_humidity": ["min", "max", "mean"],
                }).reset_index()
                
                hourly_humidity.columns = ["month", "hour", "rh_min", "rh_max", "rh_avg"]
                
                fig_diurnal = go.Figure()
                
                # Get hourly averages across ALL months for diurnal profile
                avg_hourly_rh = hourly_humidity.groupby("hour").agg({
                    "rh_min": "min",
                    "rh_max": "max",
                    "rh_avg": "mean"
                }).reset_index()
                
                # Get selected hour range
                start_hour_sel, end_hour_sel = st.session_state.get("hour_range", (8, 18))
                
                # GREYED OUT: Before selected hours
                if start_hour_sel > 0:
                    before_hours_rh = avg_hourly_rh[avg_hourly_rh["hour"] < start_hour_sel]
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours_rh["hour"],
                        y=[65] * len(before_hours_rh),
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours_rh["hour"],
                        y=[30] * len(before_hours_rh),
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Hours",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours_rh["hour"],
                        y=before_hours_rh["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours_rh["hour"],
                        y=before_hours_rh["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=before_hours_rh["hour"],
                        y=before_hours_rh["rh_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected hours
                active_hours_rh = avg_hourly_rh[(avg_hourly_rh["hour"] >= start_hour_sel) & (avg_hourly_rh["hour"] <= end_hour_sel)]
                
                # Add humidity comfort band for active hours
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours_rh["hour"],
                    y=[65] * len(active_hours_rh),
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours_rh["hour"],
                    y=[30] * len(active_hours_rh),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="Comfort band (30-65%)",
                    fillcolor="rgba(128, 128, 128, 0.2)",
                    hoverinfo="skip",
                ))
                
                # Add humidity range for active hours
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours_rh["hour"],
                    y=active_hours_rh["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours_rh["hour"],
                    y=active_hours_rh["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(0, 0, 255, 0)",
                    name="Humidity Range",
                    fillcolor="rgba(0, 150, 255, 0.3)",
                    customdata=active_hours_rh["rh_max"],
                    hovertemplate="<b>Hour %{x}:00</b><br>Min: %{y:.1f}%<br>Max: %{customdata:.1f}%<extra></extra>",
                ))
                
                # Add average for active hours
                fig_diurnal.add_trace(go.Scatter(
                    x=active_hours_rh["hour"],
                    y=active_hours_rh["rh_avg"],
                    mode="lines+markers",
                    name="Average Humidity",
                    line=dict(color="#00a8ff", width=2),
                    marker=dict(size=6),
                    hovertemplate="<b>Hour %{x}:00</b><br>Avg: %{y:.1f}%<extra></extra>",
                ))
                
                # GREYED OUT: After selected hours
                if end_hour_sel < 23:
                    after_hours_rh = avg_hourly_rh[avg_hourly_rh["hour"] > end_hour_sel]
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours_rh["hour"],
                        y=[65] * len(after_hours_rh),
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours_rh["hour"],
                        y=[30] * len(after_hours_rh),
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours_rh["hour"],
                        y=after_hours_rh["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours_rh["hour"],
                        y=after_hours_rh["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    
                    fig_diurnal.add_trace(go.Scatter(
                        x=after_hours_rh["hour"],
                        y=after_hours_rh["rh_avg"],
                        mode="lines+markers",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        marker=dict(size=4),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                fig_diurnal.update_layout(
                    title="Diurnal Humidity Profile",
                    xaxis_title="Hour of Day",
                    yaxis_title="Relative Humidity (%)",
                    hovermode="x unified",
                    showlegend=True,
                    template="plotly_white",
                    height=450,
                )
                
                st.plotly_chart(fig_diurnal, use_container_width=True)
            
            else:
                st.info("Sun Path analysis is not yet implemented.")
        
        # === TAB 4: COMFORT ANALYSIS ===
        elif st.session_state.active_tab == "Comfort Analysis":
            if selected_parameter == "Temperature":
                import plotly.graph_objects as go
                
                # Calculate day of year for selected month range for greying
                start_month_num = st.session_state.start_month_idx + 1
                end_month_num = st.session_state.end_month_idx + 1
                start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
                if end_month_num == 12:
                    end_doy = 366
                else:
                    end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
                
                # Create comfort analysis visualization
                fig_comfort = go.Figure()
                
                # GREYED OUT: Before selected range
                if start_doy > 1:
                    before_data = daily_stats[daily_stats["doy"] < start_doy]
                    fig_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Period",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    fig_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["temp_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected range
                active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
                
                # Add comfort bands for active range
                fig_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_90_upper"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["comfort_90_lower"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="ASHRAE 90% acceptability",
                    fillcolor="rgba(76, 175, 80, 0.4)",
                    hoverinfo="skip",
                ))
                
                # Add temperature data for active range
                fig_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(255, 0, 0, 0)",
                    name="Daily Temperature Range",
                    fillcolor="rgba(255, 173, 173, 0.3)",
                    customdata=active_range["temp_max"],
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
                ))
                
                # Add average for active range
                fig_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["temp_avg"],
                    mode="lines",
                    name="Average Temperature",
                    line=dict(color="#d32f2f", width=2),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
                ))
                
                # GREYED OUT: After selected range
                if end_doy < 365:
                    after_data = daily_stats[daily_stats["doy"] > end_doy]
                    if not after_data.empty:
                        fig_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_max"],
                            fill=None,
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_min"],
                            fill="tonexty",
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            fillcolor="rgba(180, 180, 180, 0.15)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["temp_avg"],
                            mode="lines",
                            line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                
                
                fig_comfort.update_layout(
                    title="Comfort Analysis – ASHRAE Adaptive Comfort",
                    xaxis_title="Day",
                    yaxis_title="Temperature (°C)",
                    hovermode="x unified",
                    showlegend=True,
                    template="plotly_white",
                    height=450,
                )
                
                st.plotly_chart(fig_comfort, use_container_width=True)
            
            elif selected_parameter == "Humidity":
                import plotly.graph_objects as go
                
                # Calculate day of year for selected month range for greying
                start_month_num = st.session_state.start_month_idx + 1
                end_month_num = st.session_state.end_month_idx + 1
                start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
                if end_month_num == 12:
                    end_doy = 366
                else:
                    end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
                
                # Create humidity comfort analysis visualization
                fig_humidity_comfort = go.Figure()
                
                # Define comfort zones for humidity
                comfort_upper = 60  # Upper comfort limit
                comfort_lower = 40  # Lower comfort limit
                
                # GREYED OUT: Before selected range
                if start_doy > 1:
                    before_data = daily_stats[daily_stats["doy"] < start_doy]
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        name="Unselected Period",
                        showlegend=True,
                        hoverinfo="skip",
                    ))
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=before_data["datetime_display"],
                        y=before_data["rh_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                
                # ACTIVE: Selected range
                active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
                
                # Add comfort band (40-60% RH) for active range
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=[comfort_upper] * len(active_range),
                    fill=None,
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=[comfort_lower] * len(active_range),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(128, 128, 128, 0)",
                    name="Comfort Band (40-60%)",
                    fillcolor="rgba(76, 175, 80, 0.4)",
                    hoverinfo="skip",
                ))
                
                # Add humidity data for active range
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(0, 150, 255, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(0, 150, 255, 0)",
                    name="Daily RH Range",
                    fillcolor="rgba(0, 150, 255, 0.3)",
                    customdata=active_range["rh_max"],
                    hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}%<br>Max: %{customdata:.2f}%<extra></extra>",
                ))
                
                # Add average for active range
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=active_range["datetime_display"],
                    y=active_range["rh_avg"],
                    mode="lines",
                    name="Average RH",
                    line=dict(color="#00a8ff", width=2),
                    hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}%<extra></extra>",
                ))
                
                # GREYED OUT: After selected range
                if end_doy < 365:
                    after_data = daily_stats[daily_stats["doy"] > end_doy]
                    if not after_data.empty:
                        fig_humidity_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["rh_max"],
                            fill=None,
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_humidity_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["rh_min"],
                            fill="tonexty",
                            mode="lines",
                            line_color="rgba(100, 100, 100, 0)",
                            fillcolor="rgba(180, 180, 180, 0.15)",
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                        fig_humidity_comfort.add_trace(go.Scatter(
                            x=after_data["datetime_display"],
                            y=after_data["rh_avg"],
                            mode="lines",
                            line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                            showlegend=False,
                            hoverinfo="skip",
                        ))
                
                fig_humidity_comfort.update_layout(
                    title="Humidity Comfort Analysis – Optimal Range (40-60%)",
                    xaxis_title="Day",
                    yaxis_title="Relative Humidity (%)",
                    hovermode="x unified",
                    showlegend=True,
                    template="plotly_white",
                    height=450,
                )
                
                st.plotly_chart(fig_humidity_comfort, use_container_width=True)
            else:
                st.info("Comfort Analysis is available for Temperature and Humidity parameters.")
        
        # === TAB 5: ENERGY METRICS ===
        elif st.session_state.active_tab == "Energy Metrics":
            if selected_parameter == "Temperature":
                # Calculate energy metrics
                filtered_df = df[
                    (df["datetime"].dt.date >= start_date) &
                    (df["datetime"].dt.date <= end_date) &
                    (df["hour"].between(start_hour, end_hour))
                ]
                
                if not filtered_df.empty:
                    # HDD18 (Heating Degree Days at 18°C base)
                    hdd18 = (18 - df["dry_bulb_temperature"]).clip(lower=0).sum()
                    
                    # CDD24 (Cooling Degree Days at 24°C base)
                    cdd24 = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                    
                    # Additional energy metrics
                    hdd18_filtered = (18 - filtered_df["dry_bulb_temperature"]).clip(lower=0).sum()
                    cdd24_filtered = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                    
                    # Degree days by month
                    monthly_hdd = df.groupby("month").apply(lambda x: (18 - x["dry_bulb_temperature"]).clip(lower=0).sum())
                    monthly_cdd = df.groupby("month").apply(lambda x: (x["dry_bulb_temperature"] - 24).clip(lower=0).sum())
                    
                    month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                    
                    # Display metrics in cards
                    st.markdown("#### Energy Performance Indicators")
                    
                    energy_col1, energy_col2, energy_col3, energy_col4 = st.columns(4)
                    
                    with energy_col1:
                        st.metric("HDD18 (Annual)", f"{hdd18:.0f}", "Heating Degree-Days")
                    
                    with energy_col2:
                        st.metric("CDD24 (Annual)", f"{cdd24:.0f}", "Cooling Degree-Days")
                    
                    with energy_col3:
                        st.metric("HDD18 (Period)", f"{hdd18_filtered:.0f}", "Heating Degree-Days")
                    
                    with energy_col4:
                        st.metric("CDD24 (Period)", f"{cdd24_filtered:.0f}", "Cooling Degree-Days")
                    
                    # Monthly breakdown chart
                    import plotly.graph_objects as go
                    from plotly.subplots import make_subplots
                    
                    fig_energy = make_subplots(specs=[[{"secondary_y": True}]])
                    
                    fig_energy.add_trace(
                        go.Bar(x=month_names, y=monthly_hdd.values, name="HDD18", marker_color="#2196F3"),
                        secondary_y=False,
                    )
                    
                    fig_energy.add_trace(
                        go.Bar(x=month_names, y=monthly_cdd.values, name="CDD24", marker_color="#FF9800"),
                        secondary_y=False,
                    )
                    
                    
                    fig_energy.update_layout(
                        title="Monthly Degree-Days Distribution",
                        xaxis_title="Month",
                        yaxis_title="Degree-Days",
                        hovermode="x unified",
                        height=400,
                        barmode="stack",
                    )
                    
                    st.plotly_chart(fig_energy, use_container_width=True)
            
            elif selected_parameter == "Humidity":
                # Calculate humidity-related metrics
                filtered_df = df[
                    (df["datetime"].dt.date >= start_date) &
                    (df["datetime"].dt.date <= end_date) &
                    (df["hour"].between(start_hour, end_hour))
                ]
                
                if not filtered_df.empty:
                    # Calculate humidity metrics for full year
                    high_humidity_annual = len(df[df["relative_humidity"] > 60])
                    condensation_risk_annual = len(df[df["relative_humidity"] > 75])
                    mold_risk_annual = len(df[df["relative_humidity"] > 60])
                    low_humidity_annual = len(df[df["relative_humidity"] < 30])
                    
                    # Filtered period metrics
                    high_humidity_filtered = len(filtered_df[filtered_df["relative_humidity"] > 60])
                    condensation_risk_filtered = len(filtered_df[filtered_df["relative_humidity"] > 75])
                    
                    # Monthly breakdown
                    monthly_high_rh = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 60]))
                    monthly_condensation = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 75]))
                    monthly_mold_risk = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 60]))
                    monthly_low_rh = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] < 30]))
                    
                    month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                    
                    # Display metrics in cards
                    st.markdown("#### Humidity Performance Indicators")
                    
                    humidity_col1, humidity_col2, humidity_col3, humidity_col4 = st.columns(4)
                    
                    with humidity_col1:
                        st.metric("High RH Hrs (Annual)", f"{high_humidity_annual:.0f}", ">60% RH")
                    
                    with humidity_col2:
                        st.metric("Condensation Risk (Annual)", f"{condensation_risk_annual:.0f}", ">75% RH")
                    
                    with humidity_col3:
                        st.metric("High RH Hrs (Period)", f"{high_humidity_filtered:.0f}", ">60% RH")
                    
                    with humidity_col4:
                        st.metric("Condensation (Period)", f"{condensation_risk_filtered:.0f}", ">75% RH")
                    
                    # Monthly breakdown chart
                    import plotly.graph_objects as go
                    from plotly.subplots import make_subplots
                    
                    fig_humidity_energy = make_subplots(specs=[[{"secondary_y": True}]])
                    
                    fig_humidity_energy.add_trace(
                        go.Bar(x=month_names, y=monthly_high_rh.values, name="High RH (>60%)", marker_color="#0099ff"),
                        secondary_y=False,
                    )
                    
                    fig_humidity_energy.add_trace(
                        go.Bar(x=month_names, y=monthly_condensation.values, name="Condensation Risk (>75%)", marker_color="#FF6B6B"),
                        secondary_y=False,
                    )
                    
                    fig_humidity_energy.add_trace(
                        go.Scatter(x=month_names, y=monthly_low_rh.values, name="Low RH (<30%)", 
                                  line=dict(color="#FFA500", width=2), mode="lines+markers"),
                        secondary_y=False,
                    )
                    
                    fig_humidity_energy.update_layout(
                        title="Monthly Humidity Risk Distribution",
                        xaxis_title="Month",
                        yaxis_title="Hours",
                        hovermode="x unified",
                        height=400,
                        barmode="group",
                    )
                    
                    st.plotly_chart(fig_humidity_energy, use_container_width=True)
            else:
                st.info("Energy Metrics is available for Temperature and Humidity parameters.")
    
    # === GENERATE REPORT LOGIC ===
    # if st.session_state.get("generate_report", False):
        # with st.spinner("Generating PowerPoint report..."):
    # try:
    #     # Get filter parameters from left column
    #     start_month_num = st.session_state.start_month_idx + 1
    #     end_month_num = st.session_state.end_month_idx + 1
    #     year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
        
    #     start_date = pd.to_datetime(f"{year}-{start_month_num}-01").date()
    #     if end_month_num == 12:
    #         end_date = pd.to_datetime(f"{year}-12-31").date()
    #     else:
    #         end_date = (pd.to_datetime(f"{year}-{end_month_num+1}-01") - pd.Timedelta(days=1)).date()
        
    #     start_hour, end_hour = st.session_state.get("hour_range", (8, 18))
        
    #     # Generate the report
    #     report_bytes = generate_pptx_report(
    #         df, 
    #         start_date, 
    #         end_date, 
    #         start_hour, 
    #         end_hour, 
    #         selected_parameter
    #     )
        
    #     # Offer download
    #     st.download_button(
    #         label="📥 Download Report (PowerPoint)",
    #         data=report_bytes,
    #         file_name=f"Climate_Analysis_Report_{start_date.strftime('%Y%m%d')}_to_{end_date.strftime('%Y%m%d')}.pptx",
    #         mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    #         key="download_report"
    #     )
        
        # Reset the flag
        # st.session_state.generate_report = False
        
#     st.error(f"❌ Failed to generate report: {str(e)}")
#     st.session_state.generate_report = False

# Adding extra space at the bottom
st.markdown("<br><br>", unsafe_allow_html=True)

# Footer Section
st.image("images/EDS-footer.png", width=2000)
st.markdown(
    """
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.4/css/all.min.css">
    <style>
        .footer {
            background-color: #f8f9fa;
            padding: 20px 0;
            color: #333;
            display: flex;
            justify-content: space-between;
            align-items: center;
            text-align: center;
        }
        .footer .logo {
            flex: 1;
        }
        .footer .logo img {
            max-width: 150px;
            height: auto;
        }
        .footer .social-media {
            flex: 2;
        }
        .footer .social-media p {
            margin: 0;
            font-size: 16px;
        }
        .footer .icons {
            margin-top: 10px;
        }
        .footer .icons a {
            margin: 0 10px;
            color: #666;
            text-decoration: none;
            transition: color 0.3s ease;
        }
        .footer .icons a:hover {
            color: #0077b5; /* LinkedIn color as default */
        }
        .footer .icons a .fab {
            font-size: 28px;
        }
        .footer .additional-content {
            margin-top: 10px;
        }
        .footer .additional-content h4 {
            margin: 0;
            font-size: 18px;
            color: #007bff;
        }
        .footer .additional-content p {
            margin: 5px 0;
            font-size: 16px;
        }
    </style>
    <div style="text-align:center; font-size:14px;">
        Email: <a href="mailto:info@edsglobal.com">info@edsglobal.com</a>   |   
        Phone: +91 . 11 . 4056 8633   |   
        <a href="https://twitter.com/edsglobal?lang=en" target="_blank"><i class="fab fa-twitter" style="color:#1DA1F2; margin:0 6px;"></i></a>
        <a href="https://www.facebook.com/Environmental.Design.Solutions/" target="_blank"><i class="fab fa-facebook" style="color:#4267B2; margin:0 6px;"></i></a>
        <a href="https://www.instagram.com/eds_global/?hl=en" target="_blank"><i class="fab fa-instagram" style="color:#E1306C; margin:0 6px;"></i></a>
        <a href="https://www.linkedin.com/company/environmental-design-solutions/" target="_blank"><i class="fab fa-linkedin" style="color:#0077b5; margin:0 6px;"></i></a>
    </div>
    """,
    unsafe_allow_html=True
)