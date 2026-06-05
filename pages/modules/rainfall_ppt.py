"""Rainfall Analysis PowerPoint report generation."""

import io
import os
import tempfile

import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modules.rainfall_module import (
    _RUNOFF_SURFACES,
    _fetch_noaa,
    _fetch_percentile_depth,
)

_MONTH_LABELS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                 "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


def generate_rainfall_pptx_report(
    station_name: str,
    station_id: str,
    year: int,
    start_month: int = 1,
    end_month: int = 12,
    heavy_rain_threshold: float = 50.0,
    surface_areas: dict = None,
    gi_percentile: int = 95,
    gi_start_year: int = 1990,
) -> io.BytesIO:
    """Generate a standalone Rainfall Analysis PPTX report."""

    # ── Data ─────────────────────────────────────────────────────────────────
    df = _fetch_noaa(station_id, year)
    if df.empty:
        raise ValueError(f"No NOAA data for station {station_id}, year {year}.")

    df_f = df[(df["month"] >= start_month) & (df["month"] <= end_month)].copy()
    if surface_areas is None:
        surface_areas = {"roof": 0.0, "paved": 0.0, "green": 0.0, "water": 0.0}

    # ── Presentation setup ────────────────────────────────────────────────────
    try:
        _here = os.path.abspath(__file__)
    except NameError:
        _here = os.getcwd()
    base_dir      = os.path.dirname(os.path.dirname(os.path.dirname(_here)))
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path     = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        from pptx.oxml.ns import qn
        prs = Presentation(template_path)
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)
    else:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK = prs.slide_layouts[6]

    TITLE_RED   = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER_CLR = RGBColor(0xC0, 0x00, 0x00)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches

    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    # ── Shared helpers ────────────────────────────────────────────────────────

    def _logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path,
                Inches(LOGO_L), Inches(LOGO_T),
                width=Inches(LOGO_W), height=Inches(LOGO_H))

    def _title(slide, text):
        tb = slide.shapes.add_textbox(
            Inches(0.27), Inches(0.13), Inches(SW - 0.57), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _divider(slide):
        line = slide.shapes.add_shape(
            1, Inches(0.27), Inches(0.62), Inches(SW - 0.54), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = DIVIDER_CLR
        line.line.fill.background()

    def _save_fig(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches="tight", facecolor="white")
            return tmp.name

    def _err(slide, msg):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = f"Visualization error: {msg}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    def _hex_rgb(h: str) -> RGBColor:
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

    def _kpi_card(slide, left, top, width, height,
                  label: str, value: str, sub: str, hex_color: str):
        """Colored KPI card with label / large value / sub-label."""
        # Background
        rect = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_rgb(hex_color)
        rect.line.fill.background()

        label_color = RGBColor(0xCC, 0xDD, 0xFF)   # soft blue-white
        sub_color   = RGBColor(0xBB, 0xCC, 0xEE)

        # Label  — small caps, top of card
        tb_l = slide.shapes.add_textbox(
            Inches(left + 0.07), Inches(top + 0.07),
            Inches(width - 0.14), Inches(0.22))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.alignment = PP_ALIGN.CENTER
        r_l = p_l.add_run()
        r_l.text = label.upper()
        r_l.font.size = Pt(8)
        r_l.font.bold = True
        r_l.font.color.rgb = label_color

        # Value — large, bold, centred vertically
        val_font = Pt(22) if height >= 1.1 else Pt(17)
        val_top  = top + (height - 0.25) / 2 - 0.08
        tb_v = slide.shapes.add_textbox(
            Inches(left + 0.07), Inches(val_top),
            Inches(width - 0.14), Inches(0.55))
        tf_v = tb_v.text_frame
        tf_v.word_wrap = True
        p_v = tf_v.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        r_v = p_v.add_run()
        r_v.text = value
        r_v.font.size = val_font
        r_v.font.bold = True
        r_v.font.color.rgb = WHITE

        # Sub — small, bottom of card
        tb_s = slide.shapes.add_textbox(
            Inches(left + 0.07), Inches(top + height - 0.23),
            Inches(width - 0.14), Inches(0.20))
        tf_s = tb_s.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.alignment = PP_ALIGN.CENTER
        r_s = p_s.add_run()
        r_s.text = sub
        r_s.font.size = Pt(7.5)
        r_s.font.color.rgb = sub_color

    def _kpi_row(slide, cards, top: float, card_h: float = 1.15):
        """Evenly distribute N KPI cards across the full slide width.
           cards = [(label, value, sub, hex_color), ...]
        """
        n       = len(cards)
        gap     = 0.12
        total_w = SW - 0.54          # 12.79" usable
        card_w  = (total_w - gap * (n - 1)) / n
        left    = 0.27
        for label, value, sub, color in cards:
            _kpi_card(slide, left, top, card_w, card_h, label, value, sub, color)
            left += card_w + gap

    def _intensity_color(mm):
        if mm < 50:    return "#93c5fd"
        elif mm < 150: return "#3b82f6"
        elif mm < 300: return "#1d4ed8"
        else:          return "#1e3a5f"

    x  = np.arange(12)
    ml = _MONTH_LABELS

    # ── COVER ─────────────────────────────────────────────────────────────────
    def _cover():
        slide = prs.slides.add_slide(BLANK)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.65), Inches(SW - 1.2), Inches(1.1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Rainfall Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.6))
        run2 = tb2.text_frame.paragraphs[0].add_run()
        run2.text = f"{station_name}  |  {year}"
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        run3 = tb3.text_frame.paragraphs[0].add_run()
        run3.text = (
            "Rainfall Patterns  ·  Intensity Distribution  "
            "·  Surface Runoff  ·  GI Recharge Balance"
        )
        run3.font.size = Pt(11)
        run3.font.color.rgb = DARK_GREY
        _logo(slide)

    _cover()

    # Chart area constants shared by most slides
    CHART_TOP  = 0.72   # inches below divider
    CHART_H    = 4.48   # inches  →  chart bottom = 5.20"
    KPI_TOP    = 5.33   # KPI row top
    KPI_H      = 1.15   # KPI card height

    # ── SLIDE 1: MONTHLY RAINFALL ─────────────────────────────────────────────
    def _monthly_rainfall():
        slide = prs.slides.add_slide(BLANK)
        _title(slide, f"Monthly Rainfall Totals  —  {station_name}, {year}")
        _divider(slide)

        try:
            monthly      = df_f.groupby("month")["prcp_mm"].sum().reindex(range(1, 13), fill_value=0)
            annual_total = float(monthly.sum())
            annual_mean  = annual_total / 12
            wettest_idx  = int(monthly.idxmax())

            fig, ax = plt.subplots(figsize=(13, 4.5), dpi=130)
            ax.bar(x, monthly.values,
                   color=[_intensity_color(v) for v in monthly.values],
                   edgecolor="none")
            ax.axhline(annual_mean, color="#ef4444", linewidth=1.4, linestyle="--")
            ax.text(10.6, annual_mean * 1.02, f"Mean: {annual_mean:.0f} mm",
                    ha="right", va="bottom", fontsize=8, color="#ef4444")
            ax.set_xticks(x)
            ax.set_xticklabels(ml, fontsize=10)
            ax.set_ylabel("Rainfall (mm)", fontsize=11, fontweight="bold")
            ax.set_title(f"Monthly Rainfall Totals – {year}",
                         fontsize=13, fontweight="bold", pad=10, color="#333")
            legend_h = [
                mpatches.Patch(color="#93c5fd", label="< 50 mm"),
                mpatches.Patch(color="#3b82f6", label="50 – 150 mm"),
                mpatches.Patch(color="#1d4ed8", label="150 – 300 mm"),
                mpatches.Patch(color="#1e3a5f", label="≥ 300 mm"),
            ]
            ax.legend(handles=legend_h, loc="upper center",
                      bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle="--", axis="y")
            ax.set_facecolor("#fafafa")
            fig.patch.set_facecolor("white")
            plt.tight_layout()
            tmp = _save_fig(fig); plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP),
                                     width=Inches(SW - 0.54), height=Inches(CHART_H))
            os.unlink(tmp)

            _kpi_row(slide, [
                ("Annual Total",  f"{annual_total:.0f} mm",           "Full year",          "#1d4ed8"),
                ("Wettest Month", ml[wettest_idx - 1],                 f"{monthly[wettest_idx]:.0f} mm", "#1e3a5f"),
                ("Mean Monthly",  f"{annual_mean:.0f} mm",            "Annual ÷ 12",        "#0891b2"),
            ], top=KPI_TOP, card_h=KPI_H)

        except Exception as e:
            _err(slide, e)
        _logo(slide)

    _monthly_rainfall()

    # ── SLIDE 2: RAINY DAYS ───────────────────────────────────────────────────
    def _rainy_days():
        slide = prs.slides.add_slide(BLANK)
        _title(slide, f"Rainy Days by Intensity  —  {station_name}, {year}")
        _divider(slide)

        try:
            rainy = df_f[df_f["prcp_mm"] > 0].copy()

            def _cnt(lo, hi=None):
                m = rainy["prcp_mm"] >= lo
                if hi is not None:
                    m &= rainy["prcp_mm"] < hi
                return rainy[m].groupby("month").size().reindex(range(1, 13), fill_value=0)

            light    = _cnt(0.001, 10)
            moderate = _cnt(10, 25)
            heavy    = _cnt(25, heavy_rain_threshold)
            extreme  = _cnt(heavy_rain_threshold)

            fig, ax = plt.subplots(figsize=(13, 4.5), dpi=130)
            ax.bar(x, light.values,    color="#bfdbfe", label="Light (< 10 mm)")
            ax.bar(x, moderate.values, bottom=light.values,
                   color="#3b82f6",  label="Moderate (10–25 mm)")
            ax.bar(x, heavy.values,    bottom=(light + moderate).values,
                   color="#1d4ed8",  label=f"Heavy (25–{heavy_rain_threshold:.0f} mm)")
            ax.bar(x, extreme.values,  bottom=(light + moderate + heavy).values,
                   color="#ef4444",  label=f"Extreme (≥ {heavy_rain_threshold:.0f} mm)")
            ax.set_xticks(x)
            ax.set_xticklabels(ml, fontsize=10)
            ax.set_ylabel("Number of Days", fontsize=11, fontweight="bold")
            ax.set_title(f"Rainy Days per Month by Intensity – {year}",
                         fontsize=13, fontweight="bold", pad=10, color="#333")
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.10),
                      ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle="--", axis="y")
            ax.set_facecolor("#fafafa")
            fig.patch.set_facecolor("white")
            plt.tight_layout()
            tmp = _save_fig(fig); plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP),
                                     width=Inches(SW - 0.54), height=Inches(CHART_H))
            os.unlink(tmp)

            total_rainy   = int(rainy.shape[0])
            n_light       = int(light.sum())
            n_mod         = int(moderate.sum())
            n_heavy       = int(heavy.sum())
            n_extreme     = int(extreme.sum())

            _kpi_row(slide, [
                ("Total Rainy Days", str(total_rainy),              "Days with rain > 0",              "#1d4ed8"),
                ("Light  < 10 mm",   str(n_light),                  "Days",                            "#0891b2"),
                ("Moderate 10–25",   str(n_mod),                    "Days",                            "#2563eb"),
                (f"Heavy 25–{heavy_rain_threshold:.0f}mm", str(n_heavy),  "Days",                    "#1e3a5f"),
                (f"Extreme ≥{heavy_rain_threshold:.0f}mm", str(n_extreme),"Days",                    "#dc2626"),
            ], top=KPI_TOP, card_h=KPI_H)

        except Exception as e:
            _err(slide, e)
        _logo(slide)

    _rainy_days()

    # ── SLIDE 3: SURFACE RUNOFF ───────────────────────────────────────────────
    def _surface_runoff():
        # 2 KPI rows → smaller chart
        CHART_H2 = 3.70
        ROW1_TOP = CHART_TOP + CHART_H2 + 0.13   # 4.55"
        ROW2_TOP = ROW1_TOP + 0.92 + 0.10         # 5.57"
        ROW_H    = 0.88

        slide = prs.slides.add_slide(BLANK)
        _title(slide, f"Surface Runoff Analysis  —  {station_name}, {year}")
        _divider(slide)

        try:
            monthly_prcp = df_f.groupby("month")["prcp_mm"].sum().reindex(range(1, 13), fill_value=0)
            monthly_vols = {}
            for surf in _RUNOFF_SURFACES:
                a = surface_areas.get(surf["key"], 0.0)
                monthly_vols[surf["key"]] = (monthly_prcp / 1000.0) * a * surf["rc"]

            fig, ax = plt.subplots(figsize=(13, 4.0), dpi=130)
            bottom = np.zeros(12)
            for surf in _RUNOFF_SURFACES:
                vals = monthly_vols[surf["key"]].values
                ax.bar(x, vals, bottom=bottom, color=surf["color"],
                       label=f"{surf['label']}  (RC {surf['rc']:.2f})")
                bottom += vals
            ax.set_xticks(x)
            ax.set_xticklabels(ml, fontsize=10)
            ax.set_ylabel("Runoff Volume (m³)", fontsize=11, fontweight="bold")
            ax.set_title(f"Monthly Surface Runoff by Type – {year}",
                         fontsize=13, fontweight="bold", pad=10, color="#333")
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.13),
                      ncol=2, frameon=True, fontsize=8)
            ax.grid(True, alpha=0.25, linestyle="--", axis="y")
            ax.set_facecolor("#fafafa")
            fig.patch.set_facecolor("white")
            plt.tight_layout()
            tmp = _save_fig(fig); plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP),
                                     width=Inches(SW - 0.54), height=Inches(CHART_H2))
            os.unlink(tmp)

            surf_df      = pd.DataFrame({s["key"]: monthly_vols[s["key"]] for s in _RUNOFF_SURFACES})
            total_m      = surf_df.sum(axis=1)
            total_annual = float(total_m.sum())
            peak_m       = int(total_m.idxmax()) if total_annual > 0 else 1

            # Row 1 — aggregate (2 large cards)
            _kpi_row(slide, [
                ("Total Annual Runoff",
                 f"{int(round(total_annual * 1000))} L",
                 "All surfaces combined",
                 "#1d4ed8"),
                ("Peak Runoff Month",
                 ml[peak_m - 1],
                 f"{int(round(float(total_m[peak_m]) * 1000))} L",
                 "#dc2626"),
            ], top=ROW1_TOP, card_h=ROW_H)

            # Row 2 — per-surface breakdown (4 cards, colors from _RUNOFF_SURFACES)
            short = ["Roof", "Paved", "Green Area", "Waterbody"]
            _kpi_row(slide, [
                (f"{short[i]}  RC {s['rc']:.2f}",
                 f"{int(round(float(monthly_vols[s['key']].sum()) * 1000))} L",
                 "Annual runoff",
                 s["color"])
                for i, s in enumerate(_RUNOFF_SURFACES)
            ], top=ROW2_TOP, card_h=ROW_H)

        except Exception as e:
            _err(slide, e)
        _logo(slide)

    _surface_runoff()

    # ── SLIDE 4: GI BALANCE ───────────────────────────────────────────────────
    def _gi_balance():
        slide = prs.slides.add_slide(BLANK)
        _title(slide, f"GI Recharge Balance  —  {station_name}, {year}")
        _divider(slide)

        try:
            gi_result = _fetch_percentile_depth(station_id, gi_percentile, gi_start_year)
            if "error" in gi_result:
                _err(slide, f"Baseline fetch failed: {gi_result['error']}")
                _logo(slide)
                return

            baseline_mm = gi_result["raw_mm"]
            daily = df.sort_values("date").copy()
            daily["stored"]   = daily["prcp_mm"].clip(upper=baseline_mm)
            daily["overflow"] = (daily["prcp_mm"] - baseline_mm).clip(lower=0)

            mgrp = daily.groupby("month").agg(
                stored=("stored",   "sum"),
                overflow=("overflow", "sum"),
            ).reindex(range(1, 13), fill_value=0)

            total_recharge = float(daily["stored"].sum())
            total_overflow = float(daily["overflow"].sum())
            overflow_days  = int((daily["overflow"] > 0).sum())
            worst_m_idx    = int(mgrp["overflow"].idxmax()) if total_overflow > 0 else 1

            fig, ax = plt.subplots(figsize=(13, 4.5), dpi=130)
            ax.bar(x,  mgrp["stored"].values,   color="#22c55e", label="Stored (L/m²)")
            ax.bar(x, -mgrp["overflow"].values,  color="#ef4444", label="Overflow (L/m²)")
            ax.axhline(0, color="#374151", linewidth=1.2, linestyle="--")
            ax.set_xticks(x)
            ax.set_xticklabels(ml, fontsize=10)
            ax.set_ylabel("Volume (L/m²)", fontsize=11, fontweight="bold")
            ax.set_title(
                f"Monthly GI Recharge Balance – {year}  "
                f"({gi_percentile}th-pctile baseline: {baseline_mm:.1f} mm/day)",
                fontsize=12, fontweight="bold", pad=10, color="#333",
            )
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.10),
                      ncol=2, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle="--", axis="y")
            ax.set_facecolor("#fafafa")
            fig.patch.set_facecolor("white")
            plt.tight_layout()
            tmp = _save_fig(fig); plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP),
                                     width=Inches(SW - 0.54), height=Inches(CHART_H))
            os.unlink(tmp)

            worst_overflow_val = int(round(float(mgrp.loc[worst_m_idx, "overflow"])))
            _kpi_row(slide, [
                ("Annual Recharge",   f"{int(round(total_recharge))} L/m²",  "Total stored by GI",          "#16a34a"),
                ("Annual Overflow",   f"{int(round(total_overflow))} L/m²",  "Excess beyond capacity",      "#dc2626"),
                ("Overflow Days",     str(overflow_days),                      f"Rain > {baseline_mm:.0f} mm", "#d97706"),
                ("Worst Month",       ml[worst_m_idx - 1],                     f"{worst_overflow_val} L/m² overflow", "#7c3aed"),
            ], top=KPI_TOP, card_h=KPI_H)

        except Exception as e:
            _err(slide, e)
        _logo(slide)

    _gi_balance()

    # ── SAVE ─────────────────────────────────────────────────────────────────
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
