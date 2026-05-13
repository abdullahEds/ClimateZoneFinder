"""Thermal Comfort PowerPoint report generation - PPT slides for thermal comfort analysis."""

import io
import os
import tempfile
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ──────────────────────────────────────────────────────────────────────────────
# Thermal Comfort Visualization Functions (non-Streamlit)
# ──────────────────────────────────────────────────────────────────────────────

def compute_psychrometric_simple(df: pd.DataFrame) -> pd.DataFrame:
    """Compute humidity ratio and wet-bulb temperature for PPT generation."""
    P_ATM = 101_325.0
    
    out = df.copy()
    T = pd.to_numeric(out["dry_bulb_temperature"], errors="coerce")
    RH = pd.to_numeric(out["relative_humidity"], errors="coerce").clip(0, 100)
    
    # Saturation vapour pressure [Pa] – Magnus-Tetens
    e_s = 611.2 * np.exp(17.67 * T / (T + 243.5))
    # Actual vapour pressure [Pa]
    e = RH / 100.0 * e_s
    # Humidity ratio [kg/kg]
    hr = 0.62198 * e / (P_ATM - e)
    out["humidity_ratio"] = hr.clip(lower=0.0)
    
    # Wet-bulb temperature (Stull 2011)
    rh_safe = RH.clip(0.5, 100)
    T_wb = (
        T * np.arctan(0.151977 * np.sqrt(rh_safe + 8.313659))
        + np.arctan(T + rh_safe)
        - np.arctan(rh_safe - 1.676331)
        + 0.00391838 * rh_safe ** 1.5 * np.arctan(0.023101 * rh_safe)
        - 4.686035
    )
    out["wet_bulb_temperature"] = T_wb
    
    return out


def compute_adaptive_comfort_simple(df: pd.DataFrame) -> pd.DataFrame:
    """Compute ASHRAE 55 adaptive comfort for PPT generation."""
    ALPHA = 0.9
    
    out = df.copy()
    T = pd.to_numeric(out["dry_bulb_temperature"], errors="coerce")
    
    if "month" not in out.columns:
        out["month"] = out["datetime"].dt.month
    if "doy" not in out.columns:
        out["doy"] = out["datetime"].dt.dayofyear
    if "hour" not in out.columns:
        out["hour"] = out["datetime"].dt.hour
    
    # Daily mean temperature
    daily_mean = out.groupby("doy")["dry_bulb_temperature"].mean()
    
    # Exponential running mean: T_pma(d) = (1–α)·T_d + α·T_pma(d-1)
    doy_sorted = sorted(daily_mean.index)
    t_pma_daily = {}
    prev = daily_mean[doy_sorted[0]]
    for d in doy_sorted:
        val = (1 - ALPHA) * daily_mean[d] + ALPHA * prev
        t_pma_daily[d] = val
        prev = val
    pma_series = pd.Series(t_pma_daily, name="t_pma")
    
    # Map back to hourly
    out["t_pma"] = out["doy"].map(pma_series)
    
    # Comfort temperature
    out["t_comf"] = 0.31 * out["t_pma"] + 17.8
    out["t_comf_80_lo"] = out["t_comf"] - 3.5
    out["t_comf_80_hi"] = out["t_comf"] + 3.5
    out["t_comf_90_lo"] = out["t_comf"] - 2.5
    out["t_comf_90_hi"] = out["t_comf"] + 2.5
    
    # Applicability (ASHRAE 55 limits prevailing mean to 10–33.5 °C)
    out["adaptive_applicable"] = out["t_pma"].between(10.0, 33.5)
    
    # Whether current hour is within bands
    out["in_80"] = (T >= out["t_comf_80_lo"]) & (T <= out["t_comf_80_hi"]) & out["adaptive_applicable"]
    out["in_90"] = (T >= out["t_comf_90_lo"]) & (T <= out["t_comf_90_hi"]) & out["adaptive_applicable"]
    
    return out


def classify_comfort_simple(df: pd.DataFrame) -> pd.DataFrame:
    """Classify comfort categories for PPT generation."""
    out = df.copy()
    T = pd.to_numeric(out["dry_bulb_temperature"], errors="coerce")
    RH = pd.to_numeric(out["relative_humidity"], errors="coerce")
    
    conditions = [
        (T >= 22) & (T <= 26) & (RH >= 30) & (RH <= 60),  # Comfortable
        T > 26,                                             # Too Hot
        T < 22,                                             # Too Cold
        (RH > 60) & (T >= 22) & (T <= 26),                # Too Humid
        (RH < 30) & (T >= 22) & (T <= 26),                # Too Dry
    ]
    choices = ["Comfortable", "Too Hot", "Too Cold", "Too Humid", "Too Dry"]
    out["comfort_cat"] = np.select(conditions, choices, default="Too Hot")
    return out


def plot_comfort_heatmap(df: pd.DataFrame) -> plt.Figure:
    """Create comfort category heatmap (Hour × Month)."""
    COMFORT_CATS = ["Comfortable", "Too Hot", "Too Cold", "Too Humid", "Too Dry"]
    COMFORT_COLORS_NUMERIC = {
        "Comfortable": 0,
        "Too Hot": 1,
        "Too Cold": 2,
        "Too Humid": 3,
        "Too Dry": 4,
    }
    
    cat_order = COMFORT_COLORS_NUMERIC
    df2 = df.copy()
    df2["cat_num"] = df2["comfort_cat"].map(cat_order).fillna(0)
    
    pivot = (
        df2.groupby(["month", "hour"])["cat_num"]
        .mean()
        .reset_index()
        .pivot(index="month", columns="hour", values="cat_num")
        .reindex(range(1, 13))
    )
    
    months_lbl = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    hours_lbl = [f"{h:02d}:00" for h in range(24)]
    
    fig, ax = plt.subplots(figsize=(13.5, 5.5), dpi=120)
    
    im = ax.imshow(pivot.values, aspect="auto", origin="upper", cmap="RdYlGn_r", vmin=0, vmax=4)
    
    ax.set_xticks(range(24))
    ax.set_xticklabels(hours_lbl, fontsize=8, rotation=45)
    ax.set_yticks(range(12))
    ax.set_yticklabels(months_lbl, fontsize=9)
    
    cbar = plt.colorbar(im, ax=ax, fraction=0.035, pad=0.03)
    cbar.set_ticks([0, 1, 2, 3, 4])
    cbar.set_ticklabels(["Comfortable", "Too Hot", "Too Cold", "Too Humid", "Too Dry"], fontsize=8)
    
    ax.set_xlabel("Hour of Day", fontsize=11, fontweight='bold')
    ax.set_ylabel("Month", fontsize=11, fontweight='bold')
    ax.set_title("Comfort Heatmap – Hour × Month (dominant category)", fontsize=12, fontweight='bold', pad=10, color='#333')
    
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    
    return fig


def plot_strategy_distribution(df: pd.DataFrame) -> plt.Figure:
    """Create bar chart of design strategy distribution."""
    STRATEGY_COLORS = {
        "Comfortable": "#2ecc71",
        "Natural Ventilation": "#3498db",
        "Evaporative Cooling": "#16a085",
        "Mechanical Cooling": "#e74c3c",
        "Heating": "#f39c12",
        "Night Flushing / Thermal Mass": "#8e44ad",
    }
    
    counts = df["strategy"].value_counts()
    total = counts.sum()
    pcts = (counts / total * 100).round(1)
    
    strategies = list(pcts.index)
    values = list(pcts.values)
    colors = [STRATEGY_COLORS.get(s, "#95a5a6") for s in strategies]
    
    fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
    
    bars = ax.barh(strategies, values, color=colors, edgecolor='#333', linewidth=0.7)
    
    # Add percentage labels
    for i, (bar, val) in enumerate(zip(bars, values)):
        ax.text(val + 1, bar.get_y() + bar.get_height()/2, f'{val:.1f}%', 
                va='center', fontsize=10, fontweight='bold')
    
    ax.set_xlabel("% of Hours", fontsize=11, fontweight='bold')
    ax.set_title("Passive Design Strategy Distribution", fontsize=12, fontweight='bold', pad=10, color='#333')
    ax.set_xlim(0, max(values) * 1.25)
    ax.grid(True, alpha=0.25, linestyle='--', axis='x')
    ax.set_facecolor('#fafafa')
    
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    
    return fig


def plot_degree_hours_monthly(df: pd.DataFrame) -> plt.Figure:
    """Create monthly cooling and heating degree hours chart."""
    CDH_BASE = 24.0
    HDH_BASE = 18.0
    
    T = pd.to_numeric(df["dry_bulb_temperature"], errors="coerce")
    month = df["month"]
    
    cdh = (T - CDH_BASE).clip(lower=0)
    hdh = (HDH_BASE - T).clip(lower=0)
    
    cdh_monthly = cdh.groupby(month).sum().reindex(range(1, 13), fill_value=0)
    hdh_monthly = hdh.groupby(month).sum().reindex(range(1, 13), fill_value=0)
    
    months_lbl = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    x = np.arange(12)
    
    fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
    
    bar_w = 0.35
    ax.bar(x - bar_w/2, cdh_monthly.values, bar_w, label='Cooling Degree Hours (CDH)', 
           color='#e74c3c', edgecolor='#333', linewidth=0.7)
    ax.bar(x + bar_w/2, hdh_monthly.values, bar_w, label='Heating Degree Hours (HDH)', 
           color='#3498db', edgecolor='#333', linewidth=0.7)
    
    ax.set_xticks(x)
    ax.set_xticklabels(months_lbl, fontsize=10)
    ax.set_ylabel("Degree Hours (°C·h)", fontsize=11, fontweight='bold')
    ax.set_title(f"Monthly Degree Hours (CDH base {CDH_BASE:.0f}°C | HDH base {HDH_BASE:.0f}°C)", 
                 fontsize=12, fontweight='bold', pad=10, color='#333')
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=2, frameon=True, fontsize=10)
    ax.grid(True, alpha=0.25, linestyle='--', axis='y')
    ax.set_facecolor('#fafafa')
    
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    
    return fig


def plot_adaptive_comfort_scatter(df: pd.DataFrame) -> plt.Figure:
    """Create adaptive comfort scatter plot with comfort bands."""
    fdf = df[df["adaptive_applicable"]].copy()
    
    if fdf.empty:
        fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
        ax.text(0.5, 0.5, 'Adaptive comfort model not applicable for this climate', 
                ha='center', va='center', fontsize=12, transform=ax.transAxes)
        fig.patch.set_facecolor('white')
        return fig
    
    # Sort for band line drawing
    sorted_pma = np.linspace(fdf["t_pma"].min() - 1, fdf["t_pma"].max() + 1, 200)
    comf_line = 0.31 * sorted_pma + 17.8
    
    fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
    
    # 80% and 90% comfort bands
    ax.fill_between(sorted_pma, comf_line - 3.5, comf_line + 3.5, alpha=0.15, 
                    color='#2ecc71', label='80% Acceptability Band')
    ax.fill_between(sorted_pma, comf_line - 2.5, comf_line + 2.5, alpha=0.25, 
                    color='#27ae60', label='90% Acceptability Band')
    
    # Comfort line
    ax.plot(sorted_pma, comf_line, color='#27ae60', linewidth=2.2, linestyle='--', 
            label='Comfort Temperature', zorder=3)
    
    # Data scatter - color by DBT
    scatter = ax.scatter(fdf["t_pma"], fdf["dry_bulb_temperature"], 
                         c=fdf["dry_bulb_temperature"], cmap='RdYlBu_r', 
                         s=5, alpha=0.4, edgecolors='none', vmin=fdf["dry_bulb_temperature"].min(),
                         vmax=fdf["dry_bulb_temperature"].max())
    
    cbar = plt.colorbar(scatter, ax=ax, fraction=0.035, pad=0.03)
    cbar.set_label('DBT (°C)', fontsize=10, fontweight='bold')
    
    ax.set_xlabel("Prevailing Mean Outdoor Temperature T_pma (°C)", fontsize=11, fontweight='bold')
    ax.set_ylabel("Dry Bulb Temperature (°C)", fontsize=11, fontweight='bold')
    ax.set_title("Adaptive Comfort – ASHRAE 55 (T_pma vs. Indoor Temperature)", 
                 fontsize=12, fontweight='bold', pad=10, color='#333')
    ax.legend(loc='upper left', fontsize=10, frameon=True)
    ax.grid(True, alpha=0.25, linestyle='--')
    ax.set_facecolor('#fafafa')
    
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    
    return fig


def plot_comfort_percentages(df: pd.DataFrame) -> plt.Figure:
    """Create summary chart of comfort statistics."""
    try:
        # Compute statistics
        pct_comfortable = (df["comfort_cat"] == "Comfortable").sum() / len(df) * 100
        pct_in_80 = df["in_80"].sum() / len(df[df["adaptive_applicable"]]) * 100 if df["adaptive_applicable"].sum() > 0 else 0
        pct_in_90 = df["in_90"].sum() / len(df[df["adaptive_applicable"]]) * 100 if df["adaptive_applicable"].sum() > 0 else 0
        
        labels = ["% Comfortable\n(Static Zone)", "% in ASHRAE\n80% Band", "% in ASHRAE\n90% Band"]
        values = [pct_comfortable, pct_in_80, pct_in_90]
        colors = ["#2ecc71", "#3498db", "#9b59b6"]
        
        fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
        
        bars = ax.bar(labels, values, color=colors, edgecolor='#333', linewidth=1.5, width=0.5)
        
        # Add percentage labels on bars
        for bar, val in zip(bars, values):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                    f'{val:.1f}%', ha='center', va='bottom', fontsize=13, fontweight='bold')
        
        ax.set_ylabel("Percentage of Hours (%)", fontsize=11, fontweight='bold')
        ax.set_ylim(0, max(values) * 1.2)
        ax.set_title("Thermal Comfort Performance Summary", fontsize=12, fontweight='bold', pad=10, color='#333')
        ax.grid(True, alpha=0.25, linestyle='--', axis='y')
        ax.set_facecolor('#fafafa')
        
        fig.patch.set_facecolor('white')
        plt.tight_layout()
        
        return fig
    except Exception:
        fig, ax = plt.subplots(figsize=(13, 5.5), dpi=120)
        ax.text(0.5, 0.5, 'Error generating comfort chart', 
                ha='center', va='center', fontsize=12, transform=ax.transAxes)
        fig.patch.set_facecolor('white')
        return fig


# ──────────────────────────────────────────────────────────────────────────────
# Main Thermal Comfort PPT Report Generator
# ──────────────────────────────────────────────────────────────────────────────

def generate_thermal_comfort_pptx_report(
    df: pd.DataFrame,
    metadata: dict,
):
    """Generate a Thermal Comfort PowerPoint report using the Voha template."""
    
    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    except NameError:
        base_dir = os.getcwd()
    
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")
    
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # Remove all slides except keep layouts
        from pptx.oxml.ns import qn
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    BLANK_LAYOUT = prs.slide_layouts[6]
    TITLE_RED = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY = RGBColor(0x40, 0x40, 0x40)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    SW = prs.slide_width.inches
    SH = prs.slide_height.inches
    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12
    
    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(LOGO_L), Inches(LOGO_T),
                                     width=Inches(LOGO_W), height=Inches(LOGO_H))
    
    def _slide_title(slide, text, top=0.13):
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(top), Inches(SW - 0.54), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED
    
    def _divider(slide, top_inches):
        bar = slide.shapes.add_shape(1, Inches(0.27), Inches(top_inches),
                                     Inches(SW - 0.54), Inches(0.03))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TITLE_RED
        bar.line.fill.background()
    
    def _err(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = TITLE_RED
    
    def _save_fig(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches="tight", facecolor="white")
            return tmp.name
    
    # Ensure necessary columns exist
    if "month" not in df.columns:
        df["month"] = df["datetime"].dt.month
    if "doy" not in df.columns:
        df["doy"] = df["datetime"].dt.dayofyear
    if "hour" not in df.columns:
        df["hour"] = df["datetime"].dt.hour
    
    # Compute thermal comfort data
    try:
        df_thermal = compute_psychrometric_simple(df)
        df_thermal = compute_adaptive_comfort_simple(df_thermal)
        df_thermal = classify_comfort_simple(df_thermal)
        
        # Add strategy column (simplified for PPT - just comfort classification for now)
        df_thermal["strategy"] = df_thermal["comfort_cat"]
    except Exception as e:
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Thermal Comfort Analysis")
        _divider(slide, 0.62)
        _err(slide, f"Data processing error: {str(e)[:50]}")
        _add_logo(slide)
        
        report_bytes = io.BytesIO()
        prs.save(report_bytes)
        report_bytes.seek(0)
        return report_bytes
    
    # ── COVER SLIDE ───────────────────────────────────────────────────────────
    def _cover():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(SW), Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(SW - 1.2), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Thermal Comfort Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.75), Inches(SW - 1.2), Inches(0.65))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        run2.text = f"Location: {_city}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
        
        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Comfort Heatmap | Design Strategies | Degree Hours | Adaptive Comfort | Performance Summary"
        run3.font.size = Pt(10)
        run3.font.color.rgb = DARK_GREY
        
        _add_logo(slide)
    
    _cover()
    
    # ── COMFORT HEATMAP SLIDE ─────────────────────────────────────────────────
    def _comfort_heatmap_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Comfort Heatmap – Hour × Month")
        _divider(slide, 0.62)
        
        try:
            fig = plot_comfort_heatmap(df_thermal)
            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, f"Comfort heatmap: {str(e)[:50]}")
        
        _add_logo(slide)
    
    _comfort_heatmap_slide()
    
    # ── STRATEGY DISTRIBUTION SLIDE ───────────────────────────────────────────
    def _strategy_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Design Strategy Opportunity Distribution")
        _divider(slide, 0.62)
        
        try:
            fig = plot_strategy_distribution(df_thermal)
            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, f"Strategy chart: {str(e)[:50]}")
        
        _add_logo(slide)
    
    _strategy_slide()
    
    # ── DEGREE HOURS SLIDE ────────────────────────────────────────────────────
    def _degree_hours_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Monthly Degree Hours – Cooling & Heating Demand")
        _divider(slide, 0.62)
        
        try:
            fig = plot_degree_hours_monthly(df_thermal)
            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, f"Degree hours: {str(e)[:50]}")
        
        _add_logo(slide)
    
    _degree_hours_slide()
    
    # ── ADAPTIVE COMFORT SLIDE ────────────────────────────────────────────────
    def _adaptive_comfort_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Adaptive Comfort – ASHRAE 55 Analysis")
        _divider(slide, 0.62)
        
        try:
            fig = plot_adaptive_comfort_scatter(df_thermal)
            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, f"Adaptive comfort: {str(e)[:50]}")
        
        _add_logo(slide)
    
    _adaptive_comfort_slide()
    
    # ── COMFORT PERFORMANCE SUMMARY SLIDE ─────────────────────────────────────
    def _performance_summary_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Thermal Comfort Performance Summary")
        _divider(slide, 0.62)
        
        try:
            fig = plot_comfort_percentages(df_thermal)
            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, f"Performance: {str(e)[:50]}")
        
        _add_logo(slide)
    
    _performance_summary_slide()
    
    # ── DESIGN RECOMMENDATIONS SLIDE ──────────────────────────────────────────
    def _recommendations_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Design Recommendations & Strategies")
        _divider(slide, 0.62)
        
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.75), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Calculate summary statistics
        try:
            pct_comfortable = (df_thermal["comfort_cat"] == "Comfortable").sum() / len(df_thermal) * 100
            pct_hot = (df_thermal["comfort_cat"] == "Too Hot").sum() / len(df_thermal) * 100
            pct_cold = (df_thermal["comfort_cat"] == "Too Cold").sum() / len(df_thermal) * 100
            mean_rh = df_thermal["relative_humidity"].mean()
            
            p = tf.paragraphs[0]
            p.text = "Climate Comfort Profile"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TITLE_RED
            p.space_after = Pt(6)
            
            p = tf.add_paragraph()
            p.text = f"• Comfortable hours: {pct_comfortable:.1f}%  |  Overheating: {pct_hot:.1f}%  |  Undercooling: {pct_cold:.1f}%"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = f"• Mean Relative Humidity: {mean_rh:.1f}%"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(10)
            
            p = tf.add_paragraph()
            p.text = "Key Design Strategies"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TITLE_RED
            p.space_before = Pt(4)
            p.space_after = Pt(6)
            
            if pct_comfortable < 40:
                if pct_hot > pct_cold:
                    p = tf.add_paragraph()
                    p.text = "• Priority: Cooling strategies – Implement high-performance envelope, external shading, and natural ventilation"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                    
                    p = tf.add_paragraph()
                    p.text = "• Consider nighttime cooling recovery and thermal mass activation"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                else:
                    p = tf.add_paragraph()
                    p.text = "• Priority: Heating strategies – Maximize solar heat gain during winter with south-facing glazing"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                    
                    p = tf.add_paragraph()
                    p.text = "• Ensure robust thermal insulation and minimize infiltration losses"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
            else:
                p = tf.add_paragraph()
                p.text = "• Climate is generally favorable – Prioritize passive design with natural ventilation and daylighting"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(3)
            
            if mean_rh > 65:
                p = tf.add_paragraph()
                p.text = "• High humidity detected – Ensure adequate dehumidification and mold risk mitigation"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(3)
            elif mean_rh < 30:
                p = tf.add_paragraph()
                p.text = "• Low humidity detected – Humidification may be required in heating season for occupant comfort"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(3)
            
            p = tf.add_paragraph()
            p.text = "• Adaptive Comfort – Leverage occupant behavior (clothing, behavior) to expand acceptable temperature ranges"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(0)
            
        except Exception as e:
            p = tf.paragraphs[0]
            p.text = f"Error generating recommendations: {str(e)[:40]}"
            p.font.size = Pt(11)
            p.font.color.rgb = TITLE_RED
        
        _add_logo(slide)
    
    _recommendations_slide()
    
    # ── ANNEXURE SLIDE ────────────────────────────────────────────────────────
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Annexure")
        _divider(slide, 0.62)
        
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = (
            "Environmental Design Solutions [EDS] is a sustainability advisory firm. "
            "Since 2002, EDS has worked on over 500 green building and energy efficiency "
            "projects worldwide. The team focuses on climate change mitigation, low-carbon "
            "design, building simulation, performance audits, and capacity building. EDS "
            "continues to contribute to the buildings community with useful tools through "
            "its IT services."
        )
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        
        for item in [
            "• Thermal comfort analysis based on ASHRAE 55 standard",
            "• Psychrometric calculations using Magnus-Tetens formula",
            "• Adaptive comfort model per ASHRAE 55-2017 §5.4",
            "• Streamlit, © Streamlit Inc., licensed under Apache 2.0",
            "• Python © Python Software Foundation, licensed under PSF License Version 2",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)
        
        _add_logo(slide)
    
    _make_annexure_slide()
    
    # Save and return report
    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes
