"""Combined PowerPoint report generation - Climate + Shading + Wind Analysis."""

import concurrent.futures
import io
import os
import tempfile
from datetime import datetime, date
from typing import Optional
# from tkinter import SW

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pytz

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .shading_helpers import (
    _ORIENTATIONS,
    build_thermal_matrix,
    compute_shading_geometry,
    build_orientation_table,
    compute_solar_angles,
    get_overheating_hours,
)


def _ppt_remove_all_slides(prs):
    """Remove every slide from an open Presentation while preserving theme/layouts."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)


def generate_combined_pptx_report(
    df: pd.DataFrame,
    start_date: date,
    end_date: date,
    start_hour: int,
    end_hour: int,
    selected_parameter: str,
    metadata: Optional[dict] = None,
    temp_threshold: float = 28.0,
    rad_threshold: float = 315.0,
    design_cutoff_angle: float = 45.0,
    n_sectors: int = 16,
    include_thermal_comfort: bool = False,
    rainfall_station_name: Optional[str] = None,
    rainfall_station_id: Optional[str] = None,
    rainfall_year: Optional[int] = None,
    rainfall_start_month: int = 1,
    rainfall_end_month: int = 12,
    rainfall_heavy_threshold: float = 50.0,
    rainfall_surface_areas: Optional[dict[str, float]] = None,
    rainfall_gi_percentile: int = 95,
    rainfall_gi_start_year: int = 1990,
    branding: Optional[dict] = None,
    solar_pv_country: str = "India",
    solar_pv_roof_size_m2: float = 100.0,
    solar_pv_roof_pct: float = 80.0,
) -> io.BytesIO:
    """Generate a combined PowerPoint report with Climate + Shading Analysis + Assumptions slide."""

    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]

    # ── Background NOAA fetch: starts now, completes before rainfall slides ────
    _noaa_future = None
    _percentile_future = None
    _rain_executor = None
    if rainfall_station_id and rainfall_year:
        try:
            from .rainfall_module import (
                _fetch_noaa as _bg_fetch_noaa,
                _fetch_percentile_depth as _bg_fetch_perc,
            )
            # Use __wrapped__ to bypass @st.cache_data — it fails in non-Streamlit threads
            _raw_fetch_noaa = getattr(_bg_fetch_noaa, '__wrapped__', _bg_fetch_noaa)
            _raw_fetch_perc = getattr(_bg_fetch_perc, '__wrapped__', _bg_fetch_perc)
            _rain_executor = concurrent.futures.ThreadPoolExecutor(max_workers=2)
            _noaa_future = _rain_executor.submit(
                _raw_fetch_noaa, rainfall_station_id, rainfall_year
            )
            _percentile_future = _rain_executor.submit(
                _raw_fetch_perc,
                rainfall_station_id,
                rainfall_gi_percentile,
                rainfall_gi_start_year,
            )
        except Exception:
            pass

    # ── pvlib solar-position cache: compute once, reuse across sun-path slides ─
    _sol_pos_cache = {}

    def _get_solar_positions(lat, lon, tz_str):
        key = (lat, lon, str(tz_str))
        if key not in _sol_pos_cache:
            from pvlib import solarposition as _sp_lib
            try:
                _tz = pytz.timezone(str(tz_str))
            except Exception:
                _tz = pytz.UTC
            times = pd.date_range(
                "2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left"
            )
            _sol_pos_cache[key] = (_tz, times, _sp_lib.get_solarposition(times, lat, lon))
        return _sol_pos_cache[key]

    TITLE_RED   = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER_CLR = RGBColor(0xC0, 0x00, 0x00)
    LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches

    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(LOGO_L), Inches(LOGO_T),
                width=Inches(LOGO_W), height=Inches(LOGO_H),
            )

    def _add_slide_title(slide, text, left=0.27, top=0.13, width=None, height=0.45):
        if width is None:
            width = SW - left - 0.3
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _add_divider(slide, top_inches):
        line = slide.shapes.add_shape(
            1, Inches(0.27), Inches(top_inches),
            Inches(SW - 0.54), Inches(0.03),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = DIVIDER_CLR
        line.line.fill.background()

    def _save_mpl_figure(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=100, bbox_inches='tight', facecolor='white')
            return tmp.name

    def _err_box(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Visualization error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    filtered_df = df[
        (df["datetime"].dt.date >= start_date) &
        (df["datetime"].dt.date <= end_date) &
        (df["hour"].between(start_hour, end_hour))
    ]

    # ── COVER SLIDE ───────────────────────────────────────────────────────────
    def _make_cover_slide():
        _branding = branding or {}
        slide = prs.slides.add_slide(BLANK_LAYOUT)

        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(SW - 1.2), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Climate & Shading Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.7))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        _location = metadata.get("location", "") if metadata else ""
        location_display = _city if _city else (_location if _location else "Location")
        project_display = _branding.get("project_name") or location_display
        run2.text = project_display
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        # Client name (if provided)
        if _branding.get("client_name"):
            tb_client = slide.shapes.add_textbox(Inches(0.6), Inches(4.55), Inches(SW - 1.2), Inches(0.45))
            run_c = tb_client.text_frame.paragraphs[0].add_run()
            run_c.text = f"Prepared for: {_branding['client_name']}"
            run_c.font.size = Pt(12)
            run_c.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        _report_date = _branding.get("report_date") or datetime.now().strftime("%d %B %Y")
        run3.text = _report_date
        run3.font.size = Pt(11)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _make_cover_slide()

    # ── ASSUMPTIONS SLIDE ─────────────────────────────────────────────────────
    def _make_assumptions_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Assumptions & Analysis Parameters")
        _add_divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.75), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        # Header
        p = tf.paragraphs[0]
        p.text = "Default Conditions & Selected Parameters"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(8)

        # Analysis Period
        p = tf.add_paragraph()
        p.text = "Analysis Period"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(2)
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f"• Date Range: {start_date.strftime('%d %b')} to {end_date.strftime('%d %b')}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = f"• Hour Range: {start_hour:02d}:00 to {end_hour:02d}:00"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(6)

        # Location & Climate Data
        p = tf.add_paragraph()
        p.text = "Location & Climate Data"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(2)
        p.space_after = Pt(4)

        _meta = metadata or {}
        p = tf.add_paragraph()
        p.text = f"• Location: {_meta.get('city', 'Unknown')}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        lat = _meta.get('latitude')
        lon = _meta.get('longitude')
        if lat is not None and lon is not None:
            p = tf.add_paragraph()
            p.text = f"• Coordinates: {lat:.2f}°N, {lon:.2f}°E"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = f"• Time Zone: {_meta.get('timezone', 'UTC')}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(6)

        # Thermal Comfort & Shading Thresholds
        p = tf.add_paragraph()
        p.text = "Thermal Comfort & Shading Thresholds"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(2)
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f"• Temperature Threshold (Overheating): {temp_threshold}°C"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = f"• Solar Radiation Threshold (GHI): {rad_threshold} W/m²"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = f"• Design Cutoff Angle (Shading): {design_cutoff_angle}°"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(6)

        # Thermal Comfort Standards
        p = tf.add_paragraph()
        p.text = "Thermal Comfort Standards Applied"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(2)
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = "• Comfort Band (Dry Bulb): 20-26°C (ASHRAE 90% acceptability)"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "• Relative Humidity (Comfortable): 30-60%"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "• Condensation Risk Threshold: RH > 75%"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GREY
        p.space_before = Pt(0)
        p.space_after = Pt(6)

        # Solar PV Parameters
        p = tf.add_paragraph()
        p.text = "Solar PV Assumptions"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(2)
        p.space_after = Pt(4)

        _eff_area = solar_pv_roof_size_m2 * (solar_pv_roof_pct / 100.0)
        _sys_kwp  = _eff_area / 10.0
        for _line in [
            f"• Country / Location: {solar_pv_country}",
            f"• Total Roof Size: {solar_pv_roof_size_m2:.0f} m²",
            f"• Solar Coverage: {solar_pv_roof_pct:.0f}% of roof area",
            f"• Effective PV Area: {_eff_area:.1f} m²  →  System Size: {_sys_kwp:.2f} kWp  (10 m² = 1 kWp)",
            "• Data Source: World Bank Global Solar Atlas 2.0 / SolarGIS PVOUT Level 1",
            "• Values represent long-term monthly average specific yield (kWh/kWp.day)",
        ]:
            p = tf.add_paragraph()
            p.text = _line
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_assumptions_slide()

    # ── SECTION 1 – DRY BULB TEMPERATURE ─────────────────────────────────────
    def _make_dbt_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                temp_min=("dry_bulb_temperature", "min"),
                temp_max=("dry_bulb_temperature", "max"),
                temp_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
            comfort_line = daily_avg.rolling(window=7, center=True).mean()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=100)
            ax.fill_between(daily_stats["doy"], comfort_line - 3.5, comfort_line + 3.5,
                            alpha=0.18, color='gray', label='ASHRAE 80% Comfort')
            ax.fill_between(daily_stats["doy"], comfort_line - 2.5, comfort_line + 2.5,
                            alpha=0.28, color='gray', label='ASHRAE 90% Comfort')
            ax.fill_between(daily_stats["doy"], daily_stats["temp_min"], daily_stats["temp_max"],
                            alpha=0.30, color='#FFB3B3', label='Daily Temp Range')
            ax.plot(daily_stats["doy"], daily_stats["temp_avg"],
                    color='#C00000', linewidth=2.2, label='Daily Average', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Annual Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=5, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_trend_slide()

    def _make_dbt_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly = df.groupby("month").agg(
                t_min=("dry_bulb_temperature", "min"),
                t_max=("dry_bulb_temperature", "max"),
                t_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=100)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly["t_min"], bar_w, color='#90CAF9', label='Min Temp')
            ax.bar(x,          monthly["t_avg"], bar_w, color='#C00000', label='Avg Temp', alpha=0.85)
            ax.bar(x + bar_w,  monthly["t_max"], bar_w, color='#EF9A9A', label='Max Temp')

            ax.axhspan(20, 26, alpha=0.10, color='green', label='Comfort Band (20–26°C)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Monthly Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min: {filtered_df['dry_bulb_temperature'].min():.1f}°C  "
                    f"Avg: {filtered_df['dry_bulb_temperature'].mean():.1f}°C  "
                    f"Max: {filtered_df['dry_bulb_temperature'].max():.1f}°C  "
                    f" |  Ann. CDD24: {(df['dry_bulb_temperature'] - 24).clip(lower=0).sum():.0f}   "
                    f"HDD18: {(18 - df['dry_bulb_temperature']).clip(lower=0).sum():.0f}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_monthly_slide()

    # ── SECTION 2 – RELATIVE HUMIDITY ─────────────────────────────────────────
    def _make_rh_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=100)
            ax.axhspan(75, 100, alpha=0.13, color='#FF6B6B', label='Condensation Risk (>75%)')
            ax.axhspan(60,  75, alpha=0.13, color='#FFA500', label='High RH (60–75%)')
            ax.axhspan(30,  60, alpha=0.13, color='#4ECDC4', label='Comfortable (30–60%)')
            ax.axhspan( 0,  30, alpha=0.13, color='#FFD93D', label='Low RH (<30%)')

            ax.fill_between(daily_stats["doy"], daily_stats["rh_min"], daily_stats["rh_max"],
                            alpha=0.28, color='#0099ff', label='Daily RH Range')
            ax.plot(daily_stats["doy"], daily_stats["rh_avg"],
                    color='#0066cc', linewidth=2.2, label='Daily Average RH', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 100)
            ax.set_title('Annual Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_trend_slide()

    def _make_rh_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly_rh = df.groupby("month").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=100)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly_rh["rh_min"], bar_w, color='#AED6F1', label='Min RH')
            ax.bar(x,          monthly_rh["rh_avg"], bar_w, color='#0066cc', label='Avg RH', alpha=0.85)
            ax.bar(x + bar_w,  monthly_rh["rh_max"], bar_w, color='#5DADE2', label='Max RH')

            ax.axhspan(30, 60, alpha=0.10, color='green', label='Comfortable (30–60%)')
            ax.axhline(75, color='#E74C3C', linewidth=1.2, linestyle='--', label='Condensation Threshold (75%)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 110)
            ax.set_title('Monthly Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min RH: {filtered_df['relative_humidity'].min():.0f}%  "
                    f"Avg RH: {filtered_df['relative_humidity'].mean():.0f}%  "
                    f"Max RH: {filtered_df['relative_humidity'].max():.0f}%  "
                    f" |  High RH hrs (>60%): {len(filtered_df[filtered_df['relative_humidity'] > 60])}  "
                    f"Condensation risk hrs (>75%): {len(filtered_df[filtered_df['relative_humidity'] > 75])}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_monthly_slide()

    # ── SECTION 3 – SUN PATH ──────────────────────────────────────────────────
    # def _make_sun_path_slide():
        # slide = prs.slides.add_slide(BLANK_LAYOUT)
        # _add_slide_title(slide, "Sun Path Diagram")
        # _add_divider(slide, 0.62)

        # _meta = metadata or {}
        # lat = _meta.get("latitude")
        # lon = _meta.get("longitude")
        # tz_str = _meta.get("timezone", "UTC")

        # if lat is None or lon is None:
        #     _err_box(slide, "Latitude/Longitude not available from EPW metadata.")
        #     _add_logo(slide)
        #     return

        # try:
        #     from pvlib import solarposition as _solpos_lib

        #     try:
        #         _tz = pytz.timezone(tz_str)
        #     except Exception:
        #         _tz = pytz.UTC

        #     times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left")
        #     sol = _solpos_lib.get_solarposition(times, lat, lon)
        #     sol = sol[sol["apparent_elevation"] > 0].copy()
        #     sol["r"] = 90 - sol["apparent_elevation"]

        #     fig = plt.figure(figsize=(9, 7.5), dpi=100, facecolor='white')
        #     ax = fig.add_subplot(111, projection='polar')
        #     ax.set_theta_zero_location('N')
        #     ax.set_theta_direction(-1)
        #     ax.set_aspect('equal', adjustable='box')
        #     ax.set_ylim(0, 90)
        #     ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
        #     ax.set_yticklabels(['90°\n(Zenith)', '75°', '60°', '45°', '30°', '15°', '0°\n(Horizon)'],
        #                        fontsize=7, color='#555')
        #     ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
        #     ax.set_xticklabels(['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'], fontsize=10, fontweight='bold')
        #     ax.set_facecolor('#F0F4F8')
        #     ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)

        #     sc = ax.scatter(
        #         np.radians(sol["azimuth"].values),
        #         sol["r"].values,
        #         c=sol.index.dayofyear,
        #         cmap='YlOrRd',
        #         s=1.0, alpha=0.55,
        #         vmin=1, vmax=365,
        #         linewidths=0, zorder=2,
        #     )
        #     cbar = fig.colorbar(sc, ax=ax, pad=0.10, fraction=0.035, shrink=0.75)
        #     cbar.set_label('Day of Year', fontsize=9)
        #     cbar.set_ticks([1, 91, 182, 273, 365])
        #     cbar.set_ticklabels(['1\n(Jan)', '91\n(Apr)', '182\n(Jul)', '273\n(Oct)', '365\n(Dec)'])

        #     key_dates = [
        #         ("Mar 21 (Spring Equinox)", "2020-03-21", "#FF9500", 1.6),
        #         ("Jun 21 (Summer Solstice)", "2020-06-21", "#CC0000", 2.0),
        #         ("Dec 21 (Winter Solstice)", "2020-12-21", "#0066CC", 2.0),
        #     ]
        #     for lbl, dstr, col, lw in key_dates:
        #         dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
        #         ks = _solpos_lib.get_solarposition(dt, lat, lon)
        #         ks = ks[ks["apparent_elevation"] > 0]
        #         if not ks.empty:
        #             ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
        #                     color=col, linewidth=lw, label=lbl, zorder=4)

        #     ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.06), ncol=3,
        #               frameon=True, fontsize=8, borderaxespad=0)
        #     ax.set_title(f'Sun Path  |  Lat: {lat:.2f}°  Lon: {lon:.2f}°',
        #                  fontsize=11, fontweight='bold', color='#333', pad=14)

        #     plt.tight_layout()
        #     tmp = _save_mpl_figure(fig)
        #     plt.close(fig)

        #     # Use square dimensions to maintain circular aspect ratio
        #     # img_size = min(SW * 0.55, SH * 0.75)
        #     # img_l = (SW - img_size) / 2
        #     img_size = min(SW*0.85 , SH*.75 )
        #     img_l = (SW - img_size) / 2
        #     img_t = 0.72
        #     slide.shapes.add_picture(tmp, Inches(img_l), Inches(img_t), width=Inches(img_size), height=Inches(img_size))
        #     os.unlink(tmp)

        # except Exception as e:
        #     _err_box(slide, e)

        # _add_logo(slide)
    def _make_sun_path_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Sun Path Diagram")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        lat = _meta.get("latitude")
        lon = _meta.get("longitude")
        tz_str = _meta.get("timezone", "UTC")

        if lat is None or lon is None:
            _err_box(slide, "Latitude/Longitude not available from EPW metadata.")
            _add_logo(slide)
            return

        try:
            from pvlib import solarposition as _solpos_lib
            _tz, times, _sol_full = _get_solar_positions(lat, lon, tz_str)
            sol = _sol_full[_sol_full["apparent_elevation"] > 0].copy()
            sol["r"] = 90 - sol["apparent_elevation"]

            # ---------- FIGURE ----------
            fig = plt.figure(figsize=(7.5, 7.5), dpi=100, facecolor='white')
            ax = fig.add_subplot(111, projection='polar')

            ax.set_theta_zero_location('N')
            ax.set_theta_direction(-1)
            ax.set_aspect('equal', adjustable='box')

            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(
                ['90°\n(Zenith)', '75°', '60°', '45°', '30°', '15°', '0°\n(Horizon)'],
                fontsize=7, color='#555'
            )

            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(
                ['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'],
                fontsize=10, fontweight='bold'
            )

            ax.set_facecolor('#F0F4F8')
            ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)

            # ---------- SCATTER ----------
            sc = ax.scatter(
                np.radians(sol["azimuth"].values),
                sol["r"].values,
                c=sol.index.dayofyear,
                cmap='YlOrRd',
                s=1.0,
                alpha=0.55,
                vmin=1,
                vmax=365,
                linewidths=0,
                zorder=2,
            )

            # ---------- COLORBAR ----------
            cbar = fig.colorbar(sc, ax=ax, pad=0.08, fraction=0.035, shrink=0.8)
            cbar.set_label('Day of Year', fontsize=9)
            cbar.set_ticks([1, 91, 182, 273, 365])
            cbar.set_ticklabels(['1\n(Jan)', '91\n(Apr)', '182\n(Jul)', '273\n(Oct)', '365\n(Dec)'])

            # ---------- KEY DATES ----------
            key_dates = [
                ("Mar 21 (Spring Equinox)", "2020-03-21", "#FF9500", 1.6),
                ("Jun 21 (Summer Solstice)", "2020-06-21", "#CC0000", 2.0),
                ("Dec 21 (Winter Solstice)", "2020-12-21", "#0066CC", 2.0),
            ]

            for lbl, dstr, col, lw in key_dates:
                dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
                ks = _solpos_lib.get_solarposition(dt, lat, lon)
                ks = ks[ks["apparent_elevation"] > 0]

                if not ks.empty:
                    ax.plot(
                        np.radians(ks["azimuth"]),
                        90 - ks["apparent_elevation"],
                        color=col,
                        linewidth=lw,
                        label=lbl,
                        zorder=4
                    )

            # ---------- LEGEND ----------
            ax.legend(
                loc='upper center',
                bbox_to_anchor=(0.5, -0.08),
                ncol=3,
                frameon=True,
                fontsize=8
            )

            # ---------- TITLE ----------
            ax.set_title(
                f'Sun Path  |  Lat: {lat:.2f}°  Lon: {lon:.2f}°',
                fontsize=11,
                fontweight='bold',
                color='#333',
                pad=14
            )

            # ---------- LAYOUT FIX (CRITICAL) ----------
            plt.tight_layout(pad=2.5)
            fig.subplots_adjust(left=0.08, right=0.88, top=0.92, bottom=0.12)

            # ---------- SAVE ----------
            tmp = _save_mpl_figure(fig)
            plt.close(fig)

            # ---------- PPT IMAGE PLACEMENT ----------
            img_size = min(SW * 0.75, SH * 0.75)   # square, no distortion
            img_l = (SW - img_size) / 2
            img_t = 0.72

            slide.shapes.add_picture(
                tmp,
                Inches(img_l),
                Inches(img_t),
                width=Inches(img_size),
                height=Inches(img_size)
            )

            os.unlink(tmp)

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)
    _make_sun_path_slide()

    # ── SECTION 4 – THERMAL & RADIATION MATRIX (Shading) ────────────────────
    def _make_thermal_matrix_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Thermal & Radiation Matrix (Shading Analysis)")
        _add_divider(slide, 0.62)

        try:
            temp_matrix, rad_matrix, overheat_mask = build_thermal_matrix(df, temp_threshold, rad_threshold)
            months_lbl = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
            hours_lbl  = [f"{h:02d}:00" for h in range(24)]

            fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.2), dpi=120)

            for ax, matrix, title, cmap, clabel in [
                (axes[0], temp_matrix, f"Mean Dry-Bulb Temp (°C)  [threshold: {temp_threshold}°C]", "RdYlBu_r", "°C"),
                (axes[1], rad_matrix,  f"Mean GHI (W/m²)  [threshold: {rad_threshold} W/m²]",     "YlOrRd",   "W/m²"),
            ]:
                im = ax.imshow(matrix.values, aspect="auto", origin="upper", cmap=cmap)
                plt.colorbar(im, ax=ax, fraction=0.035, pad=0.03, label=clabel)
                ax.set_xticks(range(12))
                ax.set_xticklabels(months_lbl, fontsize=8)
                ax.set_yticks(range(24))
                ax.set_yticklabels(hours_lbl, fontsize=7)
                ax.set_title(title, fontsize=10, fontweight="bold", pad=8)
                ax.set_xlabel("Month", fontsize=9)
                ax.set_ylabel("Hour of Day", fontsize=9)

                for h_i in range(24):
                    for m_i in range(12):
                        if overheat_mask.iloc[h_i, m_i]:
                            rect = plt.Rectangle(
                                (m_i - 0.5, h_i - 0.5), 1, 1,
                                fill=False, edgecolor="black", linewidth=1.6,
                            )
                            ax.add_patch(rect)

            fig.suptitle(
                "Overheating Hours  (black border = both thresholds exceeded)",
                fontsize=11, fontweight="bold", y=1.01, color="#333",
            )
            fig.patch.set_facecolor("white")
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_thermal_matrix_slide()

    # ── SECTION 5 – SUN PATH SHADING MODE ─────────────────────────────────────
    def _make_sun_path_shading_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Sun Path – Shading Analysis")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        _lat = _meta.get("latitude")
        _lon = _meta.get("longitude")
        _tz_str = _meta.get("timezone", "UTC")

        if _lat is None or _lon is None:
            _err_box(slide, "Latitude/Longitude not available.")
            _add_logo(slide)
            return

        try:
            from pvlib import solarposition as _sp
            tz, times, _sol_full = _get_solar_positions(_lat, _lon, _tz_str)
            sol = _sol_full[_sol_full["apparent_elevation"] > 0].copy()

            _df = df.copy()
            if "datetime" not in _df.columns:
                if isinstance(_df.index, pd.DatetimeIndex):
                    _df = _df.reset_index().rename(columns={_df.index.name or "index": "datetime"})
                else:
                    raise ValueError("EPW missing 'datetime' column")

            common_aliases = {
                'dry bulb': 'dry_bulb_temperature',
                'dry_bulb': 'dry_bulb_temperature',
                'drybulb': 'dry_bulb_temperature',
                'temperature': 'dry_bulb_temperature',
                'ghi': 'global_horizontal_irradiance',
                'global horizontal': 'global_horizontal_irradiance',
            }

            cols_lower = {c: c.lower() for c in _df.columns}
            rename_map = {}
            for alias, target in common_aliases.items():
                if target in _df.columns:
                    continue
                for orig, low in cols_lower.items():
                    if alias == low or alias in low:
                        if orig != 'datetime' and target not in _df.columns:
                            rename_map[orig] = target
                            break

            if rename_map:
                _df = _df.rename(columns=rename_map)

            epw = _df.set_index("datetime").copy()
            if epw.index.tz is None:
                epw.index = epw.index.tz_localize(tz)
            else:
                epw.index = epw.index.tz_convert(tz)
            epw.index = epw.index.map(lambda x: x.replace(year=2020))

            try:
                has_half_hour = any(t.minute != 0 for t in epw.index)
            except Exception:
                has_half_hour = False

            epw_hourly = epw
            if has_half_hour:
                candidate_cols = [c for c in [
                    "dry_bulb_temperature", "global_horizontal_irradiance"
                ] if c in epw.columns]
                if candidate_cols:
                    epw_num = epw[candidate_cols].apply(pd.to_numeric, errors='coerce')
                    epw_hourly = epw_num.resample('h').mean().dropna(how='all')
                else:
                    epw_hourly = epw.resample('h').first().dropna(how='all')

            sol = sol.join(epw_hourly[["dry_bulb_temperature", "global_horizontal_irradiance"]], how="left")
            sol["global_horizontal_irradiance"] = sol["global_horizontal_irradiance"].fillna(0)
            sol["dry_bulb_temperature"] = sol["dry_bulb_temperature"].fillna(
                sol["dry_bulb_temperature"].median()
            )
            shading_needed = (
                (sol["dry_bulb_temperature"] > temp_threshold) &
                (sol["global_horizontal_irradiance"] > rad_threshold)
            )

            fig = plt.figure(figsize=(7.5, 7.5), dpi=100, facecolor="white")
            ax = fig.add_subplot(111, projection="polar")
            ax.set_theta_zero_location("N")
            ax.set_theta_direction(-1)
            ax.set_aspect('equal', adjustable='box')
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(["90°","75°","60°","45°","30°","15°","0°"],
                               fontsize=7, color="#555")
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(["N","NE","E","SE","S","SW","W","NW"], fontsize=10, fontweight="bold")
            ax.set_facecolor("#F0F4F8")
            ax.grid(True, alpha=0.35, linestyle="--", linewidth=0.6)

            r = 90 - sol["apparent_elevation"].values
            theta = np.radians(sol["azimuth"].values)

            mask_ok = ~shading_needed.values
            ax.scatter(theta[mask_ok], r[mask_ok], c="#FFF9C4", s=1.2,
                       alpha=0.45, linewidths=0, label="No shading needed", zorder=2)

            mask_sh = shading_needed.values
            ax.scatter(theta[mask_sh], r[mask_sh], c="#E65100", s=2.5,
                       alpha=0.75, linewidths=0, label="Shading required", zorder=3)

            for lbl, dstr, col, lw in [
                ("Mar 21", "2020-03-21", "#FF9500", 1.4),
                ("Jun 21", "2020-06-21", "#CC0000", 1.8),
                ("Dec 21", "2020-12-21", "#0066CC", 1.8),
            ]:
                dt = pd.date_range(dstr, periods=288, freq="5min", tz=tz)
                ks = _sp.get_solarposition(dt, _lat, _lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            shading_pct = mask_sh.sum() / len(mask_sh) * 100 if len(mask_sh) else 0
            ax.set_title(
                f"Sun Path – Shading Mode   ({shading_pct:.1f}% of daytime hours require shading)",
                fontsize=10, fontweight="bold", color="#333", pad=14,
            )
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.06), ncol=3,
                      frameon=True, fontsize=8)

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)

            # Use square dimensions to maintain circular aspect ratio
            img_size = min(SW * 0.75, SH * 0.75)
            img_l = (SW - img_size) / 2
            slide.shapes.add_picture(tmp, Inches(img_l), Inches(0.72),
                                     width=Inches(img_size), height=Inches(img_size))
            os.unlink(tmp)

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_sun_path_shading_slide()

    # ── SECTION 6 – ORIENTATION SHADING ANALYSIS TABLE ────────────────────────
    def _make_orientation_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, f"Orientation Shading Analysis  (Design cutoff: {design_cutoff_angle}°)")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        _lat = _meta.get("latitude")
        _lon = _meta.get("longitude")
        _tz_str = _meta.get("timezone", "UTC")

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err_box(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz_str)
            if solar_pos.empty:
                _err_box(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_df = build_orientation_table(solar_pos, design_cutoff_angle)

            fig, ax = plt.subplots(figsize=(13, 5.8), dpi=120)
            ax.axis("off")

            col_labels = ["Orientation", "Rays\nHitting", "Min VSA", "Max |HSA|",
                          "D/H\nOverhang", "D/W\nFin", "Protection %"]
            table_data = []
            row_colors = []
            for _, row in orient_df.iterrows():
                pct = row["Protection (%)"]
                if pct is None:
                    c = "#f5f5f5"
                elif pct >= 95:
                    c = "#e8f5e9"
                elif pct >= 80:
                    c = "#fff3e0"
                else:
                    c = "#ffebee"
                row_colors.append([c] * 7)

                dh  = f"{row['D/H (Overhang)']:.3f}" if row["D/H (Overhang)"] is not None else "—"
                dw  = f"{row['D/W (Fin)']:.3f}"       if row["D/W (Fin)"] is not None else "—"
                vsa = f"{row['Min VSA (°)']:.1f}°"  if row["Min VSA (°)"] is not None else "—"
                hsa = f"{row['Max |HSA| (°)']:.1f}°" if row["Max |HSA| (°)"] is not None else "—"
                pct_s = f"{pct:.1f}%" if pct is not None else "—"
                table_data.append([
                    row["Orientation"], str(row["Rays Hitting"]),
                    vsa, hsa, dh, dw, pct_s,
                ])

            tbl = ax.table(cellText=table_data, colLabels=col_labels,
                           cellColours=row_colors, loc="center", cellLoc="center")
            tbl.auto_set_font_size(False)
            tbl.set_fontsize(11)
            tbl.scale(1, 2.1)

            for j in range(len(col_labels)):
                cell = tbl[0, j]
                cell.set_facecolor("#1a3a52")
                cell.set_text_props(color="white", fontweight="bold")

            for i in range(1, len(table_data) + 1):
                tbl[i, 0].get_text().set_ha("left")

            fig.patch.set_facecolor("white")
            ax.set_title(
                f"{len(solar_pos)} overheating daytime sun positions  |  "
                f"Temp > {temp_threshold}°C  &  GHI > {rad_threshold} W/m²",
                fontsize=10, color="#555", pad=10,
            )

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_orientation_slide()

    # ── SECTION 7 – SHADING MASK DIAGRAMS (2×4 grid) ─────────────────────────
    def _make_shading_masks_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Shading Mask Diagrams")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        _lat = _meta.get("latitude")
        _lon = _meta.get("longitude")
        _tz_str = _meta.get("timezone", "UTC")

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err_box(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz_str)
            if solar_pos.empty:
                _err_box(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_items = list(_ORIENTATIONS.items())
            n_cols = 4
            cell_w = (SW - 0.54) / n_cols
            cell_h = (SH - 1.05) / 2
            top_start = 0.75

            for idx, (oname, faz) in enumerate(orient_items):
                col_i = idx % n_cols
                row_i = idx // n_cols
                cell_l = 0.27 + col_i * cell_w
                cell_t = top_start + row_i * cell_h

                geom   = compute_shading_geometry(solar_pos, faz)
                facing = geom[geom["hits_facade"]]
                other  = geom[~geom["hits_facade"]]

                fig = plt.figure(figsize=(3.0, 2.8), dpi=110, facecolor="white")
                ax = fig.add_subplot(111, projection="polar")
                ax.set_theta_zero_location("N")
                ax.set_theta_direction(-1)
                ax.set_ylim(0, 90)
                ax.set_yticks([])
                ax.set_xticks(np.radians([0, 90, 180, 270]))
                ax.set_xticklabels(["N", "E", "S", "W"], fontsize=8, fontweight="bold")
                ax.set_facecolor("#f0f8ff")
                ax.grid(True, alpha=0.30, linewidth=0.5)

                if not other.empty:
                    ax.scatter(
                        np.radians(other["solar_azimuth"]), 90 - other["solar_altitude"],
                        s=2, c="lightgrey", alpha=0.5, linewidths=0, zorder=2,
                    )
                if not facing.empty:
                    ax.scatter(
                        np.radians(facing["solar_azimuth"]), 90 - facing["solar_altitude"],
                        s=4, c="#E65100", alpha=0.75, linewidths=0, zorder=3,
                    )

                rel_az_r = np.linspace(-89, 89, 179)
                tan_co = np.tan(np.radians(design_cutoff_angle))
                co_alt = np.degrees(np.arctan(tan_co * np.cos(np.radians(rel_az_r))))
                co_az  = faz + rel_az_r
                valid  = co_alt > 0
                if valid.any():
                    ax.plot(np.radians(co_az[valid]), 90 - co_alt[valid],
                            color="#1565C0", linewidth=1.4, linestyle="--", zorder=4)

                ax.plot(np.radians([faz, faz]), [0, 85],
                        color="#388E3C", linewidth=1.4, zorder=5)

                ax.set_title(oname, fontsize=7, fontweight="bold", pad=4, color="#222")
                plt.tight_layout(pad=0.3)

                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(
                    tmp, Inches(cell_l), Inches(cell_t),
                    width=Inches(cell_w - 0.05), height=Inches(cell_h - 0.05),
                )
                os.unlink(tmp)

            leg_tb = slide.shapes.add_textbox(
                Inches(0.27), Inches(SH - LOGO_H - 0.55), Inches(SW - 0.54), Inches(0.35),
            )
            leg_tf = leg_tb.text_frame
            leg_p = leg_tf.paragraphs[0]
            leg_run = leg_p.add_run()
            leg_run.text = (
                "● Overheating rays (hits facade)  "
                "● Overheating (other side)  "
                "- - Cutoff arc (VSA cut-off)  "
                "— Facade direction"
            )
            leg_run.font.size = Pt(8)
            leg_run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_shading_masks_slide()

    # ── SECTION 8 – WIND ANALYSIS SLIDES ──────────────────────────────────────
    def _prepare_wind_slides():
        """Prepare and add wind analysis slides."""
        try:
            from .wind_module import (
                prepare_wind_data, compute_wind_rose, compute_wind_statistics,
                _SPEED_LABELS, _SPEED_COLORS, _SPEED_BINS,
                _DIR_16, _DIR_8, _DIR_4, _MONTH_NAMES, _MONTH_COLORS,
            )
        except ImportError:
            return  # Skip if wind module not available

        # Prepare wind data
        months = list(range(1, 13))  # All months
        wdf = prepare_wind_data(df, months=months, n_sectors=n_sectors)

        if wdf.empty:
            return  # No wind data available

        rose_df, calm_pct = compute_wind_rose(wdf, n_sectors, exclude_calm=False)
        stats = compute_wind_statistics(wdf)

        # ── Matplotlib chart helpers (no kaleido / Chrome required) ─────────
        def _mpl_wind_rose_png():
            sector_width = 360.0 / n_sectors
            if n_sectors == 16:
                lbl = _DIR_16
            elif n_sectors == 8:
                lbl = _DIR_8
            elif n_sectors == 4:
                lbl = _DIR_4
            else:
                lbl = [f"{int(i * sector_width)}°" for i in range(n_sectors)]
            angles = np.array([np.deg2rad(i * sector_width) for i in range(n_sectors)])
            bar_w = np.deg2rad(sector_width) * 0.85
            fig2, ax = plt.subplots(subplot_kw=dict(polar=True), figsize=(9, 7))
            bottoms = np.zeros(n_sectors)
            for i, sl in enumerate(_SPEED_LABELS):
                subset = rose_df[rose_df["speed_bin"] == sl]
                fm = dict(zip(subset["direction_label"], subset["frequency_pct"]))
                freqs = np.array([fm.get(l, 0.0) for l in lbl])
                ax.bar(angles, freqs, width=bar_w, bottom=bottoms,
                       color=_SPEED_COLORS[i % len(_SPEED_COLORS)],
                       label=f"{sl} m/s", alpha=0.9, linewidth=0.3, edgecolor="white")
                bottoms += freqs
            ax.set_theta_zero_location("N")
            ax.set_theta_direction(-1)
            ax.set_xticks(angles)
            ax.set_xticklabels(lbl, fontsize=9)
            ax.tick_params(axis="y", labelsize=8)
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.1f}%"))
            ax.set_title("Wind Rose", fontsize=14, pad=20, color="#2c3e50", fontweight="bold")
            ax.annotate(f"Calm\n{calm_pct:.1f}%", xy=(0, 0), xycoords="data",
                        ha="center", va="center", fontsize=10, color="#555555")
            ax.legend(loc="lower left", bbox_to_anchor=(1.05, 0.0),
                      fontsize=9, title="Wind Speed (m/s)", title_fontsize=9)
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        def _mpl_speed_heatmap_png():
            pivot = wdf.pivot_table(values="wind_speed", index="hour",
                                    columns="month", aggfunc="mean")
            for m in range(1, 13):
                if m not in pivot.columns:
                    pivot[m] = np.nan
            pivot = pivot[sorted(pivot.columns)]
            month_labels = [_MONTH_NAMES[m - 1] for m in sorted(pivot.columns)]
            fig2, ax = plt.subplots(figsize=(12, 5))
            im = ax.imshow(pivot.values, aspect="auto", cmap="viridis", interpolation="nearest")
            ax.set_xlabel("Month", fontsize=11)
            ax.set_ylabel("Hour of Day", fontsize=11)
            ax.set_title("Wind Speed – Month × Hour", fontsize=14, color="#2c3e50", fontweight="bold")
            ax.set_xticks(range(12))
            ax.set_xticklabels(month_labels, fontsize=9)
            ax.set_yticks(range(0, 24, 3))
            ax.set_yticklabels([f"{h:02d}:00" for h in range(0, 24, 3)], fontsize=9)
            plt.colorbar(im, ax=ax, label="m/s", shrink=0.8)
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        def _mpl_direction_heatmap_png():
            tmp_df = wdf.copy()
            rad = np.deg2rad(tmp_df["wind_direction"])
            tmp_df["_u"] = np.cos(rad)
            tmp_df["_v"] = np.sin(rad)
            up = tmp_df.pivot_table(values="_u", index="hour", columns="month", aggfunc="mean")
            vp = tmp_df.pivot_table(values="_v", index="hour", columns="month", aggfunc="mean")
            up, vp = up.align(vp, join="inner")
            dir_deg = np.degrees(np.arctan2(vp.values, up.values)) % 360
            month_cols = sorted(up.columns.tolist())
            month_labels = [_MONTH_NAMES[m - 1] for m in month_cols]
            fig2, ax = plt.subplots(figsize=(12, 5))
            im = ax.imshow(dir_deg, aspect="auto", cmap="twilight",
                           vmin=0, vmax=360, interpolation="nearest")
            ax.set_xlabel("Month", fontsize=11)
            ax.set_ylabel("Hour of Day", fontsize=11)
            ax.set_title("Wind Direction – Month × Hour", fontsize=14, color="#2c3e50", fontweight="bold")
            ax.set_xticks(range(len(month_cols)))
            ax.set_xticklabels(month_labels, fontsize=9)
            ax.set_yticks(range(0, 24, 3))
            ax.set_yticklabels([f"{h:02d}:00" for h in range(0, 24, 3)], fontsize=9)
            cbar = plt.colorbar(im, ax=ax, shrink=0.8)
            cbar.set_ticks([0, 90, 180, 270, 360])
            cbar.set_ticklabels(["N 0°", "E 90°", "S 180°", "W 270°", "N 360°"])
            cbar.set_label("Direction")
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        def _mpl_speed_histogram_png():
            total = len(wdf)
            labels, pcts = [], []
            for i in range(len(_SPEED_BINS) - 1):
                lo, hi = _SPEED_BINS[i], _SPEED_BINS[i + 1]
                count = int(((wdf["wind_speed"] >= lo) & (wdf["wind_speed"] < hi)).sum())
                labels.append(_SPEED_LABELS[i])
                pcts.append(count / total * 100.0 if total > 0 else 0.0)
            fig2, ax = plt.subplots(figsize=(10, 5))
            bars = ax.bar(labels, pcts, color=_SPEED_COLORS[:len(labels)], alpha=0.9, edgecolor="white")
            for bar, pct in zip(bars, pcts):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                        f"{pct:.1f}%", ha="center", va="bottom", fontsize=10)
            ax.set_xlabel("Wind Speed Bin (m/s)", fontsize=11)
            ax.set_ylabel("Frequency (%)", fontsize=11)
            ax.set_title("Wind Speed Distribution", fontsize=14, color="#2c3e50", fontweight="bold")
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        def _mpl_climate_bubble_png():
            needed = {"dry_bulb_temperature", "relative_humidity", "wind_speed", "month"}
            fig2, ax = plt.subplots(figsize=(10, 6))
            if not needed.issubset(wdf.columns):
                ax.text(0.5, 0.5, "Missing columns for bubble chart",
                        ha="center", va="center", transform=ax.transAxes, fontsize=13)
            else:
                tmp_df = wdf.dropna(subset=["dry_bulb_temperature", "relative_humidity", "wind_speed"]).copy()
                max_spd = float(tmp_df["wind_speed"].max())
                scale = 500.0 / max_spd if max_spd > 0 else 50.0
                for m in range(1, 13):
                    mdata = tmp_df[tmp_df["month"] == m]
                    if mdata.empty:
                        continue
                    ax.scatter(mdata["dry_bulb_temperature"], mdata["relative_humidity"],
                               s=(mdata["wind_speed"] + 0.3) * scale,
                               c=_MONTH_COLORS[(m - 1) % len(_MONTH_COLORS)],
                               alpha=0.45, linewidths=0, label=_MONTH_NAMES[m - 1])
                ax.set_xlabel("Dry Bulb Temperature (°C)", fontsize=11)
                ax.set_ylabel("Relative Humidity (%)", fontsize=11)
                ax.set_ylim(0, 105)
                ax.legend(title="Month", fontsize=9, title_fontsize=9,
                          loc="center left", bbox_to_anchor=(1, 0.5))
                ax.text(0.01, 0.98, "Bubble size = wind speed (m/s)",
                        transform=ax.transAxes, fontsize=10, color="#888", va="top")
            ax.set_title("Temperature – Humidity – Wind Speed", fontsize=14,
                         color="#2c3e50", fontweight="bold")
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        # ── Wind Rose Slide ─────────────────────────────────────────────────
        def _wind_rose_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Wind Rose Analysis")
            _add_divider(slide, 0.62)

            try:
                img_path = _mpl_wind_rose_png()
                chart_w = 7.5
                slide.shapes.add_picture(img_path, Inches((SW - chart_w) / 2), Inches(0.80),
                                         width=Inches(chart_w))
                os.unlink(img_path)
            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        def _mpl_single_season_rose_png(rose_df_s, calm_pct_s, season_name):
            sector_width = 360.0 / n_sectors
            if n_sectors == 16:
                lbl = _DIR_16
            elif n_sectors == 8:
                lbl = _DIR_8
            elif n_sectors == 4:
                lbl = _DIR_4
            else:
                lbl = [f"{int(i * sector_width)}°" for i in range(n_sectors)]
            angles = np.array([np.deg2rad(i * sector_width) for i in range(n_sectors)])
            bar_w = np.deg2rad(sector_width) * 0.85
            fig2, ax = plt.subplots(subplot_kw=dict(polar=True), figsize=(9, 7))
            bottoms = np.zeros(n_sectors)
            for i, sl in enumerate(_SPEED_LABELS):
                subset = rose_df_s[rose_df_s["speed_bin"] == sl]
                fm = dict(zip(subset["direction_label"], subset["frequency_pct"]))
                freqs = np.array([fm.get(l, 0.0) for l in lbl])
                ax.bar(angles, freqs, width=bar_w, bottom=bottoms,
                       color=_SPEED_COLORS[i % len(_SPEED_COLORS)],
                       label=f"{sl} m/s", alpha=0.9, linewidth=0.3, edgecolor="white")
                bottoms += freqs
            ax.set_theta_zero_location("N")
            ax.set_theta_direction(-1)
            ax.set_xticks(angles)
            ax.set_xticklabels(lbl, fontsize=9)
            ax.tick_params(axis="y", labelsize=8)
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.1f}%"))
            ax.set_title(f"{season_name} Wind Rose", fontsize=14, pad=20,
                         color="#2c3e50", fontweight="bold")
            ax.annotate(f"Calm\n{calm_pct_s:.1f}%", xy=(0, 0), xycoords="data",
                        ha="center", va="center", fontsize=10, color="#555555")
            ax.legend(loc="lower left", bbox_to_anchor=(1.05, 0.0),
                      fontsize=9, title="Wind Speed (m/s)", title_fontsize=9)
            plt.tight_layout()
            path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            fig2.savefig(path, dpi=100, bbox_inches="tight", facecolor="white")
            plt.close(fig2)
            return path

        _wind_rose_slide()

        # ── Seasonal Wind Rose Slides (one slide per season) ────────────────
        for _sname, _smonths in [("Winter", [12, 1, 2]), ("Spring", [3, 4, 5]),
                                  ("Summer", [6, 7, 8]),  ("Fall",   [9, 10, 11])]:
            _sdf = wdf[wdf["month"].isin(_smonths)].copy()
            if _sdf.empty:
                continue
            _srose, _scalm = compute_wind_rose(_sdf, n_sectors=n_sectors, exclude_calm=False)
            _sslide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(_sslide, f"Wind Rose – {_sname}")
            _add_divider(_sslide, 0.62)
            try:
                _spath = _mpl_single_season_rose_png(_srose, _scalm, _sname)
                _cw = 7.5
                _sslide.shapes.add_picture(_spath, Inches((SW - _cw) / 2), Inches(0.80),
                                           width=Inches(_cw))
                os.unlink(_spath)
            except Exception as _se:
                _err_box(_sslide, _se)
            _add_logo(_sslide)

        # ── Wind Speed Heatmap Slide ────────────────────────────────────────
        def _speed_heatmap_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Wind Speed Heatmap (Month × Hour)")
            _add_divider(slide, 0.62)

            try:
                img_path = _mpl_speed_heatmap_png()
                slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(img_path)
            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        _speed_heatmap_slide()

        # ── Wind Direction Heatmap Slide ────────────────────────────────────
        def _direction_heatmap_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Wind Direction Heatmap (Month × Hour)")
            _add_divider(slide, 0.62)

            try:
                img_path = _mpl_direction_heatmap_png()
                slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(img_path)
            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        _direction_heatmap_slide()

        # ── Wind Speed Distribution Slide ───────────────────────────────────
        def _speed_histogram_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Wind Speed Distribution")
            _add_divider(slide, 0.62)

            try:
                img_path = _mpl_speed_histogram_png()
                slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(img_path)
            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        _speed_histogram_slide()

        # ── Climate Bubble Chart Slide ──────────────────────────────────────
        def _climate_bubble_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Temperature – Humidity – Wind Speed")
            _add_divider(slide, 0.62)

            try:
                img_path = _mpl_climate_bubble_png()
                slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(img_path)
            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        _climate_bubble_slide()

        # ── Wind Statistics Summary Slide ───────────────────────────────────
        def _wind_statistics_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Wind Statistics Summary")
            _add_divider(slide, 0.62)

            try:
                # Prepare statistics data with colors and icons
                stats_data = [
                    {
                        "label": "Prevailing Direction",
                        "value": stats.get("prevailing_direction", "N/A"),
                        "color": RGBColor(0x3B, 0x82, 0xF6),  # Blue
                        "bg_color": RGBColor(0xEF, 0xF6, 0xFF),  # Light blue
                    },
                    {
                        "label": "Mean Wind Speed",
                        "value": f"{stats.get('mean_speed', 0):.2f} m/s",
                        "color": RGBColor(0x8B, 0x5C, 0xF6),  # Purple
                        "bg_color": RGBColor(0xF5, 0xF3, 0xFF),  # Light purple
                    },
                    {
                        "label": "Maximum Wind Speed",
                        "value": f"{stats.get('max_speed', 0):.2f} m/s",
                        "color": RGBColor(0xEF, 0x44, 0x44),  # Red
                        "bg_color": RGBColor(0xFF, 0xF1, 0xF1),  # Light red
                    },
                    {
                        "label": "Calm Hours",
                        "value": f"{stats.get('calm_percent', 0):.1f}%",
                        "color": RGBColor(0xF5, 0x9E, 0x0B),  # Amber
                        "bg_color": RGBColor(0xFF, 0xF8, 0xE7),  # Light amber
                    },
                    {
                        "label": "Strongest Direction",
                        "value": stats.get("strongest_direction", "N/A"),
                        "color": RGBColor(0x06, 0xB6, 0xD4),  # Cyan
                        "bg_color": RGBColor(0xEC, 0xF8, 0xFE),  # Light cyan
                    },
                    {
                        "label": "Total Data Points",
                        "value": f"{len(wdf)} hours",
                        "color": RGBColor(0x10, 0xB9, 0x81),  # Green
                        "bg_color": RGBColor(0xF0, 0xFF, 0xF4),  # Light green
                    },
                ]

                # Create 3x2 grid of cards
                card_width = (SW - 0.8) / 3
                card_height = 1.8
                start_top = 0.75
                start_left = 0.27

                for idx, stat in enumerate(stats_data):
                    col = idx % 3
                    row = idx // 3
                    
                    left = start_left + col * (card_width + 0.05)
                    top = start_top + row * (card_height + 0.15)

                    # Add card background shape with border
                    card = slide.shapes.add_shape(
                        1,  # Rectangle
                        Inches(left),
                        Inches(top),
                        Inches(card_width),
                        Inches(card_height),
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = stat["bg_color"]
                    card.line.color.rgb = stat["color"]
                    card.line.width = Pt(2)

                    # Add label
                    label_tb = slide.shapes.add_textbox(
                        Inches(left + 0.1),
                        Inches(top + 0.1),
                        Inches(card_width - 0.2),
                        Inches(0.5),
                    )
                    label_tf = label_tb.text_frame
                    label_tf.word_wrap = True
                    p = label_tf.paragraphs[0]
                    run = p.add_run()
                    run.text = stat["label"]
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = stat["color"]

                    # Add value
                    value_tb = slide.shapes.add_textbox(
                        Inches(left + 0.1),
                        Inches(top + 0.65),
                        Inches(card_width - 0.2),
                        Inches(0.9),
                    )
                    value_tf = value_tb.text_frame
                    value_tf.word_wrap = True
                    value_tf.vertical_anchor = 1  # Middle alignment
                    p = value_tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = stat["value"]
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = DARK_GREY

            except Exception as e:
                _err_box(slide, e)

            _add_logo(slide)

        _wind_statistics_slide()

    _prepare_wind_slides()

    # ── SECTION 9 – THERMAL COMFORT ANALYSIS SLIDES ──────────────────────────
    def _prepare_thermal_comfort_slides():
        """Prepare and add thermal comfort analysis slides."""
        if not include_thermal_comfort:
            return  # Skip thermal comfort section if not requested
        try:
            from .thermal_comfort_ppt import (
                compute_psychrometric_simple,
                compute_adaptive_comfort_simple,
                classify_comfort_simple,
                plot_comfort_heatmap,
                plot_strategy_distribution,
                plot_degree_hours_monthly,
                plot_adaptive_comfort_scatter,
                plot_psychrometric_chart,
                plot_comfort_percentages
            )
        except ImportError:
            return  # Skip if thermal comfort module not available

        # Process thermal comfort data
        try:
            tdf = df.copy()
            tdf = compute_psychrometric_simple(tdf)
            tdf = compute_adaptive_comfort_simple(tdf)
            tdf = classify_comfort_simple(tdf)
            # Add strategy column based on comfort category
            tdf["strategy"] = tdf["comfort_cat"]
        except Exception as e:
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Thermal Comfort Analysis")
            _add_divider(slide, 0.62)
            _err_box(slide, f"Thermal comfort data processing error: {str(e)[:50]}")
            _add_logo(slide)
            return

        # ── Comfort Heatmap Slide ──────────────────────────────────────────
        def _comfort_heatmap_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Comfort Heatmap – Hour × Month")
            _add_divider(slide, 0.62)

            try:
                fig = plot_comfort_heatmap(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Comfort heatmap: {str(e)[:40]}")

            _add_logo(slide)

        _comfort_heatmap_slide()

        # ── Psychrometric Chart Slide ──────────────────────────────────────
        def _psychrometric_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Psychrometric Chart – Climate Data")
            _add_divider(slide, 0.62)

            try:
                fig = plot_psychrometric_chart(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Psychrometric chart: {str(e)[:40]}")

            _add_logo(slide)

        _psychrometric_slide()

        # ── Strategy Distribution Slide ────────────────────────────────────
        def _strategy_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Design Strategy Opportunities")
            _add_divider(slide, 0.62)

            try:
                fig = plot_strategy_distribution(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Strategy chart: {str(e)[:40]}")

            _add_logo(slide)

        _strategy_slide()

        # ── Degree Hours Slide ─────────────────────────────────────────────
        def _degree_hours_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Monthly Degree Hours – Cooling & Heating Demand")
            _add_divider(slide, 0.62)

            try:
                fig = plot_degree_hours_monthly(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Degree hours: {str(e)[:40]}")

            _add_logo(slide)

        _degree_hours_slide()

        # ── Adaptive Comfort Slide ─────────────────────────────────────────
        def _adaptive_comfort_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Adaptive Comfort – ASHRAE 55 Analysis")
            _add_divider(slide, 0.62)

            try:
                fig = plot_adaptive_comfort_scatter(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Adaptive comfort: {str(e)[:40]}")

            _add_logo(slide)

        _adaptive_comfort_slide()

        # ── Comfort Performance Summary Slide ───────────────────────────────
        def _performance_summary_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Thermal Comfort Performance Summary")
            _add_divider(slide, 0.62)

            try:
                fig = plot_comfort_percentages(tdf)
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                         width=Inches(SW - 0.54), height=Inches(5.9))
                os.unlink(tmp)
            except Exception as e:
                _err_box(slide, f"Performance summary: {str(e)[:40]}")

            _add_logo(slide)

        _performance_summary_slide()

        # ── Design Recommendations Slide ────────────────────────────────────
        def _design_recommendations_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, "Design Recommendations & Strategies")
            _add_divider(slide, 0.62)
            
            tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.75), Inches(SW - 0.54), Inches(6.0))
            tf = tb.text_frame
            tf.word_wrap = True
            
            # Calculate summary statistics
            try:
                pct_comfortable = (tdf["comfort_cat"] == "Comfortable").sum() / len(tdf) * 100
                pct_hot = (tdf["comfort_cat"] == "Too Hot").sum() / len(tdf) * 100
                pct_cold = (tdf["comfort_cat"] == "Too Cold").sum() / len(tdf) * 100
                mean_rh = pd.to_numeric(tdf["relative_humidity"], errors="coerce").mean()
                
                p = tf.paragraphs[0]
                p.text = "Climate Comfort Profile"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = TITLE_RED
                p.space_after = Pt(6)
                
                p = tf.add_paragraph()
                p.text = f"• Comfortable hours: {pct_comfortable:.1f}%  |  Overheating: {pct_hot:.1f}%  |  Undercooling: {pct_cold:.1f}%"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(4)
                
                p = tf.add_paragraph()
                p.text = f"• Mean Relative Humidity: {mean_rh:.1f}%"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(10)
                
                p = tf.add_paragraph()
                p.text = "Key Design Strategies"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = TITLE_RED
                p.space_before = Pt(4)
                p.space_after = Pt(6)
                
                if pct_comfortable < 40:
                    if pct_hot > pct_cold:
                        p = tf.add_paragraph()
                        p.text = "• Priority: Cooling strategies – Implement high-performance envelope, external shading, and natural ventilation"
                        p.font.size = Pt(11)
                        p.font.color.rgb = DARK_GREY
                        p.space_after = Pt(3)
                        
                        p = tf.add_paragraph()
                        p.text = "• Consider nighttime cooling recovery and thermal mass activation"
                        p.font.size = Pt(11)
                        p.font.color.rgb = DARK_GREY
                        p.space_after = Pt(3)
                    else:
                        p = tf.add_paragraph()
                        p.text = "• Priority: Heating strategies – Maximize solar heat gain during winter with south-facing glazing"
                        p.font.size = Pt(11)
                        p.font.color.rgb = DARK_GREY
                        p.space_after = Pt(3)
                        
                        p = tf.add_paragraph()
                        p.text = "• Ensure robust thermal insulation and minimize infiltration losses"
                        p.font.size = Pt(11)
                        p.font.color.rgb = DARK_GREY
                        p.space_after = Pt(3)
                else:
                    p = tf.add_paragraph()
                    p.text = "• Climate is generally favorable – Prioritize passive design with natural ventilation and daylighting"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                
                if mean_rh > 65:
                    p = tf.add_paragraph()
                    p.text = "• High humidity detected – Ensure adequate dehumidification and mold risk mitigation"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                elif mean_rh < 30:
                    p = tf.add_paragraph()
                    p.text = "• Low humidity detected – Humidification may be required in heating season for occupant comfort"
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_GREY
                    p.space_after = Pt(3)
                
                p = tf.add_paragraph()
                p.text = "• Adaptive Comfort – Leverage occupant behavior (clothing, behavior) to expand acceptable temperature ranges"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GREY
                p.space_after = Pt(0)
                
            except Exception as e:
                p = tf.paragraphs[0]
                p.text = f"Error generating recommendations: {str(e)[:40]}"
                p.font.size = Pt(11)
                p.font.color.rgb = TITLE_RED
            
            _add_logo(slide)
        
        _design_recommendations_slide()

    _prepare_thermal_comfort_slides()

    # ── SECTION 10 – RAINFALL ANALYSIS SLIDES (optional) ──────────────────────
    def _prepare_rainfall_slides():
        """Add rainfall analysis section if station/year are provided."""
        if not rainfall_station_name or not rainfall_station_id or not rainfall_year:
            return

        try:
            from .rainfall_module import _fetch_noaa, _fetch_percentile_depth, _RUNOFF_SURFACES
        except ImportError:
            return

        import matplotlib.patches as mpatches

        # Use pre-fetched data from background thread if available; fall back to direct call
        try:
            if _noaa_future is not None:
                try:
                    df_rain = _noaa_future.result(timeout=90)
                except Exception:
                    df_rain = _fetch_noaa(rainfall_station_id, rainfall_year)
            else:
                df_rain = _fetch_noaa(rainfall_station_id, rainfall_year)
            if df_rain is None or df_rain.empty:
                return
        except Exception:
            return

        _sa  = rainfall_surface_areas or {"roof": 0.0, "paved": 0.0, "green": 0.0, "water": 0.0}
        df_f = df_rain[
            (df_rain["month"] >= rainfall_start_month) &
            (df_rain["month"] <= rainfall_end_month)
        ].copy()

        ML = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
              "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
        xm = np.arange(12)

        CHART_TOP = 0.72
        CHART_H   = 4.48
        KPI_TOP   = 5.33
        KPI_H     = 0.85

        def _hex_rgb(h):
            return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

        def _fmt_L(litres, unit="L"):
            if litres >= 1_000_000:
                return f"{litres / 1_000_000:.2f}M {unit}"
            if litres >= 10_000:
                return f"{litres / 1_000:.1f}K {unit}"
            return f"{litres:,.0f} {unit}"

        def _kpi_card(slide, left, top, width, height, label, value, sub, hex_color):
            rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex_rgb(hex_color)
            rect.line.fill.background()
            label_clr = RGBColor(0xCC, 0xDD, 0xFF)
            sub_clr   = RGBColor(0xD4, 0xE8, 0xFF)
            for (tx, ty, tw, th, txt, fsize, bold, clr, wrap) in [
                (left + 0.07, top + 0.06, width - 0.14, 0.20, label, Pt(9), False, label_clr, True),
            ]:
                tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
                tf = tb.text_frame
                tf.word_wrap = wrap
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = txt
                r.font.size = fsize
                r.font.bold = bold
                r.font.color.rgb = clr
            val_top = top + (height - 0.28) / 2 - 0.05
            tb_v = slide.shapes.add_textbox(Inches(left + 0.05), Inches(val_top), Inches(width - 0.10), Inches(0.50))
            tf_v = tb_v.text_frame
            tf_v.word_wrap = False
            p_v = tf_v.paragraphs[0]
            p_v.alignment = PP_ALIGN.CENTER
            r_v = p_v.add_run()
            r_v.text = value
            r_v.font.size = Pt(26) if height >= 0.80 else Pt(22)
            r_v.font.bold = True
            r_v.font.color.rgb = WHITE
            tb_s = slide.shapes.add_textbox(Inches(left + 0.07), Inches(top + height - 0.22), Inches(width - 0.14), Inches(0.20))
            p_s = tb_s.text_frame.paragraphs[0]
            p_s.alignment = PP_ALIGN.CENTER
            r_s = p_s.add_run()
            r_s.text = sub
            r_s.font.size = Pt(8)
            r_s.font.color.rgb = sub_clr

        def _kpi_row(slide, cards, top, card_h=0.85):
            n = len(cards)
            gap = 0.12
            card_w = (SW - 0.54 - gap * (n - 1)) / n
            left = 0.27
            for lbl, val, sub, clr in cards:
                _kpi_card(slide, left, top, card_w, card_h, lbl, val, sub, clr)
                left += card_w + gap

        def _intensity_color(mm):
            if mm < 50:    return "#93c5fd"
            elif mm < 150: return "#3b82f6"
            elif mm < 300: return "#1d4ed8"
            else:          return "#1e3a5f"

        # Section header
        def _section_header():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = TITLE_RED
            bg.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(SW - 1.2), Inches(1.1))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = "Rainfall Analysis"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = WHITE
            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.6))
            run2 = tb2.text_frame.paragraphs[0].add_run()
            run2.text = f"{rainfall_station_name}  |  {rainfall_year}"
            run2.font.size = Pt(18)
            run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
            _add_logo(slide)

        _section_header()

        # Monthly Rainfall
        def _monthly_rainfall():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, f"Monthly Rainfall  —  {rainfall_station_name}, {rainfall_year}")
            _add_divider(slide, 0.62)
            try:
                monthly      = df_f.groupby("month")["prcp_mm"].sum().reindex(range(1, 13), fill_value=0)
                annual_total = float(monthly.sum())
                annual_mean  = annual_total / 12
                wettest_idx  = int(monthly.idxmax())
                fig, ax = plt.subplots(figsize=(13, 4.5), dpi=100)
                ax.bar(xm, monthly.values, color=[_intensity_color(v) for v in monthly.values], edgecolor="none")
                ax.axhline(annual_mean, color="#ef4444", linewidth=1.4, linestyle="--")
                ax.text(10.6, annual_mean * 1.02, f"Mean: {annual_mean:.0f} mm", ha="right", va="bottom", fontsize=8, color="#ef4444")
                ax.set_xticks(xm); ax.set_xticklabels(ML, fontsize=10)
                ax.set_ylabel("Rainfall (mm)", fontsize=11, fontweight="bold")
                ax.set_title(f"Monthly Rainfall Totals – {rainfall_year}", fontsize=13, fontweight="bold", pad=10, color="#333")
                ax.legend(handles=[
                    mpatches.Patch(color="#93c5fd", label="< 50 mm"),
                    mpatches.Patch(color="#3b82f6", label="50–150 mm"),
                    mpatches.Patch(color="#1d4ed8", label="150–300 mm"),
                    mpatches.Patch(color="#1e3a5f", label="≥ 300 mm"),
                ], loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
                ax.grid(True, alpha=0.25, linestyle="--", axis="y")
                ax.set_facecolor("#fafafa"); fig.patch.set_facecolor("white"); plt.tight_layout()
                tmp = _save_mpl_figure(fig); plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP), width=Inches(SW - 0.54), height=Inches(CHART_H))
                os.unlink(tmp)
                _kpi_row(slide, [
                    ("Annual Total",  f"{annual_total:.0f} mm",          "Full year",              "#1d4ed8"),
                    ("Wettest Month", ML[wettest_idx - 1],                f"{monthly[wettest_idx]:.0f} mm", "#1e3a5f"),
                    ("Mean Monthly",  f"{annual_mean:.0f} mm",           "Annual ÷ 12",            "#0891b2"),
                ], top=KPI_TOP, card_h=KPI_H)
            except Exception as e:
                _err_box(slide, e)
            _add_logo(slide)

        _monthly_rainfall()

        # Rainy Days
        def _rainy_days():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, f"Rainy Days by Intensity  —  {rainfall_station_name}, {rainfall_year}")
            _add_divider(slide, 0.62)
            try:
                rainy = df_f[df_f["prcp_mm"] > 0].copy()
                def _cnt(lo, hi=None):
                    m = rainy["prcp_mm"] >= lo
                    if hi is not None:
                        m &= rainy["prcp_mm"] < hi
                    return rainy[m].groupby("month").size().reindex(range(1, 13), fill_value=0)
                light    = _cnt(0.001, 10)
                moderate = _cnt(10, 25)
                heavy    = _cnt(25, rainfall_heavy_threshold)
                extreme  = _cnt(rainfall_heavy_threshold)
                fig, ax = plt.subplots(figsize=(13, 4.5), dpi=100)
                ax.bar(xm, light.values, color="#bfdbfe", label="Light (< 10 mm)")
                ax.bar(xm, moderate.values, bottom=light.values, color="#3b82f6", label="Moderate (10–25 mm)")
                ax.bar(xm, heavy.values, bottom=(light + moderate).values, color="#1d4ed8", label=f"Heavy (25–{rainfall_heavy_threshold:.0f} mm)")
                ax.bar(xm, extreme.values, bottom=(light + moderate + heavy).values, color="#ef4444", label=f"Extreme (≥ {rainfall_heavy_threshold:.0f} mm)")
                ax.set_xticks(xm); ax.set_xticklabels(ML, fontsize=10)
                ax.set_ylabel("Number of Days", fontsize=11, fontweight="bold")
                ax.set_title(f"Rainy Days per Month by Intensity – {rainfall_year}", fontsize=13, fontweight="bold", pad=10, color="#333")
                ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
                ax.grid(True, alpha=0.25, linestyle="--", axis="y")
                ax.set_facecolor("#fafafa"); fig.patch.set_facecolor("white"); plt.tight_layout()
                tmp = _save_mpl_figure(fig); plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP), width=Inches(SW - 0.54), height=Inches(CHART_H))
                os.unlink(tmp)
                _kpi_row(slide, [
                    ("Total Rainy Days",                       str(int(rainy.shape[0])),   "Days with rain > 0",              "#1d4ed8"),
                    ("Light < 10 mm",                          str(int(light.sum())),       "Days",                            "#0891b2"),
                    ("Moderate 10–25",                         str(int(moderate.sum())),    "Days",                            "#2563eb"),
                    (f"Heavy 25–{rainfall_heavy_threshold:.0f}mm", str(int(heavy.sum())),   "Days",                            "#1e3a5f"),
                    (f"Extreme ≥{rainfall_heavy_threshold:.0f}mm", str(int(extreme.sum())), "Days",                            "#dc2626"),
                ], top=KPI_TOP, card_h=KPI_H)
            except Exception as e:
                _err_box(slide, e)
            _add_logo(slide)

        _rainy_days()

        # GI Balance
        def _gi_balance():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, f"Rainwater Harvesting Potential  —  {rainfall_station_name}, {rainfall_year}")
            _add_divider(slide, 0.62)
            try:
                if _percentile_future is not None:
                    try:
                        gi_result = _percentile_future.result(timeout=90)
                    except Exception:
                        gi_result = _fetch_percentile_depth(rainfall_station_id, rainfall_gi_percentile, rainfall_gi_start_year)
                else:
                    gi_result = _fetch_percentile_depth(rainfall_station_id, rainfall_gi_percentile, rainfall_gi_start_year)
                if "error" in gi_result:
                    _err_box(slide, f"Baseline fetch failed: {gi_result['error']}")
                    _add_logo(slide)
                    return
                baseline_mm = gi_result["raw_mm"]
                daily = df_rain.sort_values("date").copy()
                daily["stored"]   = daily["prcp_mm"].clip(upper=baseline_mm)
                daily["overflow"] = (daily["prcp_mm"] - baseline_mm).clip(lower=0)
                mgrp = daily.groupby("month").agg(stored=("stored", "sum"), overflow=("overflow", "sum")).reindex(range(1, 13), fill_value=0)
                total_recharge = float(daily["stored"].sum())
                total_overflow = float(daily["overflow"].sum())
                overflow_days  = int((daily["overflow"] > 0).sum())
                worst_m_idx    = int(mgrp["overflow"].idxmax()) if total_overflow > 0 else 1
                fig, ax = plt.subplots(figsize=(13, 4.5), dpi=100)
                ax.bar(xm, mgrp["stored"].values, color="#22c55e", label="Stored (L/m²)")
                ax.bar(xm, -mgrp["overflow"].values, color="#ef4444", label="Overflow (L/m²)")
                ax.axhline(0, color="#374151", linewidth=1.2, linestyle="--")
                ax.set_xticks(xm); ax.set_xticklabels(ML, fontsize=10)
                ax.set_ylabel("Volume (L/m²)", fontsize=11, fontweight="bold")
                ax.set_title(
                    f"Monthly Rainwater Harvesting Potential – {rainfall_year}  "
                    f"({rainfall_gi_percentile}th-percentile baseline: {baseline_mm:.1f} mm/day)",
                    fontsize=12, fontweight="bold", pad=10, color="#333",
                )
                ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=2, frameon=True, fontsize=9)
                ax.grid(True, alpha=0.25, linestyle="--", axis="y")
                ax.set_facecolor("#fafafa"); fig.patch.set_facecolor("white"); plt.tight_layout()
                tmp = _save_mpl_figure(fig); plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP), width=Inches(SW - 0.54), height=Inches(CHART_H))
                os.unlink(tmp)
                worst_val = float(mgrp.loc[worst_m_idx, "overflow"])
                _kpi_row(slide, [
                    ("Storage Potential",    f"{total_recharge:,.0f} L/m²", "Total captured by GI",       "#16a34a"),
                    ("Recharge Potential",   f"{total_overflow:,.0f} L/m²", "Excess beyond GI capacity",   "#dc2626"),
                    ("Overflow Days",        str(overflow_days),             f"Rain > {baseline_mm:.0f} mm","#d97706"),
                    ("Worst Overflow Month", ML[worst_m_idx - 1],           f"{worst_val:,.0f} L/m²",      "#7c3aed"),
                ], top=KPI_TOP, card_h=KPI_H)
            except Exception as e:
                _err_box(slide, e)
            _add_logo(slide)

        _gi_balance()

        # Surface Runoff
        def _surface_runoff():
            CHART_H2 = 3.70
            ROW1_TOP = CHART_TOP + CHART_H2 + 0.13
            ROW2_TOP = ROW1_TOP + 0.92 + 0.10
            ROW_H    = 0.88
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, f"Surface Runoff Analysis  —  {rainfall_station_name}, {rainfall_year}")
            _add_divider(slide, 0.62)
            try:
                monthly_prcp = df_f.groupby("month")["prcp_mm"].sum().reindex(range(1, 13), fill_value=0)
                monthly_vols = {s["key"]: (monthly_prcp / 1000.0) * _sa.get(s["key"], 0.0) * s["rc"] for s in _RUNOFF_SURFACES}
                fig, ax = plt.subplots(figsize=(13, 4.0), dpi=100)
                bottom = np.zeros(12)
                for surf in _RUNOFF_SURFACES:
                    vals = monthly_vols[surf["key"]].values
                    ax.bar(xm, vals, bottom=bottom, color=surf["color"], label=f"{surf['label']}  (RC {surf['rc']:.2f})")
                    bottom += vals
                ax.set_xticks(xm); ax.set_xticklabels(ML, fontsize=10)
                ax.set_ylabel("Runoff Volume (m³)", fontsize=11, fontweight="bold")
                ax.set_title(f"Monthly Surface Runoff by Type – {rainfall_year}", fontsize=13, fontweight="bold", pad=10, color="#333")
                ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.13), ncol=2, frameon=True, fontsize=8)
                ax.grid(True, alpha=0.25, linestyle="--", axis="y")
                ax.set_facecolor("#fafafa"); fig.patch.set_facecolor("white"); plt.tight_layout()
                tmp = _save_mpl_figure(fig); plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(CHART_TOP), width=Inches(SW - 0.54), height=Inches(CHART_H2))
                os.unlink(tmp)
                surf_df      = pd.DataFrame({s["key"]: monthly_vols[s["key"]] for s in _RUNOFF_SURFACES})
                total_m      = surf_df.sum(axis=1)
                total_annual = float(total_m.sum())
                peak_m       = int(total_m.idxmax()) if total_annual > 0 else 1
                _kpi_row(slide, [
                    ("Total Annual Runoff", _fmt_L(total_annual * 1000), "All surfaces combined",           "#1d4ed8"),
                    ("Peak Runoff Month",   ML[peak_m - 1],              _fmt_L(float(total_m[peak_m]) * 1000), "#dc2626"),
                ], top=ROW1_TOP, card_h=ROW_H)
                short = ["Roof", "Paved", "Green Area", "Waterbody"]
                _kpi_row(slide, [
                    (f"{short[i]}  RC {s['rc']:.2f}", _fmt_L(float(monthly_vols[s["key"]].sum()) * 1000), "Annual runoff", s["color"])
                    for i, s in enumerate(_RUNOFF_SURFACES)
                ], top=ROW2_TOP, card_h=ROW_H)
            except Exception as e:
                _err_box(slide, e)
            _add_logo(slide)

        _surface_runoff()

    _prepare_rainfall_slides()

    if _rain_executor is not None:
        _rain_executor.shutdown(wait=False)

    # ── SECTION 11 – SOLAR PV POTENTIAL ───────────────────────────────────────
    def _prepare_solar_pv_slides():
        import pathlib

        _MONTHS_PV = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                      "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
        _DAYS_PV   = [31, 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31]

        excel_path = pathlib.Path(base_dir) / "solargis_country_pv_data.xlsx"
        if not excel_path.exists():
            return
        try:
            _raw = pd.read_excel(excel_path, sheet_name="Monthly data", header=None)
            _raw.columns = _raw.iloc[1]
            _raw = _raw.iloc[2:].reset_index(drop=True)
            _month_cols = ["January", "February", "March", "April", "May", "June",
                           "July", "August", "September", "October", "November", "December"]
            _pv = pd.DataFrame()
            _pv["country"]      = _raw["Country or region"].astype(str).str.strip()
            _pv["yearly_daily"] = pd.to_numeric(_raw["Yearly"], errors="coerce")
            for _col in _month_cols:
                _pv[_col.lower()[:3]] = pd.to_numeric(_raw[_col], errors="coerce")
            _mask = _pv["country"] == solar_pv_country
            if not _mask.any():
                return
            _row = _pv[_mask].iloc[0]
        except Exception:
            return

        _daily_vals   = [float(_row[m]) for m in ["jan", "feb", "mar", "apr", "may", "jun",
                                                    "jul", "aug", "sep", "oct", "nov", "dec"]]
        _monthly_vals = [d * days for d, days in zip(_daily_vals, _DAYS_PV)]
        _annual_daily = float(_row["yearly_daily"])
        _annual_total = sum(_monthly_vals)
        _peak_idx     = _monthly_vals.index(max(_monthly_vals))
        _low_idx      = _monthly_vals.index(min(_monthly_vals))

        _eff_area     = solar_pv_roof_size_m2 * (solar_pv_roof_pct / 100.0)
        _sys_kwp      = _eff_area / 10.0
        _monthly_kwh  = [v * _sys_kwp for v in _monthly_vals]
        _daily_kwh    = [v * _sys_kwp for v in _daily_vals]
        _annual_kwh   = sum(_monthly_kwh)
        _ann_daily_kw = _annual_daily * _sys_kwp

        _C_AMBER = RGBColor(0xF5, 0x9E, 0x0B)

        # ── KPI card helpers (scoped to solar PV section) ─────────────────────
        def _hex_rgb_pv(h):
            return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

        def _pv_kpi_card(slide, left, top, width, height, label, value, sub, hex_color):
            rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex_rgb_pv(hex_color)
            rect.line.fill.background()
            _lc = RGBColor(0xFF, 0xE0, 0xB2)
            _sc = RGBColor(0xFF, 0xF0, 0xD0)
            tb = slide.shapes.add_textbox(Inches(left + 0.07), Inches(top + 0.06),
                                          Inches(width - 0.14), Inches(0.20))
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = label; r.font.size = Pt(9); r.font.color.rgb = _lc
            val_top = top + (height - 0.28) / 2 - 0.05
            tbv = slide.shapes.add_textbox(Inches(left + 0.05), Inches(val_top),
                                           Inches(width - 0.10), Inches(0.50))
            pv = tbv.text_frame.paragraphs[0]
            pv.alignment = PP_ALIGN.CENTER
            rv = pv.add_run(); rv.text = value
            rv.font.size = Pt(26) if height >= 0.80 else Pt(22)
            rv.font.bold = True; rv.font.color.rgb = WHITE
            tbs = slide.shapes.add_textbox(Inches(left + 0.07), Inches(top + height - 0.22),
                                           Inches(width - 0.14), Inches(0.20))
            ps = tbs.text_frame.paragraphs[0]
            ps.alignment = PP_ALIGN.CENTER
            rs = ps.add_run(); rs.text = sub; rs.font.size = Pt(8); rs.font.color.rgb = _sc

        def _pv_kpi_row(slide, cards, top, card_h=0.88):
            n = len(cards)
            gap = 0.12
            card_w = (SW - 0.54 - gap * (n - 1)) / n
            left = 0.27
            for lbl, val, sub, clr in cards:
                _pv_kpi_card(slide, left, top, card_w, card_h, lbl, val, sub, clr)
                left += card_w + gap

        # ── Section header ────────────────────────────────────────────────────
        def _pv_section_header():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = _C_AMBER
            bg.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(SW - 1.2), Inches(1.1))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = "Solar PV Potential"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = WHITE
            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.7))
            run2 = tb2.text_frame.paragraphs[0].add_run()
            run2.text = (
                f"{solar_pv_country}  |  "
                f"{solar_pv_roof_size_m2:.0f} m² roof × {solar_pv_roof_pct:.0f}% coverage  "
                f"→  {_sys_kwp:.2f} kWp system"
            )
            run2.font.size = Pt(15)
            run2.font.color.rgb = RGBColor(0xFF, 0xF0, 0xCC)
            _add_logo(slide)
        _pv_section_header()

        # ── Reusable dual-axis chart builder ──────────────────────────────────
        def _pv_chart(monthly_y, daily_y, y1_lbl, y2_lbl, title):
            _H = "#f59e0b"
            _G = "#10b981"
            _R = "#ef4444"
            _B = "#3b82f6"
            bar_colors = [_G if i == _peak_idx else _R if i == _low_idx else _H for i in range(12)]
            x = np.arange(12)
            fig, ax1 = plt.subplots(figsize=(13, 4.6), dpi=100)
            ax2 = ax1.twinx()
            ax1.bar(x, monthly_y, color=bar_colors, alpha=0.88, zorder=2)
            ax2.scatter(x, daily_y, color=_B, s=80, zorder=3, edgecolors="white", linewidths=1.5)
            ax1.set_xticks(x)
            ax1.set_xticklabels(_MONTHS_PV, fontsize=10)
            ax1.set_ylabel(y1_lbl, fontsize=10, fontweight="bold")
            ax2.set_ylabel(y2_lbl, fontsize=10, fontweight="bold")
            ax1.set_title(title, fontsize=13, fontweight="bold", pad=10, color="#333")
            ax1.grid(True, alpha=0.25, linestyle="--", axis="y")
            ax1.set_facecolor("#fafafa")
            fig.patch.set_facecolor("white")
            from matplotlib.patches import Patch
            from matplotlib.lines import Line2D
            legend_elems = [
                Patch(facecolor=_G, label="Peak Month"),
                Patch(facecolor=_R, label="Lowest Month"),
                Patch(facecolor=_H, label="Monthly Total (bars)"),
                Line2D([0], [0], marker="o", color="w", markerfacecolor=_B, markersize=8, label="Daily Avg (points)"),
            ]
            ax1.legend(handles=legend_elems, loc="upper center",
                       bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
            plt.tight_layout()
            return fig

        _PV_CHART_TOP = 0.72
        _PV_CHART_H   = 4.0
        _PV_KPI_TOP   = 4.86
        _PV_KPI_H     = 0.88

        # ── Slide 1: Per-kWp reference ────────────────────────────────────────
        def _pv_per_kwp_slide():
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            _add_slide_title(slide, f"Solar PV Yield — {solar_pv_country}  (Per kWp Reference)")
            _add_divider(slide, 0.62)
            try:
                fig = _pv_chart(
                    _monthly_vals, _daily_vals,
                    "Monthly Total (kWh / kWp)", "Daily Avg (kWh / kWp.day)",
                    f"Monthly PV Yield — {solar_pv_country}",
                )
                tmp = _save_mpl_figure(fig)
                plt.close(fig)
                slide.shapes.add_picture(tmp, Inches(0.27), Inches(_PV_CHART_TOP),
                                         width=Inches(SW - 0.54), height=Inches(_PV_CHART_H))
                os.unlink(tmp)
                _pv_kpi_row(slide, [
                    ("Annual Yield",         f"{_annual_total:,.0f}",          "kWh / kWp",                           "#f97316"),
                    ("Annual Daily Average", f"{_annual_daily:.2f}",           "kWh / kWp.day",                     "#f97316"),
                    ("Peak Month Total",     f"{_monthly_vals[_peak_idx]:.0f}", f"kWh/kWp · {_MONTHS_PV[_peak_idx]}", "#10b981"),
                    ("Peak Day Average",     f"{_daily_vals[_peak_idx]:.2f}",  f"kWh/kWp.d · {_MONTHS_PV[_peak_idx]}", "#f59e0b"),
                ], top=_PV_KPI_TOP, card_h=_PV_KPI_H)
            except Exception as e:
                _err_box(slide, e)
            _add_logo(slide)
        _pv_per_kwp_slide()

        # ── Slide 2: Absolute system yield ────────────────────────────────────
        if _sys_kwp > 0:
            def _pv_absolute_slide():
                slide = prs.slides.add_slide(BLANK_LAYOUT)
                _add_slide_title(slide, f"Solar PV System Yield — {solar_pv_country}  ({_sys_kwp:.2f} kWp)")
                _add_divider(slide, 0.62)
                try:
                    fig = _pv_chart(
                        _monthly_kwh, _daily_kwh,
                        "Monthly Total (kWh)", "Daily Avg (kWh / day)",
                        f"{_sys_kwp:.2f} kWp System — Monthly Output  "
                        f"({solar_pv_roof_size_m2:.0f} m² × {solar_pv_roof_pct:.0f}% = {_eff_area:.1f} m²)",
                    )
                    tmp = _save_mpl_figure(fig)
                    plt.close(fig)
                    slide.shapes.add_picture(tmp, Inches(0.27), Inches(_PV_CHART_TOP),
                                             width=Inches(SW - 0.54), height=Inches(_PV_CHART_H))
                    os.unlink(tmp)
                    _pv_kpi_row(slide, [
                        ("Annual Yield (System)", f"{_annual_kwh:,.0f}",          "kWh / year",                           "#f97316"),
                        ("Annual Daily Average",  f"{_ann_daily_kw:.1f}",         "kWh / day",                            "#f97316"),
                        ("Peak Month Output",     f"{_monthly_kwh[_peak_idx]:,.0f}", f"kWh · {_MONTHS_PV[_peak_idx]}",    "#10b981"),
                        ("Peak Day Output",       f"{_daily_kwh[_peak_idx]:.1f}", f"kWh/d · {_MONTHS_PV[_peak_idx]}",    "#f59e0b"),
                    ], top=_PV_KPI_TOP, card_h=_PV_KPI_H)
                except Exception as e:
                    _err_box(slide, e)
                _add_logo(slide)
            _pv_absolute_slide()

    _prepare_solar_pv_slides()

    # ── ANNEXURE SLIDE ────────────────────────────────────────────────────────
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Annexure")
        _add_divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = (
            "Environmental Design Solutions [EDS] is a sustainability advisory firm. "
            "Since 2002, EDS has worked on over 500 green building and energy efficiency "
            "projects worldwide. The team focuses on climate change mitigation, low-carbon "
            "design, building simulation, performance audits, and capacity building. EDS "
            "continues to contribute to the buildings community with useful tools through "
            "its IT services."
        )
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        for item in [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "• EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "• EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "• Open-source resources used: Clima - Berkley, Streamlit, Python",
            "• EDS is not liable to inform Users about updates to the Application or underlying resources",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        for item in [
            "• Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "• Streamlit, © Streamlit Inc., licensed under Apache 2.0",
            "• Python © Python Software Foundation, licensed under PSF License Version 2",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_annexure_slide()

    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes
