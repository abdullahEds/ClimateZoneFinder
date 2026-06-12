"""Solar PV PPT Report Generator.

Generates ~6 slides for the Solar PV section, matching the style of
existing ppt_report.py generators. Called from report_api.py and
combined_report.py.

Public API
----------
    generate_solar_pv_pptx_report(epw_df, metadata, output_path=None)
        → BytesIO  (or writes to output_path)
"""

import calendar
import io
import pathlib
import tempfile

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.colors import LinearSegmentedColormap

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pages.modules.solar_pv_module import (
    compute_site_solar, compute_annual_yield, compute_economics,
    _load_solargis_data, _resolve_iso, _pr_from_ghi,
    _MONTHS, _CO2_FACTORS, _PANEL_AREA_PER_KWP,
    _pvout_label, _pvout_color,
)

# ─── Colours (match existing brand) ──────────────────────────────────────────
BRAND_ORANGE  = RGBColor(0xA8, 0x5C, 0x42)
BRAND_YELLOW  = RGBColor(0xF4, 0xB9, 0x42)
LIGHT_GREY    = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY      = RGBColor(0x88, 0x88, 0x88)
DARK          = RGBColor(0x1A, 0x1A, 0x1A)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GREEN         = RGBColor(0x27, 0xAE, 0x60)
BLUE          = RGBColor(0x29, 0x80, 0xB9)

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)

_TEMPLATE_PATH = pathlib.Path(__file__).parents[2] / \
    "Voha Hospitality Climate analysis_v4 (2).pptx"


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _pt(val): return Pt(val)
def _in(val): return Inches(val)


def _add_text(slide, text, left, top, width, height,
              font_size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT,
              wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = _pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _add_rect(slide, left, top, width, height, fill_color=LIGHT_GREY,
              line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _kpi_box(slide, label, value, unit, left, top, width=_in(2.8), height=_in(1.0),
             accent=BRAND_ORANGE):
    _add_rect(slide, left, top, width, height, fill_color=WHITE,
              line_color=RGBColor(0xE0, 0xE0, 0xE0))
    # accent bar
    _add_rect(slide, left, top, _in(0.04), height, fill_color=accent)
    _add_text(slide, label, left + _in(0.12), top + _in(0.05),
              width - _in(0.15), _in(0.25), font_size=8, color=MID_GREY)
    _add_text(slide, value, left + _in(0.12), top + _in(0.28),
              width - _in(0.15), _in(0.45), font_size=18, bold=True, color=DARK)
    _add_text(slide, unit, left + _in(0.12), top + _in(0.72),
              width - _in(0.15), _in(0.22), font_size=8, color=MID_GREY)


def _add_image(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    buf.seek(0)
    plt.close(fig)
    slide.shapes.add_picture(buf, left, top, width, height)


def _slide_title(slide, title, subtitle=None):
    _add_rect(slide, _in(0), _in(0), _SLIDE_W, _in(0.7), fill_color=BRAND_ORANGE)
    _add_text(slide, title, _in(0.3), _in(0.1), _in(10), _in(0.55),
              font_size=20, bold=True, color=WHITE)
    if subtitle:
        _add_text(slide, subtitle, _in(10.5), _in(0.22), _in(2.6), _in(0.35),
                  font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def _load_template():
    if _TEMPLATE_PATH.exists():
        prs = Presentation(str(_TEMPLATE_PATH))
        # Remove all existing slides
        xml_slides = prs.slides._sldIdLst
        for slide in list(prs.slides):
            xml_slides.remove(slide._element)
        return prs
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    return prs


def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)


# ─── Chart helpers ────────────────────────────────────────────────────────────

def _monthly_bar_line(site, title="Monthly GHI & PVOUT"):
    fig, ax1 = plt.subplots(figsize=(9, 4))
    ax2 = ax1.twinx()
    x = np.arange(12)
    ghi_vals = [site.loc[m, "ghi_daily"] for m in range(1, 13)]
    pvout_vals = [site.loc[m, "pvout"] for m in range(1, 13)]
    bars = ax1.bar(x, ghi_vals, color="#F4B942", alpha=0.85, label="GHI (kWh/m²/day)")
    line, = ax2.plot(x, pvout_vals, color="#A85C42", linewidth=2.5,
                     marker="o", markersize=5, label="PVOUT (kWh/kWp/day)")
    ax1.set_xticks(x)
    ax1.set_xticklabels(_MONTHS, fontsize=9)
    ax1.set_ylabel("GHI (kWh/m²/day)", fontsize=9)
    ax2.set_ylabel("PVOUT (kWh/kWp/day)", fontsize=9, color="#A85C42")
    ax2.tick_params(axis="y", labelcolor="#A85C42")
    ax1.set_facecolor("#fafafa")
    fig.patch.set_facecolor("white")
    fig.legend(loc="upper right", bbox_to_anchor=(1, 1), bbox_transform=ax1.transAxes,
               fontsize=8)
    ax1.set_title(title, fontsize=11, fontweight="bold", pad=6)
    fig.tight_layout()
    return fig


def _monthly_yield_bar(monthly_yield, system_kWp, title="Monthly PV Yield"):
    fig, ax = plt.subplots(figsize=(9, 4))
    x = np.arange(12)
    vals = [monthly_yield.loc[m, "monthly_yield_kWh"] if m in monthly_yield.index else 0
            for m in range(1, 13)]
    ax.bar(x, vals, color="#A85C42", alpha=0.85)
    for i, v in enumerate(vals):
        ax.text(i, v + max(vals) * 0.01, f"{v:,.0f}", ha="center",
                va="bottom", fontsize=7)
    ax.set_xticks(x)
    ax.set_xticklabels(_MONTHS, fontsize=9)
    ax.set_ylabel("Energy Yield (kWh)", fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", pad=6)
    ax.set_facecolor("#fafafa")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    return fig


def _benchmark_bar(country_monthly, site_vals, country_name):
    fig, ax = plt.subplots(figsize=(9, 4))
    x = np.arange(12)
    ax.bar(x - 0.2, site_vals, width=0.38, color="#A85C42", label="Site (EPW)", alpha=0.9)
    ax.bar(x + 0.2, country_monthly, width=0.38, color="#5B9BD5",
           label=f"Country Avg ({country_name})", alpha=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(_MONTHS, fontsize=9)
    ax.set_ylabel("PVOUT (kWh/kWp/day)", fontsize=9)
    ax.set_title("Monthly PVOUT: Site vs Country Average", fontsize=11,
                 fontweight="bold", pad=6)
    ax.legend(fontsize=8)
    ax.set_facecolor("#fafafa")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    return fig


def _payback_sensitivity(system_kWp, cost_per_kWp, annual_kWh, current_tariff):
    fig, ax = plt.subplots(figsize=(8, 3.5))
    tariffs = np.linspace(0.05, 0.30, 50)
    paybacks = [(system_kWp * cost_per_kWp) / (annual_kWh * t) for t in tariffs]
    ax.fill_between(tariffs, paybacks, alpha=0.15, color="#A85C42")
    ax.plot(tariffs, paybacks, color="#A85C42", linewidth=2)
    ax.axvline(current_tariff, color="#A85C42", linestyle="--", linewidth=1.5)
    ax.axhline(10, color="#27ae60", linestyle=":", linewidth=1.5,
               label="10-year target")
    ax.set_xlabel("Electricity Tariff (USD/kWh)", fontsize=9)
    ax.set_ylabel("Simple Payback (years)", fontsize=9)
    ax.set_title("Payback Period Sensitivity", fontsize=11, fontweight="bold", pad=6)
    ax.legend(fontsize=8)
    ax.set_facecolor("#fafafa")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    return fig


def _irradiance_heatmap(df):
    """24×12 GHI heatmap."""
    d = df.copy()
    d["month"] = d["datetime"].dt.month
    d["hour"] = d["datetime"].dt.hour
    pivot = d.groupby(["hour", "month"])["global_horizontal_irradiance"].mean().unstack("month")

    fig, ax = plt.subplots(figsize=(9, 4))
    cmap = LinearSegmentedColormap.from_list("solar", ["#fff9e6", "#F4B942", "#A85C42"])
    im = ax.imshow(pivot.values, aspect="auto", cmap=cmap, origin="upper")
    plt.colorbar(im, ax=ax, label="GHI (Wh/m²)", shrink=0.8)
    ax.set_xticks(range(12))
    ax.set_xticklabels(_MONTHS, fontsize=8)
    ax.set_yticks(range(0, 24, 2))
    ax.set_yticklabels([f"{h:02d}:00" for h in range(0, 24, 2)], fontsize=7)
    ax.set_title("Average Hourly GHI by Month", fontsize=11, fontweight="bold")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    return fig


# ─── Slide builders ───────────────────────────────────────────────────────────

def _slide_title_card(prs, city, lat, lon, pvout_annual, site):
    """Slide 1: Title + key metrics overview."""
    slide = _blank_slide(prs)
    _add_rect(slide, _in(0), _in(0), _SLIDE_W, _SLIDE_H, fill_color=LIGHT_GREY)
    _add_rect(slide, _in(0), _in(0), _in(13.33), _in(2.4), fill_color=BRAND_ORANGE)

    _add_text(slide, "Solar PV Potential Analysis", _in(0.5), _in(0.3),
              _in(10), _in(0.7), font_size=28, bold=True, color=WHITE)
    loc_str = f"📍 {city}"
    if lat and lon:
        loc_str += f"  |  {abs(lat):.2f}°{'N' if lat >= 0 else 'S'}  {abs(lon):.2f}°{'E' if lon >= 0 else 'W'}"
    _add_text(slide, loc_str, _in(0.5), _in(1.05), _in(10), _in(0.4),
              font_size=13, color=WHITE)
    _add_text(slide, "Powered by EPW data + Solargis Global PV Potential",
              _in(0.5), _in(1.55), _in(9), _in(0.35), font_size=9, color=WHITE)

    # KPIs
    kpi_data = [
        ("Annual Avg PVOUT",     f"{pvout_annual:.2f}", "kWh/kWp/day",  BRAND_ORANGE),
        ("Annual Avg GHI",       f"{site['ghi_daily'].mean():.2f}", "kWh/m²/day",  BLUE),
        ("Peak Sun Hours",       f"{site['psh'].max():.2f}", "hrs/day (best month)", RGBColor(0xE6, 0x7E, 0x22)),
        ("PV Potential Rating",  _pvout_label(pvout_annual), "", GREEN),
    ]
    for i, (label, val, unit, color) in enumerate(kpi_data):
        _kpi_box(slide, label, val, unit,
                 left=_in(0.5 + i * 3.1), top=_in(2.8),
                 width=_in(2.9), height=_in(1.1), accent=color)

    # Mini bar chart
    fig, ax = plt.subplots(figsize=(9, 2.8))
    pvout_vals = [site.loc[m, "pvout"] for m in range(1, 13)]
    colors_bar = ["#27ae60" if v >= 5 else "#F4B942" if v >= 4 else "#e67e22"
                  for v in pvout_vals]
    ax.bar(_MONTHS, pvout_vals, color=colors_bar, alpha=0.9)
    ax.axhline(pvout_annual, color="#A85C42", linestyle="--", linewidth=1.5,
               label=f"Annual avg: {pvout_annual:.2f}")
    ax.set_ylabel("PVOUT (kWh/kWp/day)", fontsize=9)
    ax.legend(fontsize=8)
    ax.set_facecolor("#fafafa")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    _add_image(slide, fig, _in(0.5), _in(4.1), _in(12.3), _in(3.0))
    return slide


def _slide_solar_resource(prs, site, df):
    """Slide 2: GHI/PVOUT bar+line + heatmap."""
    slide = _blank_slide(prs)
    _slide_title(slide, "Solar Resource Analysis", "GHI & PVOUT from EPW Data")

    fig1 = _monthly_bar_line(site)
    _add_image(slide, fig1, _in(0.3), _in(0.85), _in(6.5), _in(3.0))

    fig2 = _irradiance_heatmap(df)
    _add_image(slide, fig2, _in(6.9), _in(0.85), _in(6.1), _in(3.0))

    # Bottom data table
    tbl_data = {
        "Month": _MONTHS,
        "GHI (kWh/m²/d)": [f"{site.loc[m,'ghi_daily']:.2f}" for m in range(1, 13)],
        "PVOUT (kWh/kWp/d)": [f"{site.loc[m,'pvout']:.2f}" for m in range(1, 13)],
        "Peak Sun Hours": [f"{site.loc[m,'psh']:.2f}" for m in range(1, 13)],
    }
    y_start = _in(4.1)
    col_widths = [_in(0.7), _in(1.1), _in(1.2), _in(1.1)]
    col_starts = [_in(0.3), _in(1.1), _in(2.3), _in(3.6)]

    for c_idx, (col_name, col_start, col_w) in enumerate(
            zip(tbl_data.keys(), col_starts, col_widths)):
        _add_rect(slide, col_start, y_start, col_w, _in(0.28), BRAND_ORANGE)
        _add_text(slide, col_name, col_start + _in(0.03), y_start + _in(0.04),
                  col_w - _in(0.05), _in(0.22), font_size=7, bold=True, color=WHITE)
        vals = list(tbl_data.values())[c_idx]
        for r_idx, val in enumerate(vals):
            bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
            _add_rect(slide, col_start, y_start + _in(0.28 + r_idx * 0.22),
                      col_w, _in(0.22), bg)
            _add_text(slide, str(val),
                      col_start + _in(0.03),
                      y_start + _in(0.30 + r_idx * 0.22),
                      col_w - _in(0.05), _in(0.20), font_size=7)

    return slide


def _slide_yield(prs, site, system_kWp=100.0):
    """Slide 3: System yield for default 100 kWp."""
    slide = _blank_slide(prs)
    _slide_title(slide, "PV System Yield Estimate", f"Based on {system_kWp:.0f} kWp System")

    result = compute_annual_yield(site, system_kWp)
    monthly_yield = result["monthly"]

    kpis = [
        ("Annual Yield",     f"{result['annual_kWh']:,.0f}", "kWh/year",       BRAND_ORANGE),
        ("Specific Yield",   f"{result['specific_yield_kWh_kWp']:,.0f}", "kWh/kWp/year", BLUE),
        ("Capacity Factor",  f"{result['capacity_factor_pct']:.1f}%", "annual",  GREEN),
        ("Daily Output",     f"{result['annual_kWh']/365:,.1f}", "kWh/day avg", RGBColor(0xE6, 0x7E, 0x22)),
    ]
    for i, (label, val, unit, color) in enumerate(kpis):
        _kpi_box(slide, label, val, unit,
                 left=_in(0.3 + i * 3.25), top=_in(0.9),
                 width=_in(3.1), height=_in(1.0), accent=color)

    fig = _monthly_yield_bar(monthly_yield, system_kWp,
                              title=f"Monthly PV Yield — {system_kWp:.0f} kWp System")
    _add_image(slide, fig, _in(0.3), _in(2.1), _in(12.7), _in(5.0))
    return slide


def _slide_benchmark(prs, site, iso, monthly_sol, ci):
    """Slide 4: Country benchmark."""
    slide = _blank_slide(prs)
    _slide_title(slide, "Country & Regional Benchmark", "Solargis Global PV Potential 2020")

    site_pvout = site["pvout"].mean()
    country_row = ci[ci["ISO_A3"] == iso] if iso and not ci.empty else pd.DataFrame()
    country_name = country_row.iloc[0]["Country"] if not country_row.empty else "N/A"
    country_pvout = float(country_row.iloc[0]["PVOUT_avg"]) if not country_row.empty else None
    country_lcoe = float(country_row.iloc[0]["LCOE_USD_kWh"]) if not country_row.empty else None
    country_region = country_row.iloc[0]["Region"] if not country_row.empty else None

    global_sorted = ci["PVOUT_avg"].dropna().sort_values(ascending=False)
    rank = int((global_sorted > (country_pvout or 0)).sum()) + 1 if country_pvout else None

    kpis = [
        ("Site PVOUT (EPW)",   f"{site_pvout:.2f}", "kWh/kWp/day",   BRAND_ORANGE),
        ("Country Avg PVOUT",  f"{country_pvout:.2f}" if country_pvout else "N/A",
         country_name, BLUE),
        ("Global Rank",        f"#{rank}" if rank else "N/A",
         f"of {len(global_sorted)} countries", RGBColor(0xE6, 0x7E, 0x22)),
        ("Country LCOE",       f"${country_lcoe:.3f}/kWh" if country_lcoe else "N/A",
         "Levelised Cost of Energy", GREEN),
    ]
    for i, (label, val, unit, color) in enumerate(kpis):
        _kpi_box(slide, label, val, unit,
                 left=_in(0.3 + i * 3.25), top=_in(0.9),
                 width=_in(3.1), height=_in(1.0), accent=color)

    # Monthly comparison chart
    if iso and not monthly_sol.empty:
        row = monthly_sol[monthly_sol["ISO_A3"] == iso]
        if not row.empty:
            bench = [float(row.iloc[0][m]) for m in _MONTHS]
            site_vals = [site.loc[m, "pvout"] for m in range(1, 13)]
            fig = _benchmark_bar(bench, site_vals, country_name)
            _add_image(slide, fig, _in(0.3), _in(2.1), _in(8.5), _in(5.0))

    # Regional top-5 sidebar
    if country_region and not ci.empty:
        region_top = (ci[ci["Region"] == country_region]
                      .dropna(subset=["PVOUT_avg"])
                      .sort_values("PVOUT_avg", ascending=False)
                      .head(8))
        _add_text(slide, f"Top 8 — {country_region}", _in(9.1), _in(2.1),
                  _in(4.0), _in(0.3), font_size=9, bold=True, color=DARK)
        for idx, (_, r_row) in enumerate(region_top.iterrows()):
            bg = BRAND_ORANGE if r_row["ISO_A3"] == iso else (LIGHT_GREY if idx % 2 == 0 else WHITE)
            txt_color = WHITE if r_row["ISO_A3"] == iso else DARK
            _add_rect(slide, _in(9.1), _in(2.45 + idx * 0.55), _in(4.0), _in(0.5), bg)
            label_txt = f"{r_row['Country'][:18]}  {r_row['PVOUT_avg']:.2f}"
            _add_text(slide, label_txt, _in(9.2), _in(2.52 + idx * 0.55),
                      _in(3.8), _in(0.38), font_size=9, color=txt_color)

    return slide


def _slide_economics(prs, site, iso, ci, system_kWp=100.0,
                     cost_per_kWp=700.0):
    """Slide 5: Economics snapshot."""
    slide = _blank_slide(prs)
    _slide_title(slide, "Economic Snapshot", f"{system_kWp:.0f} kWp System")

    # Default tariff
    tariff = 0.10
    if iso and not ci.empty:
        row = ci[ci["ISO_A3"] == iso]
        if not row.empty:
            t = row.iloc[0]["Electricity_tariff_USD_cent"]
            if pd.notna(t) and float(t) > 0:
                tariff = float(t) / 100

    result = compute_annual_yield(site, system_kWp)
    econ = compute_economics(system_kWp, result["annual_kWh"],
                              cost_per_kWp, tariff, iso)

    payback_str = f"{econ['payback_years']:.1f} yrs" if econ.get("payback_years") else "N/A"
    kpis = [
        ("Total System Cost",  f"${econ['total_cost_usd']:,.0f}", f"{system_kWp:.0f} kWp @ ${cost_per_kWp}/kWp",
         RGBColor(0xE7, 0x4C, 0x3C)),
        ("Annual Savings",     f"${econ['annual_savings_usd']:,.0f}", f"@ ${tariff:.3f}/kWh tariff",
         GREEN),
        ("Simple Payback",     payback_str, "without incentives", BRAND_ORANGE),
        ("Project LCOE",       f"${econ['lcoe_usd']:.4f}/kWh" if econ.get("lcoe_usd") else "N/A",
         "25-yr lifetime", BLUE),
    ]
    for i, (label, val, unit, color) in enumerate(kpis):
        _kpi_box(slide, label, val, unit,
                 left=_in(0.3 + i * 3.25), top=_in(0.9),
                 width=_in(3.1), height=_in(1.1), accent=color)

    # Payback sensitivity
    fig1 = _payback_sensitivity(system_kWp, cost_per_kWp, result["annual_kWh"], tariff)
    _add_image(slide, fig1, _in(0.3), _in(2.2), _in(6.5), _in(3.5))

    # CO2 summary
    co2_factor = _CO2_FACTORS.get(iso or "", _CO2_FACTORS["default"])
    annual_co2 = result["annual_kWh"] * co2_factor
    lifetime_co2 = annual_co2 * 25 / 1000
    equiv_cars = annual_co2 / 1000 / 2.4

    _add_rect(slide, _in(7.1), _in(2.2), _in(5.9), _in(4.9), LIGHT_GREY)
    _add_text(slide, "Environmental Impact", _in(7.3), _in(2.35), _in(5.5), _in(0.35),
              font_size=12, bold=True, color=DARK)
    co2_items = [
        (f"{annual_co2/1000:.1f} tonnes CO₂", "avoided per year"),
        (f"{lifetime_co2:.0f} tonnes CO₂", "offset over 25-year lifetime"),
        (f"{equiv_cars:.1f} cars", "equivalent taken off road / year"),
        (f"{co2_factor:.2f} kgCO₂/kWh", "grid emission factor applied"),
    ]
    for i, (val, lbl) in enumerate(co2_items):
        _add_rect(slide, _in(7.3), _in(2.85 + i * 1.0), _in(5.5), _in(0.85), WHITE,
                  line_color=RGBColor(0xDD, 0xDD, 0xDD))
        _add_rect(slide, _in(7.3), _in(2.85 + i * 1.0), _in(0.05), _in(0.85),
                  fill_color=GREEN)
        _add_text(slide, val, _in(7.5), _in(2.9 + i * 1.0), _in(5.1), _in(0.35),
                  font_size=13, bold=True, color=DARK)
        _add_text(slide, lbl, _in(7.5), _in(3.2 + i * 1.0), _in(5.1), _in(0.25),
                  font_size=8, color=MID_GREY)

    return slide


def _slide_system_sizing(prs, site, daily_demand=50.0):
    """Slide 6: System sizing table."""
    slide = _blank_slide(prs)
    _slide_title(slide, "System Sizing Reference", f"Based on {daily_demand:.0f} kWh/day Demand")

    pvout_annual = site["pvout"].mean()
    fractions = [25, 50, 75, 100, 125]
    sizes = [daily_demand * f / 100 / pvout_annual for f in fractions]
    areas = [s * _PANEL_AREA_PER_KWP for s in sizes]
    annual_yields = [s * pvout_annual * 365 for s in sizes]

    # Table
    headers = ["Coverage (%)", "System (kWp)", "Roof Area (m²)", "Annual Yield (kWh)"]
    col_starts = [_in(1.5), _in(4.2), _in(7.0), _in(9.8)]
    col_widths  = [_in(2.5), _in(2.5), _in(2.5), _in(2.5)]
    row_height  = _in(0.55)
    y_start     = _in(1.3)

    for c_idx, (hdr, cs, cw) in enumerate(zip(headers, col_starts, col_widths)):
        _add_rect(slide, cs, y_start, cw, row_height, BRAND_ORANGE)
        _add_text(slide, hdr, cs + _in(0.05), y_start + _in(0.12),
                  cw - _in(0.1), row_height - _in(0.1), font_size=10, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    rows_data = [
        [f"{f}%", f"{s:.1f}", f"{a:,.0f}", f"{y:,.0f}"]
        for f, s, a, y in zip(fractions, sizes, areas, annual_yields)
    ]
    for r_idx, row in enumerate(rows_data):
        bg = BRAND_ORANGE if fractions[r_idx] == 100 else (LIGHT_GREY if r_idx % 2 == 0 else WHITE)
        txt_c = WHITE if fractions[r_idx] == 100 else DARK
        for c_idx, (val, cs, cw) in enumerate(zip(row, col_starts, col_widths)):
            _add_rect(slide, cs, y_start + row_height * (r_idx + 1), cw, row_height, bg)
            _add_text(slide, val, cs + _in(0.05),
                      y_start + row_height * (r_idx + 1) + _in(0.12),
                      cw - _in(0.1), row_height - _in(0.1),
                      font_size=11, color=txt_c, align=PP_ALIGN.CENTER)

    # Key note
    note = (f"Site PVOUT: {pvout_annual:.2f} kWh/kWp/day  |  "
            f"Panel area: {_PANEL_AREA_PER_KWP} m²/kWp  |  "
            f"Source: EPW hourly data")
    _add_text(slide, note, _in(0.5), _in(6.9), _in(12.0), _in(0.4),
              font_size=8, color=MID_GREY, align=PP_ALIGN.CENTER)

    return slide


# ─── Public API ───────────────────────────────────────────────────────────────

def generate_solar_pv_pptx_report(
    epw_df: pd.DataFrame,
    metadata: dict,
    output_path: str = None,
    system_kWp: float = 100.0,
    cost_per_kWp_usd: float = 700.0,
    daily_demand_kWh: float = 50.0,
) -> io.BytesIO:
    """Generate a Solar PV analysis PPTX report.

    Parameters
    ----------
    epw_df : pd.DataFrame   — parsed EPW hourly data
    metadata : dict         — from epw_parser
    output_path : str       — if given, also writes to this path
    system_kWp : float      — default system size for yield/economics slides
    cost_per_kWp_usd : float — installed cost for economics slide
    daily_demand_kWh : float — demand for sizing slide

    Returns
    -------
    io.BytesIO  with PPTX content (seeked to 0)
    """
    monthly_sol, ci = _load_solargis_data()
    iso = _resolve_iso(metadata)
    pr = _pr_from_ghi(
        epw_df.groupby(epw_df["datetime"].dt.month)["global_horizontal_irradiance"]
        .mean().mean() / 1000 * 12
        if "global_horizontal_irradiance" in epw_df.columns else 4.0
    )
    site = compute_site_solar(epw_df, pr=pr)

    city = metadata.get("city") or metadata.get("location") or "Unknown"
    lat = metadata.get("latitude")
    lon = metadata.get("longitude")
    pvout_annual = site["pvout"].mean()

    prs = _load_template()

    _slide_title_card(prs, city, lat, lon, pvout_annual, site)
    _slide_solar_resource(prs, site, epw_df)
    _slide_yield(prs, site, system_kWp)
    _slide_benchmark(prs, site, iso, monthly_sol, ci)
    _slide_economics(prs, site, iso, ci, system_kWp, cost_per_kWp_usd)
    _slide_system_sizing(prs, site, daily_demand_kWh)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    if output_path:
        with open(output_path, "wb") as f:
            f.write(buf.read())
        buf.seek(0)

    return buf
