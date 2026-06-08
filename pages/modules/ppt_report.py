"""PowerPoint report generation for Climate Analytics Dashboard."""

import io
import os
import tempfile

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pytz

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modules.shading_helpers import (
    _ORIENTATIONS,
    build_thermal_matrix,
    compute_shading_geometry,
    build_orientation_table,
    compute_solar_angles,
    get_overheating_hours,
)


def _ppt_remove_all_slides(prs):
    """Remove every slide from an open Presentation while preserving theme/layouts."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)


def generate_pptx_report(
    df: pd.DataFrame,
    start_date,
    end_date,
    start_hour: int,
    end_hour: int,
    selected_parameter: str,
    metadata: dict = None,
):
    """Generate a PowerPoint report with Dry Bulb, Relative Humidity and Sun Path sections."""

    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]

    TITLE_RED   = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER_CLR = RGBColor(0xC0, 0x00, 0x00)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches

    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(LOGO_L), Inches(LOGO_T),
                width=Inches(LOGO_W), height=Inches(LOGO_H),
            )

    def _add_slide_title(slide, text, left=0.27, top=0.13, width=None, height=0.45):
        if width is None:
            width = SW - left - 0.3
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _add_divider(slide, top_inches):
        line = slide.shapes.add_shape(
            1, Inches(0.27), Inches(top_inches),
            Inches(SW - 0.54), Inches(0.03),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = DIVIDER_CLR
        line.line.fill.background()

    def _save_mpl_figure(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches='tight', facecolor='white')
            return tmp.name

    def _err_box(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Visualization error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    filtered_df = df[
        (df["datetime"].dt.date >= start_date) &
        (df["datetime"].dt.date <= end_date) &
        (df["hour"].between(start_hour, end_hour))
    ]

    # ── COVER SLIDE ───────────────────────────────────────────────────────────
    def _make_cover_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)

        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(SW), Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(SW - 1.2), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Climate Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(SW - 1.2), Inches(0.7))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        _location = metadata.get("location", "") if metadata else ""
        location_display = _city if _city else (_location if _location else "Location")
        run2.text = f"{location_display}"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Dry Bulb Temperature  |  Relative Humidity  |  Sun Path"
        run3.font.size = Pt(11)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _make_cover_slide()

    # ── SECTION 1 – DRY BULB TEMPERATURE ─────────────────────────────────────
    def _make_dbt_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                temp_min=("dry_bulb_temperature", "min"),
                temp_max=("dry_bulb_temperature", "max"),
                temp_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
            comfort_line = daily_avg.rolling(window=7, center=True).mean()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=130)
            ax.fill_between(daily_stats["doy"], comfort_line - 3.5, comfort_line + 3.5,
                            alpha=0.18, color='gray', label='ASHRAE 80% Comfort')
            ax.fill_between(daily_stats["doy"], comfort_line - 2.5, comfort_line + 2.5,
                            alpha=0.28, color='gray', label='ASHRAE 90% Comfort')
            ax.fill_between(daily_stats["doy"], daily_stats["temp_min"], daily_stats["temp_max"],
                            alpha=0.30, color='#FFB3B3', label='Daily Temp Range')
            ax.plot(daily_stats["doy"], daily_stats["temp_avg"],
                    color='#C00000', linewidth=2.2, label='Daily Average', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Annual Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=5, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_trend_slide()

    def _make_dbt_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Dry Bulb Temperature – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly = df.groupby("month").agg(
                t_min=("dry_bulb_temperature", "min"),
                t_max=("dry_bulb_temperature", "max"),
                t_avg=("dry_bulb_temperature", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=130)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly["t_min"], bar_w, color='#90CAF9', label='Min Temp')
            ax.bar(x,          monthly["t_avg"], bar_w, color='#C00000', label='Avg Temp', alpha=0.85)
            ax.bar(x + bar_w,  monthly["t_max"], bar_w, color='#EF9A9A', label='Max Temp')

            ax.axhspan(20, 26, alpha=0.10, color='green', label='Comfort Band (20–26°C)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Monthly Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min: {filtered_df['dry_bulb_temperature'].min():.1f}°C  "
                    f"Avg: {filtered_df['dry_bulb_temperature'].mean():.1f}°C  "
                    f"Max: {filtered_df['dry_bulb_temperature'].max():.1f}°C  "
                    f" |  Ann. CDD24: {(df['dry_bulb_temperature'] - 24).clip(lower=0).sum():.0f}   "
                    f"HDD18: {(18 - df['dry_bulb_temperature']).clip(lower=0).sum():.0f}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_dbt_monthly_slide()

    # ── SECTION 2 – RELATIVE HUMIDITY ─────────────────────────────────────────
    def _make_rh_trend_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity")
        _add_divider(slide, 0.62)

        try:
            daily_stats = df.groupby("doy").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            start_doy = pd.to_datetime(f"2024-{start_date.month:02d}-01").dayofyear
            end_doy = (
                366 if end_date.month == 12
                else pd.to_datetime(f"2024-{end_date.month+1:02d}-01").dayofyear - 1
            )

            fig, ax = plt.subplots(figsize=(13, 5.4), dpi=130)
            ax.axhspan(75, 100, alpha=0.13, color='#FF6B6B', label='Condensation Risk (>75%)')
            ax.axhspan(60,  75, alpha=0.13, color='#FFA500', label='High RH (60–75%)')
            ax.axhspan(30,  60, alpha=0.13, color='#4ECDC4', label='Comfortable (30–60%)')
            ax.axhspan( 0,  30, alpha=0.13, color='#FFD93D', label='Low RH (<30%)')

            ax.fill_between(daily_stats["doy"], daily_stats["rh_min"], daily_stats["rh_max"],
                            alpha=0.28, color='#0099ff', label='Daily RH Range')
            ax.plot(daily_stats["doy"], daily_stats["rh_avg"],
                    color='#0066cc', linewidth=2.2, label='Daily Average RH', zorder=3)
            ax.axvspan(start_doy, end_doy, alpha=0.07, color='#2c5aa0', label='Selected Period')

            months_doy = [1, 32, 60, 91, 121, 152, 182, 213, 244, 274, 305, 335]
            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            ax.set_xticks(months_doy)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 100)
            ax.set_title('Annual Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.10), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            plt.tight_layout()

            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.8))
            os.unlink(tmp)
        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_trend_slide()

    def _make_rh_monthly_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Relative Humidity – Monthly Summary")
        _add_divider(slide, 0.62)

        try:
            monthly_rh = df.groupby("month").agg(
                rh_min=("relative_humidity", "min"),
                rh_max=("relative_humidity", "max"),
                rh_avg=("relative_humidity", "mean"),
            ).reset_index()

            months_lbl = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
            x = np.arange(12)

            fig, ax = plt.subplots(figsize=(13, 5.0), dpi=130)
            bar_w = 0.30
            ax.bar(x - bar_w, monthly_rh["rh_min"], bar_w, color='#AED6F1', label='Min RH')
            ax.bar(x,          monthly_rh["rh_avg"], bar_w, color='#0066cc', label='Avg RH', alpha=0.85)
            ax.bar(x + bar_w,  monthly_rh["rh_max"], bar_w, color='#5DADE2', label='Max RH')

            ax.axhspan(30, 60, alpha=0.10, color='green', label='Comfortable (30–60%)')
            ax.axhline(75, color='#E74C3C', linewidth=1.2, linestyle='--', label='Condensation Threshold (75%)')

            ax.set_xticks(x)
            ax.set_xticklabels(months_lbl, fontsize=10)
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_ylim(0, 110)
            ax.set_title('Monthly Relative Humidity Trend', fontsize=13, fontweight='bold', pad=10, color='#333')
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.09), ncol=4, frameon=True, fontsize=9)
            ax.grid(True, alpha=0.25, linestyle='--', axis='y')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')

            for m in range(start_date.month, end_date.month + 1):
                ax.axvspan(m - 1 - 0.5, m - 1 + 0.5, alpha=0.06, color='navy')

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72), width=Inches(SW - 0.54), height=Inches(5.5))
            os.unlink(tmp)

            if not filtered_df.empty:
                stats_txt = (
                    f"Selected Period   Min RH: {filtered_df['relative_humidity'].min():.0f}%  "
                    f"Avg RH: {filtered_df['relative_humidity'].mean():.0f}%  "
                    f"Max RH: {filtered_df['relative_humidity'].max():.0f}%  "
                    f" |  High RH hrs (>60%): {len(filtered_df[filtered_df['relative_humidity'] > 60])}  "
                    f"Condensation risk hrs (>75%): {len(filtered_df[filtered_df['relative_humidity'] > 75])}"
                )
                tb = slide.shapes.add_textbox(Inches(0.27), Inches(6.35), Inches(SW - 0.54), Inches(0.35))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = stats_txt
                run.font.size = Pt(9)
                run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_rh_monthly_slide()

    # ── SECTION 3 – SUN PATH ──────────────────────────────────────────────────
    def _make_sun_path_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Sun Path Diagram")
        _add_divider(slide, 0.62)

        _meta = metadata or {}
        lat = _meta.get("latitude")
        lon = _meta.get("longitude")
        tz_str = _meta.get("timezone", "UTC")

        if lat is None or lon is None:
            _err_box(slide, "Latitude/Longitude not available from EPW metadata.")
            _add_logo(slide)
            return

        try:
            from pvlib import solarposition as _solpos_lib

            try:
                _tz = pytz.timezone(tz_str)
            except Exception:
                _tz = pytz.UTC

            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left")
            sol = _solpos_lib.get_solarposition(times, lat, lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()
            sol["r"] = 90 - sol["apparent_elevation"]

            fig = plt.figure(figsize=(7.5, 7.5), dpi=130, facecolor='white')
            ax = fig.add_subplot(111, projection='polar')
            ax.set_theta_zero_location('N')
            ax.set_theta_direction(-1)
            ax.set_aspect('equal', adjustable='box')
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(['90°\n(Zenith)', '75°', '60°', '45°', '30°', '15°', '0°\n(Horizon)'],
                               fontsize=7, color='#555')
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'], fontsize=10, fontweight='bold')
            ax.set_facecolor('#F0F4F8')
            ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)

            sc = ax.scatter(
                np.radians(sol["azimuth"].values),
                sol["r"].values,
                c=sol.index.dayofyear,
                cmap='YlOrRd',
                s=1.0, alpha=0.55,
                vmin=1, vmax=365,
                linewidths=0, zorder=2,
            )
            cbar = fig.colorbar(sc, ax=ax, pad=0.10, fraction=0.035, shrink=0.75)
            cbar.set_label('Day of Year', fontsize=9)
            cbar.set_ticks([1, 91, 182, 273, 365])
            cbar.set_ticklabels(['1\n(Jan)', '91\n(Apr)', '182\n(Jul)', '273\n(Oct)', '365\n(Dec)'])

            key_dates = [
                ("Mar 21 (Spring Equinox)", "2020-03-21", "#FF9500", 1.6),
                ("Jun 21 (Summer Solstice)", "2020-06-21", "#CC0000", 2.0),
                ("Dec 21 (Winter Solstice)", "2020-12-21", "#0066CC", 2.0),
            ]
            for lbl, dstr, col, lw in key_dates:
                dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
                ks = _solpos_lib.get_solarposition(dt, lat, lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.06), ncol=3,
                      frameon=True, fontsize=8, borderaxespad=0)
            ax.set_title(f'Sun Path  |  Lat: {lat:.2f}°  Lon: {lon:.2f}°',
                         fontsize=11, fontweight='bold', color='#333', pad=14)

            plt.tight_layout()
            tmp = _save_mpl_figure(fig)
            plt.close(fig)

            # Use square dimensions to maintain circular aspect ratio
            img_size = min(SW * 0.55, SH * 0.75)
            img_l = (SW - img_size) / 2
            img_t = 0.72
            slide.shapes.add_picture(tmp, Inches(img_l), Inches(img_t), width=Inches(img_size), height=Inches(img_size))
            os.unlink(tmp)

        except Exception as e:
            _err_box(slide, e)

        _add_logo(slide)

    _make_sun_path_slide()

    # ── SHADING ANALYSIS SLIDE ────────────────────────────────────────────────
    def _plot_sun_path_shading(lat, lon, tz_str):
        """Generate a sun path diagram with horizontal overhang shading profile."""
        try:
            from pvlib import solarposition as _sp

            try:
                _tz = pytz.timezone(tz_str)
            except Exception:
                _tz = pytz.UTC

            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=_tz, inclusive="left")
            sol = _sp.get_solarposition(times, lat, lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()
            sol["r"] = 90 - sol["apparent_elevation"]

            # fig = plt.figure(figsize=(8.0, 7.0), dpi=130, facecolor='white')
            fig = plt.figure(figsize=(7.0, 7.0), dpi=130, facecolor='white')
            ax = fig.add_subplot(111, projection='polar')
            ax.set_theta_zero_location('N')
            ax.set_theta_direction(-1)
            ax.set_aspect('equal', adjustable='box')
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(['90°', '75°', '60°', '45°', '30°', '15°', '0°'], fontsize=7, color='#555')
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(['N', 'NE', 'E', 'SE', 'S', 'SW', 'W', 'NW'], fontsize=9, fontweight='bold')
            ax.set_facecolor('#F0F4F8')
            ax.grid(True, alpha=0.35, linestyle='--', linewidth=0.6)

            sc = ax.scatter(
                np.radians(sol["azimuth"].values),
                sol["r"].values,
                c=sol.index.dayofyear,
                cmap='YlOrRd',
                s=0.8, alpha=0.6,
                vmin=1, vmax=365,
                linewidths=0, zorder=2,
            )
            cbar = fig.colorbar(sc, ax=ax, pad=0.10, fraction=0.035, shrink=0.7)
            cbar.set_label('Day of Year', fontsize=8)
            cbar.set_ticks([1, 91, 182, 273, 365])
            cbar.set_ticklabels(['Jan', 'Apr', 'Jul', 'Oct', 'Dec'], fontsize=7)

            key_dates = [
                ("Equinox", "2020-03-21", "#FF9500", 1.5),
                ("Summer",  "2020-06-21", "#CC0000", 1.8),
                ("Winter",  "2020-12-21", "#0066CC", 1.8),
            ]
            for lbl, dstr, col, lw in key_dates:
                dt = pd.date_range(dstr, periods=288, freq='5min', tz=_tz)
                ks = _sp.get_solarposition(dt, lat, lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            overhang_altitude = 35
            theta_range = np.radians(np.linspace(45, 315, 100))
            shading_altitude = np.ones_like(theta_range) * overhang_altitude

            ax.fill_between(theta_range, shading_altitude, 90, alpha=0.12, color='#8B4513',
                            label='Shading Zone', zorder=1)
            ax.plot(theta_range, shading_altitude, color='#654321', linewidth=2.5,
                   label='Overhang Profile', linestyle='--', zorder=3)

            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.08), ncol=4,
                     frameon=True, fontsize=7, borderaxespad=0)
            ax.set_title(f'Sun Path with Shading Profile\nLat: {lat:.2f}°  Lon: {lon:.2f}°',
                        fontsize=10, fontweight='bold', color='#333', pad=12)

            # plt.tight_layout()
            plt.subplots_adjust(top=0.88, bottom=0.12)
            return fig
        except Exception as e:
            print(f"Shading diagram error: {e}")
            return None

    def _make_shading_summary_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Shading Strategy")
        _add_divider(slide, 0.62)

        diagram_added = False
        _meta = metadata or {}
        _lat = _meta.get("latitude")
        _lon = _meta.get("longitude")
        _tz_str = _meta.get("timezone", "UTC")

        if _lat is not None and _lon is not None:
            try:
                shading_fig = _plot_sun_path_shading(_lat, _lon, _tz_str)
                if shading_fig is not None:
                    tmp_shading = _save_mpl_figure(shading_fig)
                    plt.close(shading_fig)
                    diagram_width = (SW - 0.54) * 0.45
                    slide.shapes.add_picture(tmp_shading, Inches(0.27), Inches(0.75),
                                           width=Inches(diagram_width), height=Inches(5.8))
                    os.unlink(tmp_shading)
                    diagram_added = True
            except Exception as e:
                print(f"Shading diagram error: {e}")

        if diagram_added:
            diagram_width = (SW - 0.54) * 0.45
            text_left = 0.27 + diagram_width
            text_width = (SW - 0.54) * 0.55
            tb = slide.shapes.add_textbox(Inches(text_left), Inches(0.80),
                                          Inches(text_width), Inches(5.8))
        else:
            tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))

        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Solar Geometry & Shading Analysis"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(4)

        try:
            _meta = metadata or {}
            _lat = _meta.get("latitude")
            _lon = _meta.get("longitude")
            _tz_str = _meta.get("timezone", "UTC")

            if _lat is None or _lon is None:
                raise ValueError("Latitude/Longitude not available")

            from pvlib import solarposition as _solpos_lib_shade

            try:
                _tz_obj = pytz.timezone(_tz_str)
            except Exception:
                _tz_obj = pytz.UTC

            summer_dt = pd.date_range("2020-06-21", periods=24, freq="h", tz=_tz_obj)
            ssum = _solpos_lib_shade.get_solarposition(summer_dt, _lat, _lon)
            noon_alt_sum = float(ssum.iloc[12]["apparent_elevation"])

            winter_dt = pd.date_range("2020-12-21", periods=24, freq="h", tz=_tz_obj)
            swin = _solpos_lib_shade.get_solarposition(winter_dt, _lat, _lon)
            noon_alt_win = float(swin.iloc[12]["apparent_elevation"])

            p = tf.add_paragraph()
            p.text = f"• Summer Solstice (Jun 21): Solar altitude at noon = {noon_alt_sum:.1f}°"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(3)

            p = tf.add_paragraph()
            p.text = f"• Winter Solstice (Dec 21): Solar altitude at noon = {noon_alt_win:.1f}°"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(3)

            ghi_col = filtered_df.get("global_horizontal_irradiance", pd.Series(0))
            if len(ghi_col) == 0:
                ghi_col = pd.Series(0, index=filtered_df.index)
            ghi_col = ghi_col.fillna(0)

            temp_col = filtered_df.get("dry_bulb_temperature", pd.Series(20))
            if len(temp_col) == 0:
                temp_col = pd.Series(20, index=filtered_df.index)
            temp_col = temp_col.fillna(20)

            shading_needed = (temp_col > 28) & (ghi_col > 315)
            shading_hours = shading_needed.sum() / 2
            total_observation_hours = len(filtered_df) / 2

            if total_observation_hours > 0:
                shading_pct = (shading_hours / total_observation_hours) * 100
            else:
                shading_pct = 0

            p = tf.add_paragraph()
            p.text = f"• Shading required: {shading_hours:.0f} hours ({shading_pct:.1f}% of period)"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(6)
        except Exception:
            p = tf.add_paragraph()
            p.text = "• Solar altitude and shading data not available"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "Shading Recommendations"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        recommendations = [
            "• South-facing facades: Use horizontal overhangs (louvers) or shading devices to block summer sun while allowing winter sunlight penetration",
            "• East/West facades: Use vertical fins or combination of overhangs and fins to minimize morning/afternoon heat gain",
            "• North-facing facades: Minimal shading required; prioritize daylighting and views",
            "• Use high-performance glazing with low solar heat gain coefficient (SHGC) in high solar radiation areas",
            "• Consider automated shading systems for dynamic climate response throughout the year",
        ]
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.0
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "Design Considerations"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(3)

        considerations = [
            "• Optimal window-to-wall ratio: 30-40% for climate comfort; balance daylighting with thermal performance",
            "• Depth of shading device: D/H ratio (depth to height) between 0.5-1.0 for effective summer shading",
            "• Material selection: High-albedo surfaces reflect solar radiation; low-emissivity coatings minimize thermal transmission",
        ]
        for cons in considerations:
            p = tf.add_paragraph()
            p.text = cons
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.0
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_shading_summary_slide()

    # ── ANNEXURE SLIDE ────────────────────────────────────────────────────────
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _add_slide_title(slide, "Annexure")
        _add_divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = (
            "Environmental Design Solutions [EDS] is a sustainability advisory firm. "
            "Since 2002, EDS has worked on over 500 green building and energy efficiency "
            "projects worldwide. The team focuses on climate change mitigation, low-carbon "
            "design, building simulation, performance audits, and capacity building. EDS "
            "continues to contribute to the buildings community with useful tools through "
            "its IT services."
        )
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        for item in [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "\u2022  EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "\u2022  EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "\u2022  Open-source resources used: Clima - Berkley, Streamlit, Python",
            "\u2022  EDS is not liable to inform Users about updates to the Application or underlying resources",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        for item in [
            "\u2022  Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "\u2022  Streamlit, \u00a9 Streamlit Inc., licensed under Apache 2.0",
            "\u2022  Python \u00a9 Python Software Foundation, licensed under PSF License Version 2",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_annexure_slide()

    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes


# ─────────────────────────────────────────────────────────────────────────────


def generate_shading_pptx_report(
    df: pd.DataFrame,
    metadata: dict,
    temp_threshold: float = 28.0,
    rad_threshold: float = 315.0,
    lat: float = None,
    lon: float = None,
    tz_str: str = "UTC",
    design_cutoff_angle: float = 45.0,
):
    """Generate a Shading Analysis PowerPoint report using the Voha template."""

    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]
    TITLE_RED = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY  = RGBColor(0x40, 0x40, 0x40)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches
    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(LOGO_L), Inches(LOGO_T),
                                     width=Inches(LOGO_W), height=Inches(LOGO_H))

    def _slide_title(slide, text, top=0.13):
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(top), Inches(SW - 0.54), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _divider(slide, top_inches):
        bar = slide.shapes.add_shape(1, Inches(0.27), Inches(top_inches),
                                     Inches(SW - 0.54), Inches(0.03))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TITLE_RED
        bar.line.fill.background()

    def _err(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = TITLE_RED

    def _save_fig(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches="tight", facecolor="white")
            return tmp.name

    _lat  = lat  if lat  is not None else (metadata.get("latitude")  or 0.0)
    _lon  = lon  if lon  is not None else (metadata.get("longitude") or 0.0)
    _tz   = tz_str or metadata.get("timezone", "UTC")

    # ── COVER ─────────────────────────────────────────────────────────────────
    def _cover():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(SW), Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(SW - 1.2), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Shading Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.75), Inches(SW - 1.2), Inches(0.65))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        run2.text = f"Location: {_city}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Thermal & Radiation Matrix  |  Sun Path (Shading Mode)  |  Orientation Analysis  |  Shading Masks"
        run3.font.size = Pt(10)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _cover()

    # ── THERMAL & RADIATION MATRIX ────────────────────────────────────────────
    def _thermal_matrix_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Thermal & Radiation Matrix")
        _divider(slide, 0.62)

        try:
            temp_matrix, rad_matrix, overheat_mask = build_thermal_matrix(df, temp_threshold, rad_threshold)
            months_lbl = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
            hours_lbl  = [f"{h:02d}:00" for h in range(24)]

            fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.2), dpi=120)

            for ax, matrix, title, cmap, clabel in [
                (axes[0], temp_matrix, f"Mean Dry-Bulb Temp (\u00b0C)  [threshold: {temp_threshold}\u00b0C]", "RdYlBu_r", "\u00b0C"),
                (axes[1], rad_matrix,  f"Mean GHI (W/m\u00b2)  [threshold: {rad_threshold} W/m\u00b2]",     "YlOrRd",   "W/m\u00b2"),
            ]:
                im = ax.imshow(matrix.values, aspect="auto", origin="upper", cmap=cmap)
                plt.colorbar(im, ax=ax, fraction=0.035, pad=0.03, label=clabel)
                ax.set_xticks(range(12))
                ax.set_xticklabels(months_lbl, fontsize=8)
                ax.set_yticks(range(24))
                ax.set_yticklabels(hours_lbl, fontsize=7)
                ax.set_title(title, fontsize=10, fontweight="bold", pad=8)
                ax.set_xlabel("Month", fontsize=9)
                ax.set_ylabel("Hour of Day", fontsize=9)

                for h_i in range(24):
                    for m_i in range(12):
                        if overheat_mask.iloc[h_i, m_i]:
                            rect = plt.Rectangle(
                                (m_i - 0.5, h_i - 0.5), 1, 1,
                                fill=False, edgecolor="black", linewidth=1.6,
                            )
                            ax.add_patch(rect)

            fig.suptitle(
                "Overheating Hours  (black border = both thresholds exceeded)",
                fontsize=11, fontweight="bold", y=1.01, color="#333",
            )
            fig.patch.set_facecolor("white")
            plt.tight_layout()

            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)
        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _thermal_matrix_slide()

    # ── SUN PATH (SHADING MODE) ───────────────────────────────────────────────
    def _sun_path_shading_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Sun Path – Shading Analysis")
        _divider(slide, 0.62)

        try:
            from pvlib import solarposition as _sp

            try:
                tz = pytz.timezone(_tz)
            except Exception:
                tz = pytz.UTC

            times = pd.date_range("2020-01-01", "2021-01-01", freq="h", tz=tz, inclusive="left")
            sol = _sp.get_solarposition(times, _lat, _lon)
            sol = sol[sol["apparent_elevation"] > 0].copy()

            # Robust EPW handling: ensure datetime, map common aliases, and resample half-hour to hourly
            _df = df.copy()
            if "datetime" not in _df.columns:
                if isinstance(_df.index, pd.DatetimeIndex):
                    _df = _df.reset_index().rename(columns={_df.index.name or "index": "datetime"})
                else:
                    raise ValueError("EPW missing 'datetime' column")

            common_aliases = {
                'dni': 'direct_normal_irradiance',
                'direct normal': 'direct_normal_irradiance',
                'direct_normal': 'direct_normal_irradiance',
                'dhi': 'diffuse_horizontal_irradiance',
                'diffuse horizontal': 'diffuse_horizontal_irradiance',
                'ghi': 'global_horizontal_irradiance',
                'global horizontal': 'global_horizontal_irradiance',
                'dry bulb': 'dry_bulb_temperature',
                'dry_bulb': 'dry_bulb_temperature',
                'drybulb': 'dry_bulb_temperature',
                'temperature': 'dry_bulb_temperature',
                'temp': 'dry_bulb_temperature',
            }

            cols_lower = {c: c.lower() for c in _df.columns}
            rename_map = {}
            for alias, target in common_aliases.items():
                if target in _df.columns:
                    continue
                for orig, low in cols_lower.items():
                    if alias == low or alias in low:
                        if orig != 'datetime' and target not in _df.columns:
                            rename_map[orig] = target
                            break

            if rename_map:
                _df = _df.rename(columns=rename_map)

            epw = _df.set_index("datetime").copy()
            if epw.index.tz is None:
                epw.index = epw.index.tz_localize(tz)
            else:
                epw.index = epw.index.tz_convert(tz)
            epw.index = epw.index.map(lambda x: x.replace(year=2020))

            # Resample half-hour EPW to hourly for consistent join
            try:
                has_half_hour = any(t.minute != 0 for t in epw.index)
            except Exception:
                has_half_hour = False

            epw_hourly = epw
            if has_half_hour:
                candidate_cols = [c for c in [
                    "dry_bulb_temperature", "direct_normal_irradiance",
                    "diffuse_horizontal_irradiance", "global_horizontal_irradiance"
                ] if c in epw.columns]
                if candidate_cols:
                    epw_num = epw[candidate_cols].apply(pd.to_numeric, errors='coerce')
                    epw_hourly = epw_num.resample('h').mean().dropna(how='all')
                else:
                    epw_hourly = epw.resample('h').first().dropna(how='all')

            sol = sol.join(epw_hourly[["dry_bulb_temperature", "global_horizontal_irradiance"]], how="left")
            sol["global_horizontal_irradiance"] = sol["global_horizontal_irradiance"].fillna(0)
            sol["dry_bulb_temperature"] = sol["dry_bulb_temperature"].fillna(
                sol["dry_bulb_temperature"].median()
            )
            shading_needed = (
                (sol["dry_bulb_temperature"] > temp_threshold) &
                (sol["global_horizontal_irradiance"] > rad_threshold)
            )

            fig = plt.figure(figsize=(7.5, 7.5), dpi=130, facecolor="white")
            ax = fig.add_subplot(111, projection="polar")
            ax.set_theta_zero_location("N")
            ax.set_theta_direction(-1)
            ax.set_aspect('equal', adjustable='box')
            ax.set_ylim(0, 90)
            ax.set_yticks([0, 15, 30, 45, 60, 75, 90])
            ax.set_yticklabels(["90\u00b0","75\u00b0","60\u00b0","45\u00b0","30\u00b0","15\u00b0","0\u00b0"],
                               fontsize=7, color="#555")
            ax.set_xticks(np.radians([0, 45, 90, 135, 180, 225, 270, 315]))
            ax.set_xticklabels(["N","NE","E","SE","S","SW","W","NW"], fontsize=10, fontweight="bold")
            ax.set_facecolor("#F0F4F8")
            ax.grid(True, alpha=0.35, linestyle="--", linewidth=0.6)

            r = 90 - sol["apparent_elevation"].values
            theta = np.radians(sol["azimuth"].values)

            mask_ok = ~shading_needed.values
            ax.scatter(theta[mask_ok], r[mask_ok], c="#FFF9C4", s=1.2,
                       alpha=0.45, linewidths=0, label="No shading needed", zorder=2)

            mask_sh = shading_needed.values
            ax.scatter(theta[mask_sh], r[mask_sh], c="#E65100", s=2.5,
                       alpha=0.75, linewidths=0, label="Shading required", zorder=3)

            for lbl, dstr, col, lw in [
                ("Mar 21", "2020-03-21", "#FF9500", 1.4),
                ("Jun 21", "2020-06-21", "#CC0000", 1.8),
                ("Dec 21", "2020-12-21", "#0066CC", 1.8),
            ]:
                dt = pd.date_range(dstr, periods=288, freq="5min", tz=tz)
                ks = _sp.get_solarposition(dt, _lat, _lon)
                ks = ks[ks["apparent_elevation"] > 0]
                if not ks.empty:
                    ax.plot(np.radians(ks["azimuth"]), 90 - ks["apparent_elevation"],
                            color=col, linewidth=lw, label=lbl, zorder=4)

            shading_pct = mask_sh.sum() / len(mask_sh) * 100 if len(mask_sh) else 0
            ax.set_title(
                f"Sun Path – Shading Mode   ({shading_pct:.1f}% of daytime hours require shading)",
                fontsize=10, fontweight="bold", color="#333", pad=14,
            )
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.06), ncol=3,
                      frameon=True, fontsize=8)

            plt.tight_layout()
            tmp = _save_fig(fig)
            plt.close(fig)

            # Use square dimensions to maintain circular aspect ratio
            img_size = min(SW * 0.50, SH * 0.75)
            img_l = (SW - img_size) / 2
            slide.shapes.add_picture(tmp, Inches(img_l), Inches(0.72),
                                     width=Inches(img_size), height=Inches(img_size))
            os.unlink(tmp)

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _sun_path_shading_slide()

    # ── ORIENTATION SHADING ANALYSIS TABLE ───────────────────────────────────
    def _orientation_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, f"Orientation Shading Analysis  (Design cutoff: {design_cutoff_angle}\u00b0)")
        _divider(slide, 0.62)

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz)
            if solar_pos.empty:
                _err(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_df = build_orientation_table(solar_pos, design_cutoff_angle)

            fig, ax = plt.subplots(figsize=(13, 5.8), dpi=120)
            ax.axis("off")

            col_labels = ["Orientation", "Rays\nHitting", "Min VSA", "Max |HSA|",
                          "D/H\nOverhang", "D/W\nFin", "Protection %"]
            table_data = []
            row_colors = []
            for _, row in orient_df.iterrows():
                pct = row["Protection (%)"]
                if pct is None:
                    c = "#f5f5f5"
                elif pct >= 95:
                    c = "#e8f5e9"
                elif pct >= 80:
                    c = "#fff3e0"
                else:
                    c = "#ffebee"
                row_colors.append([c] * 7)

                dh  = f"{row['D/H (Overhang)']:.3f}" if row["D/H (Overhang)"] is not None else "—"
                dw  = f"{row['D/W (Fin)']:.3f}"       if row["D/W (Fin)"] is not None else "—"
                vsa = f"{row['Min VSA (°)']:.1f}\u00b0"  if row["Min VSA (°)"] is not None else "—"
                hsa = f"{row['Max |HSA| (°)']:.1f}\u00b0" if row["Max |HSA| (°)"] is not None else "—"
                pct_s = f"{pct:.1f}%" if pct is not None else "—"
                table_data.append([
                    row["Orientation"], str(row["Rays Hitting"]),
                    vsa, hsa, dh, dw, pct_s,
                ])

            tbl = ax.table(cellText=table_data, colLabels=col_labels,
                           cellColours=row_colors, loc="center", cellLoc="center")
            tbl.auto_set_font_size(False)
            tbl.set_fontsize(11)
            tbl.scale(1, 2.1)

            for j in range(len(col_labels)):
                cell = tbl[0, j]
                cell.set_facecolor("#1a3a52")
                cell.set_text_props(color="white", fontweight="bold")

            for i in range(1, len(table_data) + 1):
                tbl[i, 0].get_text().set_ha("left")

            fig.patch.set_facecolor("white")
            ax.set_title(
                f"{len(solar_pos)} overheating daytime sun positions  |  "
                f"Temp > {temp_threshold}\u00b0C  &  GHI > {rad_threshold} W/m\u00b2",
                fontsize=10, color="#555", pad=10,
            )

            tmp = _save_fig(fig)
            plt.close(fig)
            slide.shapes.add_picture(tmp, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(tmp)

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _orientation_slide()

    # ── SHADING MASK DIAGRAMS (2×4 grid) ─────────────────────────────────────
    def _shading_masks_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Shading Mask Diagrams")
        _divider(slide, 0.62)

        try:
            overheat_df = get_overheating_hours(df, temp_threshold, rad_threshold)
            if overheat_df.empty:
                _err(slide, "No overheating hours found with current thresholds.")
                _add_logo(slide)
                return

            solar_pos = compute_solar_angles(overheat_df, _lat, _lon, _tz)
            if solar_pos.empty:
                _err(slide, "No daytime overheating sun positions found.")
                _add_logo(slide)
                return

            orient_items = list(_ORIENTATIONS.items())
            n_cols = 4
            cell_w = (SW - 0.54) / n_cols
            cell_h = (SH - 1.05) / 2
            top_start = 0.75

            for idx, (oname, faz) in enumerate(orient_items):
                col_i = idx % n_cols
                row_i = idx // n_cols
                cell_l = 0.27 + col_i * cell_w
                cell_t = top_start + row_i * cell_h

                geom   = compute_shading_geometry(solar_pos, faz)
                facing = geom[geom["hits_facade"]]
                other  = geom[~geom["hits_facade"]]

                fig = plt.figure(figsize=(3.0, 2.8), dpi=110, facecolor="white")
                ax = fig.add_subplot(111, projection="polar")
                ax.set_theta_zero_location("N")
                ax.set_theta_direction(-1)
                ax.set_ylim(0, 90)
                ax.set_yticks([])
                ax.set_xticks(np.radians([0, 90, 180, 270]))
                ax.set_xticklabels(["N", "E", "S", "W"], fontsize=8, fontweight="bold")
                ax.set_facecolor("#f0f8ff")
                ax.grid(True, alpha=0.30, linewidth=0.5)

                if not other.empty:
                    ax.scatter(
                        np.radians(other["solar_azimuth"]), 90 - other["solar_altitude"],
                        s=2, c="lightgrey", alpha=0.5, linewidths=0, zorder=2,
                    )
                if not facing.empty:
                    ax.scatter(
                        np.radians(facing["solar_azimuth"]), 90 - facing["solar_altitude"],
                        s=4, c="#E65100", alpha=0.75, linewidths=0, zorder=3,
                    )

                rel_az_r = np.linspace(-89, 89, 179)
                tan_co = np.tan(np.radians(design_cutoff_angle))
                co_alt = np.degrees(np.arctan(tan_co * np.cos(np.radians(rel_az_r))))
                co_az  = faz + rel_az_r
                valid  = co_alt > 0
                if valid.any():
                    ax.plot(np.radians(co_az[valid]), 90 - co_alt[valid],
                            color="#1565C0", linewidth=1.4, linestyle="--", zorder=4)

                ax.plot(np.radians([faz, faz]), [0, 85],
                        color="#388E3C", linewidth=1.4, zorder=5)

                ax.set_title(oname, fontsize=7, fontweight="bold", pad=4, color="#222")
                plt.tight_layout(pad=0.3)

                tmp = _save_fig(fig)
                plt.close(fig)
                slide.shapes.add_picture(
                    tmp, Inches(cell_l), Inches(cell_t),
                    width=Inches(cell_w - 0.05), height=Inches(cell_h - 0.05),
                )
                os.unlink(tmp)

            leg_tb = slide.shapes.add_textbox(
                Inches(0.27), Inches(SH - LOGO_H - 0.55), Inches(SW - 0.54), Inches(0.35),
            )
            leg_tf = leg_tb.text_frame
            leg_p = leg_tf.paragraphs[0]
            leg_run = leg_p.add_run()
            leg_run.text = (
                "\u25cf Overheating rays (hits facade)  "
                "\u25cf Overheating (other side)  "
                "- - Cutoff arc (VSA cut-off)  "
                "\u2014 Facade direction"
            )
            leg_run.font.size = Pt(8)
            leg_run.font.color.rgb = DARK_GREY

        except Exception as e:
            _err(slide, e)

        _add_logo(slide)

    _shading_masks_slide()

    # ── ANNEXURE ──────────────────────────────────────────────────────────────
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Annexure")
        _divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = (
            "Environmental Design Solutions [EDS] is a sustainability advisory firm. "
            "Since 2002, EDS has worked on over 500 green building and energy efficiency "
            "projects worldwide. The team focuses on climate change mitigation, low-carbon "
            "design, building simulation, performance audits, and capacity building. EDS "
            "continues to contribute to the buildings community with useful tools through "
            "its IT services."
        )
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        for item in [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "\u2022  EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "\u2022  EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "\u2022  Open-source resources used: Clima - Berkley, Streamlit, Python",
            "\u2022  EDS is not liable to inform Users about updates to the Application or underlying resources",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        for item in [
            "\u2022  Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "\u2022  Streamlit, \u00a9 Streamlit Inc., licensed under Apache 2.0",
            "\u2022  Python \u00a9 Python Software Foundation, licensed under PSF License Version 2",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_annexure_slide()

    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes


def generate_wind_pptx_report(
    df: pd.DataFrame,
    metadata: dict,
    n_sectors: int = 16,
):
    """Generate a Wind Analysis PowerPoint report using the Voha template."""

    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    except NameError:
        base_dir = os.getcwd()
    template_path = os.path.join(base_dir, "Voha Hospitality Climate analysis_v4 (2).pptx")
    logo_path = os.path.join(base_dir, "EDSlogo.png")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        _ppt_remove_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    BLANK_LAYOUT = prs.slide_layouts[6]
    TITLE_RED = RGBColor(0xC0, 0x00, 0x00)
    DARK_GREY  = RGBColor(0x40, 0x40, 0x40)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    SW = prs.slide_width.inches
    SH = prs.slide_height.inches
    LOGO_H = 0.40
    LOGO_W = LOGO_H * (550 / 308)
    LOGO_L = 0.18
    LOGO_T = SH - LOGO_H - 0.12

    def _add_logo(slide):
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(LOGO_L), Inches(LOGO_T),
                                     width=Inches(LOGO_W), height=Inches(LOGO_H))

    def _slide_title(slide, text, top=0.13):
        tb = slide.shapes.add_textbox(Inches(0.27), Inches(top), Inches(SW - 0.54), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_RED

    def _divider(slide, top_inches):
        bar = slide.shapes.add_shape(1, Inches(0.27), Inches(top_inches),
                                     Inches(SW - 0.54), Inches(0.03))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TITLE_RED
        bar.line.fill.background()

    def _err(slide, err):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(SW - 1), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Error: {err}"
        run.font.size = Pt(10)
        run.font.color.rgb = TITLE_RED

    def _save_fig(fig) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name, dpi=130, bbox_inches="tight", facecolor="white")
            return tmp.name

    # ── COVER ─────────────────────────────────────────────────────────────────
    def _cover():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        bg = slide.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(SW), Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_RED
        bg.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(SW - 1.2), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Wind Analysis Report"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE

        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.75), Inches(SW - 1.2), Inches(0.65))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        _city = metadata.get("city", "") if metadata else ""
        run2.text = f"Location: {_city}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(SW - 1.2), Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Sections: Wind Rose | Speed & Direction Heatmaps | Statistics | Climate Integration"
        run3.font.size = Pt(10)
        run3.font.color.rgb = DARK_GREY

        _add_logo(slide)

    _cover()

    # Import wind module utilities
    try:
        from modules.wind_module import (
            prepare_wind_data, compute_wind_rose, compute_wind_statistics,
            _SPEED_LABELS, _SPEED_COLORS, _SPEED_BINS,
            _DIR_16, _DIR_8, _DIR_4, _MONTH_NAMES, _MONTH_COLORS,
        )
    except ImportError as e:
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Analysis")
        _divider(slide, 0.62)
        _err(slide, f"Could not import wind module: {str(e)[:50]}")
        _add_logo(slide)
        
        report_bytes = io.BytesIO()
        prs.save(report_bytes)
        report_bytes.seek(0)
        return report_bytes

    # Prepare wind data
    months = list(range(1, 13))  # All months
    wdf = prepare_wind_data(df, months=months, n_sectors=n_sectors)

    if wdf.empty:
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Analysis")
        _divider(slide, 0.62)
        _err(slide, "No wind data available for the analysis period.")
        _add_logo(slide)
        
        report_bytes = io.BytesIO()
        prs.save(report_bytes)
        report_bytes.seek(0)
        return report_bytes

    rose_df, calm_pct = compute_wind_rose(wdf, n_sectors=n_sectors, exclude_calm=False)
    stats = compute_wind_statistics(wdf)

    # ── Matplotlib chart helpers (no kaleido / Chrome required) ──────────────
    def _mpl_wind_rose_png():
        sector_width = 360.0 / n_sectors
        if n_sectors == 16:
            lbl = _DIR_16
        elif n_sectors == 8:
            lbl = _DIR_8
        elif n_sectors == 4:
            lbl = _DIR_4
        else:
            lbl = [f"{int(i * sector_width)}°" for i in range(n_sectors)]
        angles = np.array([np.deg2rad(i * sector_width) for i in range(n_sectors)])
        bar_w = np.deg2rad(sector_width) * 0.85
        fig2, ax = plt.subplots(subplot_kw=dict(polar=True), figsize=(9, 7))
        bottoms = np.zeros(n_sectors)
        for i, sl in enumerate(_SPEED_LABELS):
            subset = rose_df[rose_df["speed_bin"] == sl]
            fm = dict(zip(subset["direction_label"], subset["frequency_pct"]))
            freqs = np.array([fm.get(l, 0.0) for l in lbl])
            ax.bar(angles, freqs, width=bar_w, bottom=bottoms,
                   color=_SPEED_COLORS[i % len(_SPEED_COLORS)],
                   label=f"{sl} m/s", alpha=0.9, linewidth=0.3, edgecolor="white")
            bottoms += freqs
        ax.set_theta_zero_location("N")
        ax.set_theta_direction(-1)
        ax.set_xticks(angles)
        ax.set_xticklabels(lbl, fontsize=9)
        ax.tick_params(axis="y", labelsize=8)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.1f}%"))
        ax.set_title("Wind Rose", fontsize=14, pad=20, color="#2c3e50", fontweight="bold")
        ax.annotate(f"Calm\n{calm_pct:.1f}%", xy=(0, 0), xycoords="data",
                    ha="center", va="center", fontsize=10, color="#555555")
        ax.legend(loc="lower left", bbox_to_anchor=(1.05, 0.0),
                  fontsize=9, title="Wind Speed (m/s)", title_fontsize=9)
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    def _mpl_speed_heatmap_png():
        pivot = wdf.pivot_table(values="wind_speed", index="hour",
                                columns="month", aggfunc="mean")
        for m in range(1, 13):
            if m not in pivot.columns:
                pivot[m] = np.nan
        pivot = pivot[sorted(pivot.columns)]
        month_labels = [_MONTH_NAMES[m - 1] for m in sorted(pivot.columns)]
        fig2, ax = plt.subplots(figsize=(12, 5))
        im = ax.imshow(pivot.values, aspect="auto", cmap="viridis", interpolation="nearest")
        ax.set_xlabel("Month", fontsize=11)
        ax.set_ylabel("Hour of Day", fontsize=11)
        ax.set_title("Wind Speed – Month × Hour", fontsize=14, color="#2c3e50", fontweight="bold")
        ax.set_xticks(range(12))
        ax.set_xticklabels(month_labels, fontsize=9)
        ax.set_yticks(range(0, 24, 3))
        ax.set_yticklabels([f"{h:02d}:00" for h in range(0, 24, 3)], fontsize=9)
        plt.colorbar(im, ax=ax, label="m/s", shrink=0.8)
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    def _mpl_direction_heatmap_png():
        tmp_df = wdf.copy()
        rad = np.deg2rad(tmp_df["wind_direction"])
        tmp_df["_u"] = np.cos(rad)
        tmp_df["_v"] = np.sin(rad)
        up = tmp_df.pivot_table(values="_u", index="hour", columns="month", aggfunc="mean")
        vp = tmp_df.pivot_table(values="_v", index="hour", columns="month", aggfunc="mean")
        up, vp = up.align(vp, join="inner")
        dir_deg = np.degrees(np.arctan2(vp.values, up.values)) % 360
        month_cols = sorted(up.columns.tolist())
        month_labels = [_MONTH_NAMES[m - 1] for m in month_cols]
        fig2, ax = plt.subplots(figsize=(12, 5))
        im = ax.imshow(dir_deg, aspect="auto", cmap="twilight",
                       vmin=0, vmax=360, interpolation="nearest")
        ax.set_xlabel("Month", fontsize=11)
        ax.set_ylabel("Hour of Day", fontsize=11)
        ax.set_title("Wind Direction – Month × Hour", fontsize=14, color="#2c3e50", fontweight="bold")
        ax.set_xticks(range(len(month_cols)))
        ax.set_xticklabels(month_labels, fontsize=9)
        ax.set_yticks(range(0, 24, 3))
        ax.set_yticklabels([f"{h:02d}:00" for h in range(0, 24, 3)], fontsize=9)
        cbar = plt.colorbar(im, ax=ax, shrink=0.8)
        cbar.set_ticks([0, 90, 180, 270, 360])
        cbar.set_ticklabels(["N 0°", "E 90°", "S 180°", "W 270°", "N 360°"])
        cbar.set_label("Direction")
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    def _mpl_speed_histogram_png():
        total = len(wdf)
        labels, pcts = [], []
        for i in range(len(_SPEED_BINS) - 1):
            lo, hi = _SPEED_BINS[i], _SPEED_BINS[i + 1]
            count = int(((wdf["wind_speed"] >= lo) & (wdf["wind_speed"] < hi)).sum())
            labels.append(_SPEED_LABELS[i])
            pcts.append(count / total * 100.0 if total > 0 else 0.0)
        fig2, ax = plt.subplots(figsize=(10, 5))
        bars = ax.bar(labels, pcts, color=_SPEED_COLORS[:len(labels)], alpha=0.9, edgecolor="white")
        for bar, pct in zip(bars, pcts):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                    f"{pct:.1f}%", ha="center", va="bottom", fontsize=10)
        ax.set_xlabel("Wind Speed Bin (m/s)", fontsize=11)
        ax.set_ylabel("Frequency (%)", fontsize=11)
        ax.set_title("Wind Speed Distribution", fontsize=14, color="#2c3e50", fontweight="bold")
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    def _mpl_climate_bubble_png():
        needed = {"dry_bulb_temperature", "relative_humidity", "wind_speed", "month"}
        fig2, ax = plt.subplots(figsize=(10, 6))
        if not needed.issubset(wdf.columns):
            ax.text(0.5, 0.5, "Missing columns for bubble chart",
                    ha="center", va="center", transform=ax.transAxes, fontsize=13)
        else:
            tmp_df = wdf.dropna(subset=["dry_bulb_temperature", "relative_humidity", "wind_speed"]).copy()
            max_spd = float(tmp_df["wind_speed"].max())
            scale = 500.0 / max_spd if max_spd > 0 else 50.0
            for m in range(1, 13):
                mdata = tmp_df[tmp_df["month"] == m]
                if mdata.empty:
                    continue
                ax.scatter(mdata["dry_bulb_temperature"], mdata["relative_humidity"],
                           s=(mdata["wind_speed"] + 0.3) * scale,
                           c=_MONTH_COLORS[(m - 1) % len(_MONTH_COLORS)],
                           alpha=0.45, linewidths=0, label=_MONTH_NAMES[m - 1])
            ax.set_xlabel("Dry Bulb Temperature (°C)", fontsize=11)
            ax.set_ylabel("Relative Humidity (%)", fontsize=11)
            ax.set_ylim(0, 105)
            ax.legend(title="Month", fontsize=9, title_fontsize=9,
                      loc="center left", bbox_to_anchor=(1, 0.5))
            ax.text(0.01, 0.98, "Bubble size = wind speed (m/s)",
                    transform=ax.transAxes, fontsize=10, color="#888", va="top")
        ax.set_title("Temperature – Humidity – Wind Speed", fontsize=14,
                     color="#2c3e50", fontweight="bold")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    def _mpl_single_season_rose_png(rose_df_s, calm_pct_s, season_name):
        sector_width = 360.0 / n_sectors
        if n_sectors == 16:
            lbl = _DIR_16
        elif n_sectors == 8:
            lbl = _DIR_8
        elif n_sectors == 4:
            lbl = _DIR_4
        else:
            lbl = [f"{int(i * sector_width)}°" for i in range(n_sectors)]
        angles = np.array([np.deg2rad(i * sector_width) for i in range(n_sectors)])
        bar_w = np.deg2rad(sector_width) * 0.85
        fig2, ax = plt.subplots(subplot_kw=dict(polar=True), figsize=(9, 7))
        bottoms = np.zeros(n_sectors)
        for i, sl in enumerate(_SPEED_LABELS):
            subset = rose_df_s[rose_df_s["speed_bin"] == sl]
            fm = dict(zip(subset["direction_label"], subset["frequency_pct"]))
            freqs = np.array([fm.get(l, 0.0) for l in lbl])
            ax.bar(angles, freqs, width=bar_w, bottom=bottoms,
                   color=_SPEED_COLORS[i % len(_SPEED_COLORS)],
                   label=f"{sl} m/s", alpha=0.9, linewidth=0.3, edgecolor="white")
            bottoms += freqs
        ax.set_theta_zero_location("N")
        ax.set_theta_direction(-1)
        ax.set_xticks(angles)
        ax.set_xticklabels(lbl, fontsize=9)
        ax.tick_params(axis="y", labelsize=8)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.1f}%"))
        ax.set_title(f"{season_name} Wind Rose", fontsize=14, pad=20,
                     color="#2c3e50", fontweight="bold")
        ax.annotate(f"Calm\n{calm_pct_s:.1f}%", xy=(0, 0), xycoords="data",
                    ha="center", va="center", fontsize=10, color="#555555")
        ax.legend(loc="lower left", bbox_to_anchor=(1.05, 0.0),
                  fontsize=9, title="Wind Speed (m/s)", title_fontsize=9)
        plt.tight_layout()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        fig2.savefig(path, dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig2)
        return path

    # ── WIND ROSE SLIDE ───────────────────────────────────────────────────────
    def _wind_rose_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Rose Analysis")
        _divider(slide, 0.62)

        try:
            img_path = _mpl_wind_rose_png()
            chart_w = 7.5
            slide.shapes.add_picture(img_path, Inches((SW - chart_w) / 2), Inches(0.80),
                                     width=Inches(chart_w))
            os.unlink(img_path)
        except Exception as e:
            _err(slide, f"Wind rose: {str(e)[:50]}")

        _add_logo(slide)

    _wind_rose_slide()

    # ── SEASONAL WIND ROSE SLIDES (one slide per season) ─────────────────────
    for _sname, _smonths in [("Winter", [12, 1, 2]), ("Spring", [3, 4, 5]),
                              ("Summer", [6, 7, 8]),  ("Fall",   [9, 10, 11])]:
        _sdf = wdf[wdf["month"].isin(_smonths)].copy()
        if _sdf.empty:
            continue
        _srose, _scalm = compute_wind_rose(_sdf, n_sectors=n_sectors, exclude_calm=False)
        _sslide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(_sslide, f"Wind Rose – {_sname}")
        _divider(_sslide, 0.62)
        try:
            _spath = _mpl_single_season_rose_png(_srose, _scalm, _sname)
            _cw = 7.5
            _sslide.shapes.add_picture(_spath, Inches((SW - _cw) / 2), Inches(0.80),
                                       width=Inches(_cw))
            os.unlink(_spath)
        except Exception as _se:
            _err(_sslide, f"{_sname} wind rose: {str(_se)[:50]}")
        _add_logo(_sslide)

    # ── WIND SPEED HEATMAP SLIDE ──────────────────────────────────────────────
    def _speed_heatmap_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Speed Heatmap (Month × Hour)")
        _divider(slide, 0.62)

        try:
            img_path = _mpl_speed_heatmap_png()
            slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(img_path)
        except Exception as e:
            _err(slide, f"Speed heatmap: {str(e)[:50]}")

        _add_logo(slide)

    _speed_heatmap_slide()

    # ── WIND DIRECTION HEATMAP SLIDE ──────────────────────────────────────────
    def _direction_heatmap_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Direction Heatmap (Month × Hour)")
        _divider(slide, 0.62)

        try:
            img_path = _mpl_direction_heatmap_png()
            slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(img_path)
        except Exception as e:
            _err(slide, f"Direction heatmap: {str(e)[:50]}")

        _add_logo(slide)

    _direction_heatmap_slide()

    # ── WIND SPEED HISTOGRAM SLIDE ────────────────────────────────────────────
    def _speed_histogram_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Speed Distribution")
        _divider(slide, 0.62)

        try:
            img_path = _mpl_speed_histogram_png()
            slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(img_path)
        except Exception as e:
            _err(slide, f"Speed histogram: {str(e)[:50]}")

        _add_logo(slide)

    _speed_histogram_slide()

    # ── CLIMATE BUBBLE CHART SLIDE ────────────────────────────────────────────
    def _climate_bubble_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Temperature – Humidity – Wind Speed")
        _divider(slide, 0.62)

        try:
            img_path = _mpl_climate_bubble_png()
            slide.shapes.add_picture(img_path, Inches(0.27), Inches(0.72),
                                     width=Inches(SW - 0.54), height=Inches(5.9))
            os.unlink(img_path)
        except Exception as e:
            _err(slide, f"Climate bubble: {str(e)[:50]}")

        _add_logo(slide)

    _climate_bubble_slide()

    # ── WIND STATISTICS SUMMARY SLIDE ─────────────────────────────────────────
    def _statistics_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Wind Statistics Summary")
        _divider(slide, 0.62)

        try:
            # Prepare statistics data
            stat_labels = [
                "Prevailing Direction",
                "Mean Wind Speed",
                "Maximum Wind Speed",
                "Calm Hours (%)",
                "Strongest Wind Direction",
                "Total Data Points",
            ]

            stat_values = [
                stats.get("prevailing_direction", "N/A"),
                f"{stats.get('mean_speed', 0):.2f} m/s",
                f"{stats.get('max_speed', 0):.2f} m/s",
                f"{stats.get('calm_percent', 0):.1f}%",
                stats.get("strongest_direction", "N/A"),
                f"{len(wdf)} hours",
            ]

            # Create a two-column layout
            col_width = (SW - 1.0) / 2
            left_x = 0.5
            right_x = left_x + col_width + 0.2

            # Left column
            left_tb = slide.shapes.add_textbox(Inches(left_x), Inches(0.8), Inches(col_width), Inches(5.5))
            left_tf = left_tb.text_frame
            left_tf.word_wrap = True

            for i in range(0, 3):
                label = stat_labels[i]
                value = stat_values[i]
                
                p = left_tf.add_paragraph() if i > 0 else left_tf.paragraphs[0]
                p.text = ""
                p.space_before = Pt(12)
                p.space_after = Pt(2)

                run = p.add_run()
                run.text = label
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = TITLE_RED

                p2 = left_tf.add_paragraph()
                p2.text = value
                p2.font.size = Pt(11)
                p2.font.color.rgb = DARK_GREY
                p2.space_after = Pt(0)

            # Right column
            right_tb = slide.shapes.add_textbox(Inches(right_x), Inches(0.8), Inches(col_width), Inches(5.5))
            right_tf = right_tb.text_frame
            right_tf.word_wrap = True

            for i in range(3, 6):
                label = stat_labels[i]
                value = stat_values[i]
                
                p = right_tf.add_paragraph() if i == 3 else right_tf.add_paragraph()
                p.text = ""
                p.space_before = Pt(12)
                p.space_after = Pt(2)

                run = p.add_run()
                run.text = label
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = TITLE_RED

                p2 = right_tf.add_paragraph()
                p2.text = value
                p2.font.size = Pt(11)
                p2.font.color.rgb = DARK_GREY
                p2.space_after = Pt(0)

        except Exception as e:
            _err(slide, f"Statistics: {str(e)[:50]}")

        _add_logo(slide)

    _statistics_slide()

    # ── ANNEXURE ──────────────────────────────────────────────────────────────
    def _make_annexure_slide():
        slide = prs.slides.add_slide(BLANK_LAYOUT)
        _slide_title(slide, "Annexure")
        _divider(slide, 0.62)

        tb = slide.shapes.add_textbox(Inches(0.27), Inches(0.80), Inches(SW - 0.54), Inches(6.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "About EDS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = (
            "Environmental Design Solutions [EDS] is a sustainability advisory firm. "
            "Since 2002, EDS has worked on over 500 green building and energy efficiency "
            "projects worldwide. The team focuses on climate change mitigation, low-carbon "
            "design, building simulation, performance audits, and capacity building. EDS "
            "continues to contribute to the buildings community with useful tools through "
            "its IT services."
        )
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GREY
        p.line_spacing = 1.2
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = "Disclaimer"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        for item in [
            "Climate Zone Analyser is an outcome of the best efforts of building simulation experts at EDS.",
            "\u2022  EDS does not assume responsibility for outcomes from its use. By using this Application, the User indemnifies EDS against any damages.",
            "\u2022  EDS does not guarantee uninterrupted availability. By using this Application, the User agrees to share uploaded information with EDS for analysis and research purposes.",
            "\u2022  Open-source resources used: Clima - Berkley, Streamlit, Python",
            "\u2022  EDS is not liable to inform Users about updates to the Application or underlying resources",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = "Acknowledgement"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_RED
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        for item in [
            "\u2022  Betti, G., et al. CBE Clima Tool Build. Simul. (2023). https://doi.org/10.1007/s12273-023-1090-5",
            "\u2022  Streamlit, \u00a9 Streamlit Inc., licensed under Apache 2.0",
            "\u2022  Python \u00a9 Python Software Foundation, licensed under PSF License Version 2",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GREY
            p.line_spacing = 1.1
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        _add_logo(slide)

    _make_annexure_slide()

    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    return report_bytes
